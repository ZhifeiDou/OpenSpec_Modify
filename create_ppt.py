"""
Generate PPT: OpenSpec 技术研究汇报
Target audience: 有技术背景、使用过 Claude Code 但未接触 OpenSpec 的技术领导
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Image Paths ──
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DASHBOARD = os.path.join(BASE_DIR, "OpenSpec-main", "assets", "openspec_dashboard.png")
IMG_WORKFLOW = os.path.join(BASE_DIR, "ppt_images", "workflow_diagram.jpg")
IMG_SEQUENCE = os.path.join(BASE_DIR, "ppt_images", "openspec_sequence.png")
IMG_LINEAR = os.path.join(BASE_DIR, "ppt_images", "linear_workflow.jpg")

# ── Color Palette ──
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_RED_BG = RGBColor(0xFD, 0xE8, 0xE8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        spPr = shape._element.spPr
        solidFill_elem = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill_elem is not None:
            srgbClr = solidFill_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr is not None:
                alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=DARK_GRAY, line_spacing=1.5, font_name="Microsoft YaHei",
                      alignment=PP_ALIGN.LEFT):
    """lines: list of (text, bold, color_override, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, c, fs = item, False, color, font_size
        elif len(item) == 2:
            text, bold = item
            c, fs = color, font_size
        elif len(item) == 3:
            text, bold, c = item
            fs = font_size
        else:
            text, bold, c, fs = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=ACCENT_BLUE, bg_color=WHITE, border_color=None):
    """Add a card-like shape with title and body."""
    card = add_rect(slide, left, top, width, height, bg_color)
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_multiline_box(slide, left + Inches(0.2), top + Inches(0.6),
                      width - Inches(0.4), height - Inches(0.8),
                      body_lines, font_size=13, color=DARK_GRAY, line_spacing=1.4)
    return card


def add_page_number(slide, num, total):
    add_text_box(slide, W - Inches(1.2), H - Inches(0.5),
                 Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, text):
    """Top accent bar + section title."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 text, font_size=28, color=BLACK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(1.2), Inches(0.05), ACCENT_BLUE)


TOTAL_SLIDES = 15

# ════════════════════════════════════════════════════════════════════════
# Slide 1: Cover
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "OpenSpec：规范化 Agentic Coding 的关键基础设施",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "基于 Claude Code 的企业级 AI 辅助开发规范化研究",
             font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
             "技术研究汇报  |  2026年2月",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 2: 目录
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "目录")

toc_items = [
    ("01", "背景与问题", "Agentic Coding 现状与挑战"),
    ("02", "OpenSpec 核心概念", "框架介绍、运行界面与工作流程"),
    ("03", "实测对比与洞察", "量化交易系统 — 有/无 OpenSpec 对比"),
    ("04", "企业级价值", "规范化、可追溯、可协作"),
    ("05", "商业价值与 ROI", "市场机会与投资回报分析"),
    ("06", "研究计划与建议", "落地路径与资源需求"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(0.85) * i
    add_rect(slide, Inches(1.2), y, Inches(0.7), Inches(0.6), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(0.7), Inches(0.4),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.2), y + Inches(0.02), Inches(4), Inches(0.35),
                 title, font_size=18, color=BLACK, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.35), Inches(8), Inches(0.3),
                 desc, font_size=13, color=MID_GRAY)

add_page_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 3: 背景 — Agentic Coding 现状
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "背景：Agentic Coding 正在改变软件开发")

add_multiline_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.8), [
    ("Claude Code 等 AI 编程工具已经能够自主完成复杂的软件开发任务。", False, DARK_GRAY, 16),
    ('但在企业环境中，"能写代码"只是第一步 - 更关键的是"写对的代码"。', False, DARK_GRAY, 16),
], line_spacing=1.6)

# Pain points
pain_points = [
    ("需求失控", "AI 容易自作主张添加/遗漏功能\n需求只存在于聊天记录中，无法追溯", ACCENT_RED, LIGHT_RED_BG),
    ("质量不稳定", "同样的需求，不同对话产出差异大\n缺乏验证标准，靠人肉 review", ACCENT_ORANGE, LIGHT_ORANGE_BG),
    ("协作困难", "一个人的 AI 对话上下文无法共享\n多人并行修改容易冲突", ACCENT_BLUE, LIGHT_BLUE_BG),
    ("知识流失", "为什么做、怎么做的决策过程随对话消失\n新人接手困难，重复踩坑", DARK_GRAY, LIGHT_GRAY),
]

for i, (title, desc, tc, bgc) in enumerate(pain_points):
    x = Inches(0.8) + Inches(3.0) * i
    y = Inches(3.0)
    card = add_rect(slide, x, y, Inches(2.8), Inches(3.5), bgc)
    card.line.color.rgb = tc
    card.line.width = Pt(1.5)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.5),
                 title, font_size=20, color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(2.2),
                      desc.split('\n'), font_size=13, color=DARK_GRAY, line_spacing=1.6,
                      alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             'Claude Code 解决了"怎么写"，但没有解决"写什么"和"写得对不对"',
             font_size=15, color=ACCENT_RED, bold=True)

add_page_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 4: OpenSpec 是什么
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "OpenSpec：AI 辅助开发的规格框架")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
             "OpenSpec 是一个轻量级的规格驱动框架，让人和 AI 在写代码之前先达成共识。",
             font_size=16, color=DARK_GRAY)

# Core principles
principles = [
    ("流动而非僵化", "没有阶段门控，按需推进"),
    ("迭代而非瀑布", "边做边学，持续优化"),
    ("增量而非全量", "用 Delta 记录变更，不重写整体"),
    ("存量优先", "为已有项目设计，不是只给新项目"),
]

for i, (title, desc) in enumerate(principles):
    x = Inches(0.8) + Inches(3.0) * i
    y = Inches(2.3)
    add_rect(slide, x, y, Inches(2.8), Inches(1.2), LIGHT_BLUE_BG)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 title, font_size=15, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.5), Inches(0.4),
                 desc, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Directory structure
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.4),
             "项目结构（两个核心目录）", font_size=16, color=BLACK, bold=True)

struct_text = """openspec/
├── specs/              ← 当前行为的真相源
│   └── <domain>/spec.md    （按领域组织的需求）
└── changes/            ← 提议的变更
    ├── <change-name>/
    │   ├── proposal.md     （为什么改）
    │   ├── usecases.md     （用户场景）
    │   ├── specs/          （Delta: 增/改/删）
    │   ├── design.md       （怎么改）
    │   └── tasks.md        （实现清单）
    └── archive/            （已完成变更的归档）"""

add_text_box(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(3.0),
             struct_text, font_size=11, color=DARK_GRAY,
             font_name="Consolas")

# Key insight box
add_rect(slide, Inches(7.0), Inches(3.9), Inches(5.5), Inches(3.0), LIGHT_GREEN_BG)
add_text_box(slide, Inches(7.2), Inches(4.1), Inches(5.1), Inches(0.4),
             "核心创新：Delta Spec（增量规格）", font_size=15, color=ACCENT_GREEN, bold=True)
add_multiline_box(slide, Inches(7.2), Inches(4.6), Inches(5.1), Inches(2.2), [
    ("不是重写整个规格，而是只描述变化：", False, DARK_GRAY, 13),
    ("", False),
    ("## ADDED Requirements      ← 新增需求", False, ACCENT_GREEN, 12),
    ("## MODIFIED Requirements  ← 修改需求", False, ACCENT_ORANGE, 12),
    ("## REMOVED Requirements   ← 删除需求", False, ACCENT_RED, 12),
    ("", False),
    ("每个需求附带 Given/When/Then 验证场景", False, DARK_GRAY, 13),
    ("归档时自动合并到主规格 → 活文档", False, DARK_GRAY, 13),
], line_spacing=1.3, font_name="Microsoft YaHei")

add_page_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 5: OpenSpec 实际运行效果（图片页）
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "OpenSpec 实际运行效果")

# Left: Dashboard screenshot
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "Dashboard 总览（终端界面）", font_size=15, color=ACCENT_BLUE, bold=True)
if os.path.exists(IMG_DASHBOARD):
    slide.shapes.add_picture(IMG_DASHBOARD, Inches(0.8), Inches(1.9),
                             Inches(5.2), Inches(4.5))
# Right: Workflow diagram
add_text_box(slide, Inches(6.5), Inches(1.4), Inches(6), Inches(0.4),
             "Change Workflow 全景（人机协作流程）", font_size=15, color=ACCENT_GREEN, bold=True)
if os.path.exists(IMG_WORKFLOW):
    slide.shapes.add_picture(IMG_WORKFLOW, Inches(6.5), Inches(1.9),
                             Inches(6.3), Inches(4.5))

# Caption
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.6),
             "左图：openspec view 命令展示项目规格状态、活跃变更和任务进度  |  右图：完整的 Propose → Review → Apply → Archive 工作流",
             font_size=12, color=MID_GRAY)

# Source attribution
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.3),
             "图片来源：github.com/Fission-AI/OpenSpec  |  harikrishnan.io",
             font_size=10, color=MID_GRAY)

add_page_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 6: 工作流程
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "OpenSpec 工作流程：从想法到归档")

# Workflow steps as a horizontal pipeline
steps = [
    ("Explore", "探索想法\n明确方向", ACCENT_BLUE, "可选"),
    ("Proposal", "为什么改\n改什么", ACCENT_BLUE, ""),
    ("Use Cases", "用户场景\n主流程+异常", ACCENT_BLUE, ""),
    ("Specs", "Delta 规格\nGiven/When/Then", ACCENT_GREEN, ""),
    ("Design", "技术方案\n架构决策", ACCENT_GREEN, "可选"),
    ("Tasks", "实现清单\n可追踪进度", ACCENT_ORANGE, ""),
    ("Apply", "AI 按清单\n写代码", ACCENT_ORANGE, ""),
    ("Verify", "验证实现\n匹配规格", ACCENT_RED, ""),
    ("Archive", "合并规格\n归档变更", ACCENT_RED, ""),
]

for i, (name, desc, color, note) in enumerate(steps):
    x = Inches(0.5) + Inches(1.38) * i
    y = Inches(2.0)
    # Box
    add_rect(slide, x, y, Inches(1.25), Inches(2.2), color)
    add_text_box(slide, x, y + Inches(0.15), Inches(1.25), Inches(0.4),
                 name, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.05), y + Inches(0.6), Inches(1.15), Inches(1.2),
                      desc.split('\n'), font_size=11, color=WHITE, line_spacing=1.3,
                      alignment=PP_ALIGN.CENTER)
    if note:
        add_text_box(slide, x, y + Inches(1.8), Inches(1.25), Inches(0.3),
                     note, font_size=10, color=RGBColor(0xFF, 0xFF, 0xCC),
                     alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(steps) - 1:
        arrow_x = x + Inches(1.25)
        add_text_box(slide, arrow_x, y + Inches(0.85), Inches(0.13), Inches(0.3),
                     "→", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Key points below
add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.04), LIGHT_GRAY)

key_points = [
    ("一键加速", "/opsx:ff 可以一步生成所有工件，省去逐步创建", ACCENT_BLUE),
    ("灵活不僵化", "工件之间是依赖图（DAG），不是固定流水线，可跳过可选步骤", ACCENT_GREEN),
    ("闭环演进", "Archive 后 specs 自动更新 → 下次变更基于最新状态 → 持续积累", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(key_points):
    y = Inches(5.1) + Inches(0.7) * i
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.5), color)
    add_text_box(slide, Inches(0.8), y + Inches(0.02), Inches(2), Inches(0.4),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.02), Inches(10), Inches(0.4),
                 desc, font_size=14, color=DARK_GRAY)

add_page_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 7: 协作模式详解（序列图 + Linear 集成）
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "OpenSpec 协作模式：Project Owner / AI / Developer")

# Left: Sequence diagram
if os.path.exists(IMG_SEQUENCE):
    slide.shapes.add_picture(IMG_SEQUENCE, Inches(0.5), Inches(1.5),
                             Inches(5.8), Inches(5.5))

# Right: Linear integration diagram
add_text_box(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
             "企业级集成示例：Linear + OpenSpec", font_size=15, color=ACCENT_GREEN, bold=True)
if os.path.exists(IMG_LINEAR):
    slide.shapes.add_picture(IMG_LINEAR, Inches(6.8), Inches(1.9),
                             Inches(5.8), Inches(4.0))

# Key insights on the right, below Linear diagram
add_multiline_box(slide, Inches(6.8), Inches(6.1), Inches(5.8), Inches(1.2), [
    ("三方协作：Owner 定义需求 → OpenSpec 管理规格 → Developer/AI 实现", False, DARK_GRAY, 12),
    ("全程可追踪：每个阶段的产出都是持久化的 Markdown 文件", False, DARK_GRAY, 12),
    ("可集成现有工具链：Linear、Jira、GitHub Issues 等项目管理工具", False, ACCENT_BLUE, 12),
], line_spacing=1.4)

# Source attribution
add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "图片来源：hashrocket.com/blog  |  intent-driven.dev",
             font_size=10, color=MID_GRAY)

add_page_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 8: 实测对比 — 概览
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "实测对比：同一需求，有/无 OpenSpec 的差异")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "测试项目：A股有色金属量化选股交易系统（相同需求描述，分别用两种方式让 Claude Code 实现）",
             font_size=15, color=DARK_GRAY)

# Comparison table header
table_y = Inches(2.2)
add_rect(slide, Inches(0.8), table_y, Inches(3.5), Inches(0.6), MID_GRAY)
add_text_box(slide, Inches(0.8), table_y + Inches(0.1), Inches(3.5), Inches(0.4),
             "对比维度", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.3), table_y, Inches(4.2), Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(4.3), table_y + Inches(0.1), Inches(4.2), Inches(0.4),
             "使用 OpenSpec", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(8.5), table_y, Inches(4.2), Inches(0.6), ACCENT_ORANGE)
add_text_box(slide, Inches(8.5), table_y + Inches(0.1), Inches(4.2), Inches(0.4),
             "不使用 OpenSpec", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ("代码总量", "~3,900 行", "~1,200 行"),
    ("文件数量", "40+ 文件，模块化设计", "5 个文件，单体结构"),
    ("设计模式", "适配器模式、工厂模式、依赖注入", "直接实现，全局变量"),
    ("测试覆盖", "18 个测试文件（单元+集成+端到端）", "无测试"),
    ("配置管理", "外部 YAML 配置文件", "Python 常量硬编码"),
    ("错误处理", "多数据源 fallback、显式异常处理", "基本 try-catch"),
    ("文档", "8 个领域规格 + BDD 验证场景", "仅代码内注释"),
    ("可维护性", "高 — 明确的模块契约", "低 — 紧耦合"),
]

for i, (dim, openspec_val, no_openspec_val) in enumerate(rows):
    y = table_y + Inches(0.6) + Inches(0.52) * i
    bgc = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(3.5), Inches(0.52), bgc)
    add_rect(slide, Inches(4.3), y, Inches(4.2), Inches(0.52), bgc)
    add_rect(slide, Inches(8.5), y, Inches(4.2), Inches(0.52), bgc)

    add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(3.1), Inches(0.36),
                 dim, font_size=13, color=BLACK, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.08), Inches(3.8), Inches(0.36),
                 openspec_val, font_size=12, color=DARK_GRAY)
    add_text_box(slide, Inches(8.7), y + Inches(0.08), Inches(3.8), Inches(0.36),
                 no_openspec_val, font_size=12, color=DARK_GRAY)

add_page_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 9: 实测对比 — 关键洞察
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "实测洞察：OpenSpec 不只是让代码变多，而是让代码变好")

# Insight cards
insights = [
    ("架构质量跃升", ACCENT_BLUE, LIGHT_BLUE_BG, [
        "OpenSpec 的 specs 和 design 工件迫使 AI 在",
        "写代码前思考架构。结果：",
        "",
        "• 多数据源适配器 vs 硬编码单一数据源",
        "• 12 类配置外部化 vs Python 常量",
        "• 模块间通过接口通信 vs 直接引用",
    ]),
    ("测试从 0 到 18", ACCENT_GREEN, LIGHT_GREEN_BG, [
        "任务清单（tasks.md）明确要求测试，",
        "AI 无法跳过。相比之下：",
        "",
        "• OpenSpec: 18 个测试文件，覆盖",
        "  单元测试、集成测试、端到端测试",
        "• 无 OpenSpec: 0 个测试文件",
    ]),
    ("可维护性差距", ACCENT_ORANGE, LIGHT_ORANGE_BG, [
        "6 个月后需要修改因子计算逻辑：",
        "",
        "• OpenSpec: 查 specs → 改 delta →",
        "  更新 tasks → 精确修改模块",
        "• 无 OpenSpec: 在 1200 行单体文件",
        "  里全局搜索，祈祷不破坏其他功能",
    ]),
    ("知识沉淀", ACCENT_RED, LIGHT_RED_BG, [
        "为什么选 AKShare 而非 Tushare？",
        "",
        "• OpenSpec: design.md 记录了决策理由",
        "  和替代方案的权衡",
        "• 无 OpenSpec: 没人知道，当时的对话",
        "  记录已经找不到了",
    ]),
]

for i, (title, tc, bgc, lines) in enumerate(insights):
    x = Inches(0.5) + Inches(3.15) * i
    y = Inches(1.6)
    card = add_rect(slide, x, y, Inches(3.0), Inches(4.8), bgc)
    card.line.color.rgb = tc
    card.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), Inches(2.6), Inches(0.5),
                 title, font_size=17, color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.6), Inches(3.6),
                      lines, font_size=12, color=DARK_GRAY, line_spacing=1.3)

add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
             "结论：OpenSpec 的价值不在于写更多代码，而在于让 AI 写出工程级别的代码",
             font_size=15, color=ACCENT_BLUE, bold=True)

add_page_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 10: 企业级价值
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "企业级价值：为什么公司需要 OpenSpec")

value_items = [
    ("规范化", ACCENT_BLUE, [
        ("统一 AI 编码标准", "所有团队遵循相同的 proposal → spec → design → task → implement 流程"),
        ("需求有据可查", "每个变更都有 proposal（为什么）、spec（是什么）、design（怎么做）"),
        ("质量有标准可验", "BDD 场景（Given/When/Then）提供明确的验收标准"),
    ]),
    ("可追溯", ACCENT_GREEN, [
        ("变更全链路追踪", "从需求 → 设计 → 任务 → 代码，每一步都有记录"),
        ("审计友好", "归档保留完整上下文，满足合规要求"),
        ("决策可回溯", "design.md 记录技术选型理由和替代方案的权衡"),
    ]),
    ("可协作", ACCENT_ORANGE, [
        ("多人并行不冲突", "Delta spec 机制让不同变更可以并行，各自只描述自己的增量"),
        ("上下文可共享", "规格文件是团队共享的，不是某个人的聊天记录"),
        ("新人可快速接手", "读 specs + design + archived changes 即可理解系统演化历程"),
    ]),
]

for i, (title, color, items) in enumerate(value_items):
    x = Inches(0.6) + Inches(4.1) * i
    y = Inches(1.5)
    add_rect(slide, x, y, Inches(3.9), Inches(0.65), color)
    add_text_box(slide, x, y + Inches(0.12), Inches(3.9), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, (sub_title, sub_desc) in enumerate(items):
        sy = y + Inches(0.85) + Inches(1.5) * j
        add_text_box(slide, x + Inches(0.15), sy, Inches(3.6), Inches(0.35),
                     sub_title, font_size=14, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), sy + Inches(0.35), Inches(3.6), Inches(1.0),
                     sub_desc, font_size=12, color=DARK_GRAY)

add_page_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 11: 商业价值 — 市场分析
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "商业价值：Agentic Coding 治理是蓝海市场")

# Market context
add_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.5), LIGHT_BLUE_BG)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
             "市场趋势", font_size=16, color=ACCENT_BLUE, bold=True)
add_multiline_box(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.8), [
    ('AI 编程工具正从"辅助"走向"自主"（Agentic）', False, DARK_GRAY, 13),
    ("• Cursor, Windsurf, Claude Code 等工具快速普及", False, DARK_GRAY, 13),
    ("• 企业采用率快速增长，但治理和规范严重缺失", False, DARK_GRAY, 13),
    ("• 当前阶段类比：代码版本管理出现前的混乱期", False, ACCENT_RED, 13),
    ("• OpenSpec 定位 = Agentic Coding 时代的 Git", False, ACCENT_BLUE, 14),
], line_spacing=1.5)

# Why now
add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5), LIGHT_GREEN_BG)
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
             "为什么现在值得研究", font_size=16, color=ACCENT_GREEN, bold=True)
add_multiline_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(1.8), [
    ("• OpenSpec 仍处于早期阶段，先发优势明显", False, DARK_GRAY, 13),
    ("• 与 Claude Code 深度集成，Anthropic 生态优势", False, DARK_GRAY, 13),
    ("• 开源项目，可以基于自身需求定制和贡献", False, DARK_GRAY, 13),
    ("• 规格驱动思想与现有 DevOps 体系天然兼容", False, DARK_GRAY, 13),
    ("• 掌握核心方法论 → 可输出咨询/工具/培训服务", False, ACCENT_GREEN, 14),
], line_spacing=1.5)

# Business opportunities
add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.5),
             "商业化方向", font_size=16, color=BLACK, bold=True)

biz_items = [
    ("企业内部效率提升", "统一 AI 编码规范\n减少返工、提升交付质量\n降低代码审查成本", ACCENT_BLUE),
    ("工具产品化", "基于 OpenSpec 的企业级\n管理平台（Web UI、审批\n流程、数据统计）", ACCENT_GREEN),
    ("咨询与培训", "Agentic Coding 治理方法论\n企业落地咨询、团队培训\n最佳实践输出", ACCENT_ORANGE),
    ("生态贡献", "自定义 Schema 市场\n行业特定工作流模板\n插件和集成开发", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(biz_items):
    x = Inches(0.5) + Inches(3.2) * i
    y = Inches(5.0)
    add_rect(slide, x, y, Inches(3.0), Inches(2.2), color)
    add_text_box(slide, x, y + Inches(0.15), Inches(3.0), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(1.4),
                      desc.split('\n'), font_size=12, color=WHITE, line_spacing=1.4,
                      alignment=PP_ALIGN.CENTER)

add_page_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 12: ROI 分析
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "ROI 分析：规范化带来的投资回报")

# Before/After comparison
add_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), LIGHT_RED_BG)
add_text_box(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.5),
             "Without OpenSpec（当前状态）", font_size=18, color=ACCENT_RED, bold=True,
             alignment=PP_ALIGN.CENTER)

without_items = [
    "需求理解偏差率高 → 返工率 30-50%",
    "无测试 → 上线后 bug 修复成本高 (10-100x)",
    "代码审查耗时 → 审查者需从头理解意图",
    "知识流失 → 新人上手周期长 (2-4 周)",
    "多人协作冲突 → 集成时间占 20-30%",
    "无法审计 → 合规风险",
]
for i, item in enumerate(without_items):
    add_text_box(slide, Inches(0.8), Inches(2.3) + Inches(0.65) * i,
                 Inches(5.4), Inches(0.5),
                 f"✗  {item}", font_size=13, color=DARK_GRAY)

add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), LIGHT_GREEN_BG)
add_text_box(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
             "With OpenSpec（目标状态）", font_size=18, color=ACCENT_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)

with_items = [
    "规格前置 → 返工率降至 10-15%",
    "任务清单强制测试 → bug 在开发阶段发现",
    "审查者读 proposal + spec → 审查效率提升 50%",
    "活文档 + 归档 → 新人上手周期缩短至 3-5 天",
    "Delta spec 隔离 → 并行开发无冲突",
    "全链路追踪 → 天然满足审计需求",
]
for i, item in enumerate(with_items):
    add_text_box(slide, Inches(7.1), Inches(2.3) + Inches(0.65) * i,
                 Inches(5.4), Inches(0.5),
                 f"✓  {item}", font_size=13, color=DARK_GRAY)

add_page_number(slide, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 13: 与现有工具链的关系
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "定位：OpenSpec 与现有工具链的互补关系")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             'OpenSpec 不是替代品，而是缺失的一环 - 填补 AI 编程的"治理真空"',
             font_size=15, color=DARK_GRAY)

# Layer diagram
layers = [
    ("需求管理层", "OpenSpec Specs + Changes", ACCENT_BLUE, "WHO / WHAT / WHY"),
    ("AI 执行层", "Claude Code / Cursor / Copilot", ACCENT_GREEN, "HOW（写代码）"),
    ("版本控制层", "Git / GitHub / GitLab", ACCENT_ORANGE, "版本管理 + 协作"),
    ("CI/CD 层", "Jenkins / GitHub Actions", MID_GRAY, "构建 + 部署"),
]

for i, (layer, tool, color, desc) in enumerate(layers):
    y = Inches(2.2) + Inches(1.2) * i
    # Layer bar
    add_rect(slide, Inches(1.5), y, Inches(2.5), Inches(0.9), color)
    add_text_box(slide, Inches(1.5), y + Inches(0.2), Inches(2.5), Inches(0.5),
                 layer, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Tool
    add_rect(slide, Inches(4.2), y, Inches(4.5), Inches(0.9), LIGHT_GRAY)
    add_text_box(slide, Inches(4.2), y + Inches(0.2), Inches(4.5), Inches(0.5),
                 tool, font_size=14, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, Inches(9.0), y + Inches(0.2), Inches(3.5), Inches(0.5),
                 desc, font_size=13, color=MID_GRAY)

# Highlight the first layer
add_rect(slide, Inches(1.2), Inches(2.05), Inches(8.0), Inches(1.2), ACCENT_BLUE, alpha=10)

add_text_box(slide, Inches(1.5), Inches(7.0) - Inches(0.6), Inches(10), Inches(0.5),
             'OpenSpec 管理的是 AI 的"输入质量" - 输入越精确，AI 的输出越可控',
             font_size=15, color=ACCENT_BLUE, bold=True)

add_page_number(slide, 13, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 14: 研究计划
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "研究计划与落地建议")

phases = [
    ("第一阶段：深度评估", "2 周", ACCENT_BLUE, [
        "完成 3-5 个不同复杂度项目的对比测试",
        "量化评估代码质量指标（测试覆盖率、模块化程度、缺陷密度）",
        "评估 Token 消耗 / 时间成本的 overhead",
        "输出：《OpenSpec 技术评估报告》",
    ]),
    ("第二阶段：试点落地", "4 周", ACCENT_GREEN, [
        "选择 1 个实际业务项目进行试点",
        "定制适合团队的 Schema 工作流",
        "制定团队 Agentic Coding 规范",
        "输出：《企业 AI 编码规范 v1.0》",
    ]),
    ("第三阶段：推广优化", "4 周", ACCENT_ORANGE, [
        "团队培训和推广",
        "收集反馈，迭代优化工作流",
        "探索工具产品化可行性",
        "输出：推广方案 + 产品化可行性分析",
    ]),
]

for i, (phase, duration, color, items) in enumerate(phases):
    x = Inches(0.5) + Inches(4.2) * i
    y = Inches(1.5)

    add_rect(slide, x, y, Inches(3.9), Inches(0.7), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.0), Inches(0.4),
                 phase, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(2.8), y + Inches(0.1), Inches(0.95), Inches(0.4),
                 duration, font_size=13, color=WHITE, alignment=PP_ALIGN.RIGHT)

    for j, item in enumerate(items):
        iy = y + Inches(0.9) + Inches(0.55) * j
        add_text_box(slide, x + Inches(0.15), iy, Inches(3.6), Inches(0.5),
                     f"• {item}", font_size=12, color=DARK_GRAY)

# Resource ask
add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.04), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(5.0), Inches(5), Inches(0.4),
             "资源需求", font_size=16, color=BLACK, bold=True)

add_multiline_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), [
    ("• 人力：1 名研究员（全职） + 1 名高级工程师（兼职 review）", False, DARK_GRAY, 14),
    ("• 工具：Claude Code Pro 账号（API 用量约 $200-500/月）", False, DARK_GRAY, 14),
    ("• 时间：10 周完成三个阶段", False, DARK_GRAY, 14),
    ("• 预期产出：技术评估报告 + 企业编码规范 + 试点项目 + 产品化方案", False, ACCENT_BLUE, 14),
], line_spacing=1.5)

add_page_number(slide, 14, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════════════════
# Slide 15: 总结
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
             "总结", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

summary_items = [
    ("Agentic Coding 已是趋势", "Claude Code 等工具让 AI 自主写代码成为现实，但企业缺少治理手段"),
    ("OpenSpec 填补治理空白", '在"AI 能写代码"和"AI 写对代码"之间架起桥梁，用规格驱动 AI 行为'),
    ("实测验证：质量提升显著", "3.26x 代码模块化、18 vs 0 测试文件、完整架构设计 vs 单体代码"),
    ("商业价值：蓝海赛道", "先发优势 + 开源可定制 + 咨询/工具/培训多种变现路径"),
    ("建议：立项研究", "10 周、低成本投入，高潜力回报"),
]

for i, (title, desc) in enumerate(summary_items):
    y = Inches(2.2) + Inches(0.9) * i
    add_rect(slide, Inches(1.5), y, Inches(0.08), Inches(0.65), ACCENT_BLUE)
    add_text_box(slide, Inches(1.8), y + Inches(0.02), Inches(4), Inches(0.4),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(5.8), y + Inches(0.02), Inches(6), Inches(0.6),
                 desc, font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))

add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.5),
             "谢谢  |  Q&A",
             font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 15, TOTAL_SLIDES)

# ── Save ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "OpenSpec_技术研究汇报.pptx")
prs.save(output_path)
print("PPT saved successfully!")
