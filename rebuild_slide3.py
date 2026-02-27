#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Rebuild slide 3 of the PPT with a better layout."""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

prs = Presentation('OpenSpec_技术研究汇报_精简版.pptx')
slide = prs.slides[2]

# ── Remove all existing shapes ──
spTree = slide.shapes._spTree
for sp in list(spTree):
    tag = etree.QName(sp.tag).localname
    if tag in ('sp', 'pic', 'grpSp', 'cxnSp'):
        spTree.remove(sp)

print(f"Shapes after clear: {len(slide.shapes)}")

# ── Helper functions ──
def _in(val):
    """Convert inches to EMU."""
    return Inches(val)

def add_rect(left, top, width, height, fill_color=None, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _in(left), _in(top), _in(width), _in(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    else:
        shape.fill.background()
    if not line:
        shape.line.fill.background()
    return shape

def add_rounded_rect(left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _in(left), _in(top), _in(width), _in(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = RGBColor.from_string(line_color)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(left, top, width, height, text, size=12, color='333333', bold=False,
             align=PP_ALIGN.LEFT, font_name='Microsoft YaHei', wrap=True, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(_in(left), _in(top), _in(width), _in(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = font_name
    p.font.color.rgb = RGBColor.from_string(color)
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline(left, top, width, height, lines, size=9, color='333333',
                  font_name='Microsoft YaHei', line_spacing=1.15):
    """lines: list of (text, color, bold) tuples."""
    txBox = slide.shapes.add_textbox(_in(left), _in(top), _in(width), _in(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, c, b) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.name = font_name
        p.font.color.rgb = RGBColor.from_string(c if c else color)
        p.font.bold = b
        p.space_after = Pt(1)
        try:
            p.line_spacing = line_spacing
        except:
            pass
    return txBox


# ========================================================================
# LAYOUT CONSTANTS
# ========================================================================
SLIDE_W = 13.33
SLIDE_H = 7.50

# Margins
M_LEFT = 0.35
M_RIGHT = 0.35
M_TOP_CONTENT = 1.25
M_BOTTOM = 0.20

# Columns
LEFT_COL_W = 3.70
GAP = 0.20
LEFT_COL_X = M_LEFT
RIGHT_COL_X = LEFT_COL_X + LEFT_COL_W + GAP
RIGHT_COL_W = SLIDE_W - M_RIGHT - RIGHT_COL_X

# Section header
HEADER_H = 0.33
CONTENT_TOP = M_TOP_CONTENT + HEADER_H + 0.04

# ========================================================================
# 1. TOP DECORATIVE BARS
# ========================================================================
add_rect(0, 0, SLIDE_W, 0.06, fill_color='1A1A2E')
add_rect(0, SLIDE_H - 0.06, SLIDE_W, 0.06, fill_color='1A1A2E')

# ========================================================================
# 2. TITLE AREA
# ========================================================================
add_text(0.80, 0.28, 11.48, 0.46, '工具生态洞察与 OpenSpec 详解',
         size=28, color='1A1A2E', bold=True)
# Blue accent bar
add_rect(0.80, 0.82, 1.20, 0.05, fill_color='0078D4')
# Subtitle
add_text(0.80, 0.93, 11.48, 0.22,
         'Agentic Coding 工具横向对比  ·  OpenSpec 核心机制与工作流  ·  领域适用性预留',
         size=12, color='555555')

# ========================================================================
# 3. LEFT COLUMN — Agentic Coding 工具洞察
# ========================================================================
# Section header bar
add_rect(LEFT_COL_X, M_TOP_CONTENT, LEFT_COL_W, HEADER_H, fill_color='1A1A2E')
add_text(LEFT_COL_X, M_TOP_CONTENT + 0.02, LEFT_COL_W, HEADER_H - 0.04,
         'Agentic Coding 工具洞察', size=14, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)

# Big container background
container_top = CONTENT_TOP
container_h = SLIDE_H - M_BOTTOM - CONTENT_TOP - 0.06
add_rounded_rect(LEFT_COL_X, container_top, LEFT_COL_W, container_h, fill_color='F5F5F5')

# Tool cards data
tools = [
    {
        'name': 'OpenSpec',
        'accent': '0078D4',
        'subtitle': '开源 · Anthropic 社区 · 规格驱动',
        'points': [
            ('▸ 唯一覆盖完整 SE 生命周期（7 步工作流）', '333333', False),
            ('▸ Delta Spec 只记录变化，不重写全部规格', '333333', False),
            ('▸ 归档自动合并 → 活文档，知识不丢失', '333333', False),
            ('▸ Given/When/Then 验证，Spec-Code 可追溯', '333333', False),
        ]
    },
    {
        'name': 'GitHub Spec Kit',
        'accent': '68217A',
        'subtitle': 'GitHub 官方 · Copilot 生态 · 2025',
        'points': [
            ('▸ 从 Issue 一键生成 Spec，与 PR 直接联动', '333333', False),
            ('▸ 原生 GitHub Actions 驱动，集成度最高', '333333', False),
            ('▸ 轻量快速，适合中小项目快速迭代', '333333', False),
            ('▸ 无增量管理，无跨 artifact 一致性验证', '555555', False),
        ]
    },
    {
        'name': 'Superpowers',
        'accent': '008080',
        'subtitle': 'Claude Code 技能增强 · 结构化工作流',
        'points': [
            ('▸ TDD / 调试 / 代码审查 强制标准化流程', '333333', False),
            ('▸ 过程质量保证，有效减少低级实现失误', '333333', False),
            ('▸ 不管理需求规格，专注代码实现过程', '555555', False),
            ('▸ 可与 OpenSpec 互补：规格 + 实现双保险', '333333', False),
        ]
    },
    {
        'name': 'Everything as Claude Code',
        'accent': 'E86C00',
        'subtitle': '全能 CLI · 对话驱动 · 零配置',
        'points': [
            ('▸ 自然语言直接驱动，零学习成本', '333333', False),
            ('▸ 最大灵活性，极适合快速原型和探索', '333333', False),
            ('▸ 需求随对话消失，无持久化知识管理', '555555', False),
            ('▸ 质量随 Prompt 波动，结果难以稳定复现', '555555', False),
        ]
    },
]

card_inner_x = LEFT_COL_X + 0.06
card_inner_w = LEFT_COL_W - 0.12
card_gap = 0.07
card_h = (container_h - 0.10 - card_gap * 3) / 4
card_y = container_top + 0.05

for tool in tools:
    # Card background
    add_rounded_rect(card_inner_x, card_y, card_inner_w, card_h, fill_color='FFFFFF')
    # Accent bar at top
    add_rect(card_inner_x, card_y, card_inner_w, 0.04, fill_color=tool['accent'])
    # Tool name
    add_text(card_inner_x + 0.04, card_y + 0.08, card_inner_w - 0.08, 0.22,
             tool['name'], size=13, color=tool['accent'], bold=True)
    # Subtitle
    add_text(card_inner_x + 0.04, card_y + 0.28, card_inner_w - 0.08, 0.14,
             tool['subtitle'], size=8, color='999999')
    # Bullet points
    add_multiline(card_inner_x + 0.04, card_y + 0.42, card_inner_w - 0.08, card_h - 0.48,
                  tool['points'], size=9, line_spacing=1.2)
    card_y += card_h + card_gap


# ========================================================================
# 4. RIGHT SECTION — OpenSpec 详细介绍
# ========================================================================
# Section header bar
add_rect(RIGHT_COL_X, M_TOP_CONTENT, RIGHT_COL_W, HEADER_H, fill_color='0078D4')
add_text(RIGHT_COL_X, M_TOP_CONTENT + 0.02, RIGHT_COL_W, HEADER_H - 0.04,
         'OpenSpec 详细介绍', size=14, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)

# Big container background
add_rounded_rect(RIGHT_COL_X, CONTENT_TOP, RIGHT_COL_W, container_h, fill_color='F7F7F7')

# ── 4a. Workflow row (7 steps) ──
wf_top = CONTENT_TOP + 0.10
wf_h = 0.95

# Workflow step data
steps = [
    ('Proposal', '为什么改', '改什么', '5BA0D0'),
    ('Specs', '规格定义', 'Given/When/Then', '509950'),
    ('Design', '技术方案', '架构决策', '509950'),
    ('Tasks', '实现清单', '可追踪进度', 'CC9040'),
    ('Apply', 'AI 按清单', '写代码', 'CC9040'),
    ('Verify', '验证实现', '匹配规格', 'C06060'),
    ('Archive', '合并规格', '归档变更', 'C06060'),
]

n_steps = len(steps)
arrow_w = 0.28
inner_margin = 0.15
available_w = RIGHT_COL_W - 2 * inner_margin
box_w = (available_w - (n_steps - 1) * arrow_w) / n_steps
step_x = RIGHT_COL_X + inner_margin

for i, (name, line1, line2, color) in enumerate(steps):
    # Step box
    add_rounded_rect(step_x, wf_top, box_w, wf_h, fill_color=color)
    # Step name
    add_text(step_x, wf_top + 0.06, box_w, 0.24, name,
             size=12, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
    # Line 1
    add_text(step_x, wf_top + 0.32, box_w, 0.20, line1,
             size=9.5, color='FFFFFF', align=PP_ALIGN.CENTER)
    # Line 2
    add_text(step_x, wf_top + 0.52, box_w, 0.20, line2,
             size=9.5, color='FFFFFF', align=PP_ALIGN.CENTER)

    # Arrow after box (except last)
    if i < n_steps - 1:
        arrow_x = step_x + box_w
        add_text(arrow_x, wf_top + 0.25, arrow_w, 0.35, '→',
                 size=14, color='888888', bold=True, align=PP_ALIGN.CENTER)

    step_x += box_w + arrow_w

# ── Divider line after workflow ──
div_y = wf_top + wf_h + 0.12
add_rect(RIGHT_COL_X + inner_margin, div_y, available_w * 0.6, 0.01, fill_color='DDDDDD')

# ── 4b. 定位 + 四大设计原则 ──
text_y = div_y + 0.10
# 定位
add_text(RIGHT_COL_X + inner_margin, text_y, available_w, 0.18,
         '定位', size=11, color='1A1A2E', bold=True)
add_text(RIGHT_COL_X + inner_margin, text_y + 0.17, available_w, 0.20,
         '轻量级规格驱动框架，让人和 AI 在写代码之前先达成共识，确保工程标准水平的产出',
         size=9, color='555555')

# 四大设计原则
principle_y = text_y + 0.40
add_text(RIGHT_COL_X + inner_margin, principle_y, available_w, 0.18,
         '四大设计原则', size=11, color='1A1A2E', bold=True)

principles = [
    ('✦ 流动而非僵化 — 没有阶段门控，按需推进', '333333', False),
    ('✦ 迭代而非瀑布 — 边做边学，持续优化', '333333', False),
    ('✦ 增量而非全量 — 用 Delta 记录变更，不重写整体', '333333', False),
    ('✦ 存量优先 — 为已有项目设计，不只给新项目', '333333', False),
]
add_multiline(RIGHT_COL_X + inner_margin, principle_y + 0.17, available_w, 0.70,
              principles, size=8.5, line_spacing=1.35)

# ── 4c. Bottom three cards ──
cards_y = principle_y + 0.82
cards_gap = 0.15
card_w = (available_w - 2 * cards_gap) / 3
cards_bottom = CONTENT_TOP + container_h - 0.10
cards_h = cards_bottom - cards_y

cx = RIGHT_COL_X + inner_margin

# ── Card 1: 项目结构 ──
add_rounded_rect(cx, cards_y, card_w, cards_h, fill_color='FFFFFF', line_color='E0E0E0')
add_text(cx + 0.08, cards_y + 0.06, card_w - 0.16, 0.20,
         '项目结构', size=11, color='1A1A2E', bold=True)

tree_lines = [
    ('openspec/', '333333', True),
    ('├── specs/        ← 真相源', '555555', False),
    ('│   └── <domain>/spec.md', '555555', False),
    ('└── changes/      ← 变更', '555555', False),
    ('    ├── proposal.md  (为什么改)', '555555', False),
    ('    ├── usecases.md  (用户场景)', '555555', False),
    ('    ├── specs/        (Delta)', '555555', False),
    ('    ├── design.md    (怎么改)', '555555', False),
    ('    └── tasks.md     (实现清单)', '555555', False),
]
txBox = add_multiline(cx + 0.08, cards_y + 0.28, card_w - 0.16, cards_h - 0.34,
                      tree_lines, size=7, line_spacing=1.25)
# Change font to Consolas for code tree
for para in txBox.text_frame.paragraphs:
    para.font.name = 'Consolas'

# ── Card 2: 核心创新 Delta Spec ──
cx2 = cx + card_w + cards_gap
add_rounded_rect(cx2, cards_y, card_w, cards_h, fill_color='FFFFFF', line_color='E0E0E0')
add_text(cx2 + 0.08, cards_y + 0.06, card_w - 0.16, 0.20,
         '核心创新：Delta Spec', size=11, color='1A1A2E', bold=True)

delta_lines = [
    ('不重写全部，只描述变化：', '333333', False),
    ('', '333333', False),
    ('  ADDED — 新增需求', '34A853', True),
    ('  MODIFIED — 修改需求', 'FBBC04', True),
    ('  REMOVED — 删除需求', 'EA4335', True),
    ('', '333333', False),
    ('每个需求附带 G/W/T 验证', '333333', False),
    ('归档时自动合并 → 活文档', '333333', False),
]
add_multiline(cx2 + 0.08, cards_y + 0.28, card_w - 0.16, cards_h - 0.34,
              delta_lines, size=8.5, line_spacing=1.35)

# ── Card 3: 领域缺陷分析 ──
cx3 = cx2 + card_w + cards_gap
add_rounded_rect(cx3, cards_y, card_w, cards_h, fill_color='FFFFFF', line_color='E0E0E0')
# Orange accent bar at top
add_rect(cx3, cards_y, card_w, 0.04, fill_color='E86C00')
add_text(cx3 + 0.08, cards_y + 0.08, card_w - 0.16, 0.20,
         '领域缺陷分析', size=11, color='E86C00', bold=True)
add_text(cx3 + 0.08, cards_y + 0.28, card_w - 0.16, 0.16,
         'OpenSpec 在我们领域的具体缺陷与改进方向', size=8, color='999999')

domain_items = [
    ('·  华为方法论适配', 'AAAAAA', False),
    ('·  大规模项目性能', 'AAAAAA', False),
    ('·  多团队协作支持', 'AAAAAA', False),
    ('·  工具链集成', 'AAAAAA', False),
    ('·  安全合规要求', 'AAAAAA', False),
]
add_multiline(cx3 + 0.08, cards_y + 0.46, card_w - 0.16, cards_h - 0.56,
              domain_items, size=9, line_spacing=1.4)

add_text(cx3 + 0.08, cards_y + cards_h - 0.22, card_w - 0.16, 0.16,
         '[ 此区域预留供后续填充 ]', size=8, color='CCCCCC', align=PP_ALIGN.CENTER)


# ========================================================================
# SAVE
# ========================================================================
out_name = 'OpenSpec_技术研究汇报_精简版_v2.pptx'
prs.save(out_name)
print(f"✓ Slide 3 rebuilt and saved to {out_name}")
