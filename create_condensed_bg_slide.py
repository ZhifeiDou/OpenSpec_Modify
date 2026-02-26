import sys, os, shutil
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.dml.color import RGBColor
from lxml import etree

src = 'OpenSpec_技术研究汇报_精简版_update.pptx'
if not os.path.exists(src):
    shutil.copy2('OpenSpec_技术研究汇报_精简版.pptx', src)

prs = Presentation(src)
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

blank_layout = None
for layout in prs.slide_layouts:
    if 'blank' in layout.name.lower():
        blank_layout = layout; break
if blank_layout is None:
    blank_layout = prs.slide_layouts[-1]
slide = prs.slides.add_slide(blank_layout)

# Colors
DARK   = RGBColor(0x1A,0x1A,0x2E)
BODY   = RGBColor(0x33,0x33,0x33)
GRAY   = RGBColor(0x88,0x88,0x88)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
RED    = RGBColor(0xD1,0x34,0x38)
GREEN  = RGBColor(0x10,0x7C,0x10)
BLUE   = RGBColor(0x00,0x78,0xD4)
TEAL   = RGBColor(0x00,0x8B,0x8B)
ORANGE = RGBColor(0xE8,0x6C,0x00)

# Layout - 3 columns: 33% Why + 42% What + 25% Placeholder
SW, SH = 13.33, 7.50
LM = 0.20; TM = 0.72; GAP = 0.12
STRIP = 0.26
AVAIL_W = SW - LM*2 - GAP*2
COL1_W = AVAIL_W * 0.33
COL2_W = AVAIL_W * 0.42
COL3_W = AVAIL_W * 0.25
COL_H = SH - TM - 0.12

COL1_X = LM
COL2_X = COL1_X + COL1_W + GAP
COL3_X = COL2_X + COL2_W + GAP

# Chrome bars (top and bottom)
for yy in [0.05, SH-0.05]:
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(yy), Inches(SW), Inches(0.05))
    b.fill.solid(); b.fill.fore_color.rgb = BLUE; b.line.fill.background()

# Slide title
ttx = slide.shapes.add_textbox(Inches(0.30), Inches(0.14), Inches(10), Inches(0.48))
p = ttx.text_frame.paragraphs[0]
r = p.add_run(); r.text = u'背景与框架：Agentic Coding 与 OpenSpec'
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = DARK

def set_radius(shape, v=4000):
    pg = shape._element.find('.//{%s}prstGeom' % NS)
    if pg is not None:
        al = pg.find('{%s}avLst' % NS)
        if al is None: al = etree.SubElement(pg, '{%s}avLst' % NS)
        else:
            for c in list(al): al.remove(c)
        gd = etree.SubElement(al, '{%s}gd' % NS)
        gd.set('name','adj'); gd.set('fmla','val %d' % v)

def add_content_panel(slide, x, y, w, h, title, accent, light, paras):
    """Content panel: title strip + single text box with paragraph-level formatting."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid(); bg.fill.fore_color.rgb = light; bg.line.fill.background()
    set_radius(bg)

    st = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(STRIP))
    st.fill.solid(); st.fill.fore_color.rgb = accent; st.line.fill.background()
    set_radius(st)

    stx = slide.shapes.add_textbox(Inches(x+0.08), Inches(y), Inches(w-0.16), Inches(STRIP))
    stf = stx.text_frame; stf.word_wrap = True
    bp = stf._txBody.find('{%s}bodyPr' % NS)
    if bp is not None: bp.set('anchor','ctr')
    sp = stf.paragraphs[0]
    r = sp.add_run(); r.text = title; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

    pad = 0.06
    cty = y + STRIP + 0.03
    cth = h - STRIP - 0.06
    txb = slide.shapes.add_textbox(Inches(x+pad), Inches(cty), Inches(w-pad*2), Inches(cth))
    tf = txb.text_frame; tf.word_wrap = True
    bp2 = tf._txBody.find('{%s}bodyPr' % NS)
    if bp2 is not None:
        bp2.set('lIns', str(int(Inches(0.04))))
        bp2.set('rIns', str(int(Inches(0.04))))
        bp2.set('tIns', str(int(Inches(0.02))))
        bp2.set('bIns', str(int(Inches(0.02))))
        etree.SubElement(bp2, '{%s}normAutofit' % NS)

    first = True
    for item in paras:
        text, ptype = item[0], item[1]
        color = item[2] if len(item) > 2 else None

        if first:
            pp = tf.paragraphs[0]; first = False
        else:
            pp = tf.add_paragraph()

        if ptype == 'header':
            pp.space_before = Pt(5); pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.bold = True; rr.font.size = Pt(9)
            rr.font.color.rgb = color or accent
        elif ptype == 'warn':
            pp.space_before = Pt(4); pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(10)
            rr.font.bold = True; rr.font.color.rgb = color or RED
        elif ptype == 'code':
            pp.space_before = Pt(0); pp.space_after = Pt(0)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(7.5)
            rr.font.bold = False; rr.font.color.rgb = color or BODY
            rr.font.name = 'Consolas'
        elif ptype == 'empty':
            pp.space_before = Pt(2); pp.space_after = Pt(0)
            rr = pp.add_run(); rr.text = ''; rr.font.size = Pt(4)
        else:  # body
            pp.space_before = Pt(1); pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(9)
            rr.font.bold = False; rr.font.color.rgb = color or BODY

# ═══════ Panel 1: Why Agentic Coding needs OpenSpec ═══════
add_content_panel(slide, COL1_X, TM, COL1_W, COL_H,
    u'为什么需要 OpenSpec', RED, RGBColor(0xFF,0xF0,0xF0), [
    (u'Agentic Coding 现状', 'header', RED),
    (u'Claude Code 等 AI 编程工具已能自主完成', 'body'),
    (u'复杂的软件开发任务。但在企业环境中，', 'body'),
    (u'能写代码只是第一步。', 'body'),
    ('', 'empty'),
    (u'❌ 需求失控', 'header', RED),
    (u'AI 容易自作主张添加/遗漏功能', 'body'),
    (u'需求只存在于聊天记录中，无法追溯', 'body'),
    (u'❌ 质量不稳定', 'header', RED),
    (u'同样需求，不同对话产出差异大', 'body'),
    (u'缺乏验证标准，靠人肉 review', 'body'),
    (u'❌ 协作困难', 'header', RED),
    (u'AI 对话上下文无法共享', 'body'),
    (u'多人并行修改容易冲突', 'body'),
    (u'❌ 知识流失', 'header', RED),
    (u'决策过程随对话消失', 'body'),
    (u'新人接手困难，重复踩坑', 'body'),
    ('', 'empty'),
    (u'⚠ 解决了"怎么写"', 'warn', RED),
    (u'但没解决"写什么"和"写得对不对"', 'warn', RED),
])

# ═══════ Panel 2: What is OpenSpec ═══════
add_content_panel(slide, COL2_X, TM, COL2_W, COL_H,
    u'OpenSpec：规格驱动的人机协同框架', BLUE, RGBColor(0xF0,0xF5,0xFF), [
    (u'定位', 'header', BLUE),
    (u'轻量级规格驱动框架，让人和 AI 在写代码之前先达成共识', 'body'),
    ('', 'empty'),
    (u'四大设计原则', 'header', BLUE),
    (u'✦ 流动而非僵化 — 没有阶段门控，按需推进', 'body'),
    (u'✦ 迭代而非瀑布 — 边做边学，持续优化', 'body'),
    (u'✦ 增量而非全量 — 用 Delta 记录变更，不重写整体', 'body'),
    (u'✦ 存量优先 — 为已有项目设计，不只给新项目', 'body'),
    ('', 'empty'),
    (u'项目结构（两个核心目录）', 'header', DARK),
    (u'openspec/', 'code'),
    (u'├── specs/           ← 当前行为的真相源', 'code'),
    (u'│   └── <domain>/spec.md', 'code'),
    (u'└── changes/         ← 提议的变更', 'code'),
    (u'    ├── proposal.md  （为什么改）', 'code'),
    (u'    ├── usecases.md  （用户场景）', 'code'),
    (u'    ├── specs/       （Delta: 增/改/删）', 'code'),
    (u'    ├── design.md    （怎么改）', 'code'),
    (u'    └── tasks.md     （实现清单）', 'code'),
    ('', 'empty'),
    (u'核心创新：Delta Spec（增量规格）', 'header', GREEN),
    (u'不是重写整个规格，而是只描述变化：', 'body'),
    (u'## ADDED Requirements      ← 新增需求', 'body', GREEN),
    (u'## MODIFIED Requirements  ← 修改需求', 'body', ORANGE),
    (u'## REMOVED Requirements   ← 删除需求', 'body', RED),
    (u'每个需求附带 Given/When/Then 验证场景', 'body'),
    (u'归档时自动合并到主规格 → 活文档', 'body'),
])

# ═══════ Panel 3: Placeholder ═══════
# Title strip
st = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(COL3_X), Inches(TM), Inches(COL3_W), Inches(STRIP))
st.fill.solid(); st.fill.fore_color.rgb = TEAL; st.line.fill.background()
set_radius(st)

stx = slide.shapes.add_textbox(Inches(COL3_X+0.08), Inches(TM), Inches(COL3_W-0.16), Inches(STRIP))
stf = stx.text_frame; stf.word_wrap = True
bp = stf._txBody.find('{%s}bodyPr' % NS)
if bp is not None: bp.set('anchor','ctr')
sp = stf.paragraphs[0]
r = sp.add_run(); r.text = u'更大的背景'; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

# Dashed border box
ph_body_y = TM + STRIP + 0.06
ph_body_h = COL_H - STRIP - 0.12
ph_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(COL3_X + 0.08), Inches(ph_body_y),
    Inches(COL3_W - 0.16), Inches(ph_body_h))
ph_box.fill.solid(); ph_box.fill.fore_color.rgb = RGBColor(0xF0,0xFA,0xFA)
ph_box.line.color.rgb = TEAL
ph_box.line.width = Pt(1.5)
ph_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
set_radius(ph_box)

# Centered placeholder text
ph_txb = slide.shapes.add_textbox(
    Inches(COL3_X + 0.15), Inches(ph_body_y + ph_body_h * 0.25),
    Inches(COL3_W - 0.30), Inches(ph_body_h * 0.5))
tf = ph_txb.text_frame; tf.word_wrap = True
bp2 = tf._txBody.find('{%s}bodyPr' % NS)
if bp2 is not None: bp2.set('anchor', 'ctr')

p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run(); r1.text = u'人与 AI 智能体的'
r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = TEAL

p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = u'工程能力建设'
r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = TEAL

p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(16)
r3 = p3.add_run(); r3.text = u'（讨论内容）'
r3.font.size = Pt(11); r3.font.bold = False; r3.font.color.rgb = GRAY

# ═══════ Reposition: after slide 6 (0-indexed=5) ═══════
sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
new_id = all_ids[-1]
ref_id = all_ids[5]
sldIdLst.remove(new_id)
ref_id.addnext(new_id)

prs.save('OpenSpec_技术研究汇报_精简版_update.pptx')
print('Done! New slide inserted after slide 6.')
