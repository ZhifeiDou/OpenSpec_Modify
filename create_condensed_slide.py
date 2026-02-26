import sys, shutil
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

shutil.copy2('OpenSpec_技术研究汇报_精简版.pptx', 'OpenSpec_技术研究汇报_精简版_update.pptx')
prs = Presentation('OpenSpec_技术研究汇报_精简版_update.pptx')
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

blank_layout = None
for layout in prs.slide_layouts:
    if 'blank' in layout.name.lower():
        blank_layout = layout; break
if blank_layout is None:
    blank_layout = prs.slide_layouts[-1]
slide = prs.slides.add_slide(blank_layout)

# Colors
DARK   = RGBColor(0x1A,0x1A,0x2E)
BODY   = RGBColor(0x33,0x33,0x33)
GRAY   = RGBColor(0x88,0x88,0x88)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
RED    = RGBColor(0xC0,0x39,0x2B)
GREEN  = RGBColor(0x1B,0x7A,0x2B)
TOPBAR = RGBColor(0x00,0x78,0xD4)

DIRS = {
    1: {'main': RGBColor(0x00,0x78,0xD4), 'light': RGBColor(0xF0,0xF5,0xFF)},
    2: {'main': RGBColor(0x10,0x7C,0x10), 'light': RGBColor(0xEF,0xF7,0xEF)},
    3: {'main': RGBColor(0xE8,0x6C,0x00), 'light': RGBColor(0xFF,0xF6,0xE8)},
    4: {'main': RGBColor(0x00,0x8B,0x8B), 'light': RGBColor(0xEB,0xF5,0xF5)},
}

# Layout - maximize space usage
SW, SH = 13.33, 7.50
LM = 0.20; TM = 0.72; GAP = 0.12
CW = (SW - LM*2 - GAP) / 2
CH = (SH - TM - 0.12 - GAP) / 2
STRIP = 0.26

CELLS = {
    1: (LM, TM),
    2: (LM+CW+GAP, TM),
    3: (LM, TM+CH+GAP),
    4: (LM+CW+GAP, TM+CH+GAP),
}

# Chrome
for yy in [0.05, SH-0.05]:
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(yy), Inches(SW), Inches(0.05))
    b.fill.solid(); b.fill.fore_color.rgb = TOPBAR; b.line.fill.background()

ttx = slide.shapes.add_textbox(Inches(0.30), Inches(0.14), Inches(10), Inches(0.48))
p = ttx.text_frame.paragraphs[0]
r = p.add_run(); r.text = u'\u6f5c\u5728\u7814\u7a76\u65b9\u5411\uff1a\u8be6\u7ec6\u65b9\u6848\u4e0e\u5b66\u672f\u4f9d\u636e'
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = DARK

def set_radius(shape, v=4000):
    pg = shape._element.find('.//{%s}prstGeom' % NS)
    if pg is not None:
        al = pg.find('{%s}avLst' % NS)
        if al is None: al = etree.SubElement(pg, '{%s}avLst' % NS)
        else:
            for c in list(al): al.remove(c)
        gd = etree.SubElement(al, '{%s}gd' % NS)
        gd.set('name','adj'); gd.set('fmla','val %d' % v)

def add_cell(slide, dn, title, paras):
    """
    paras: list of (text, type) where type is:
      'h_prob' = problem header, 'h_meth' = method header,
      'h_val' = value header, 'h_ref' = ref header, 'h_path' = path header,
      'body' = normal body, 'ref' = reference text, 'val' = value highlight
    """
    cx, cy = CELLS[dn]; dc = DIRS[dn]

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx), Inches(cy), Inches(CW), Inches(CH))
    bg.fill.solid(); bg.fill.fore_color.rgb = dc['light']; bg.line.fill.background()
    set_radius(bg)

    # Title strip
    st = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx), Inches(cy), Inches(CW), Inches(STRIP))
    st.fill.solid(); st.fill.fore_color.rgb = dc['main']; st.line.fill.background()
    set_radius(st)

    stx = slide.shapes.add_textbox(Inches(cx+0.08), Inches(cy), Inches(CW-0.16), Inches(STRIP))
    stf = stx.text_frame; stf.word_wrap = True
    bp = stf._txBody.find('{%s}bodyPr' % NS)
    if bp is not None: bp.set('anchor','ctr')
    sp = stf.paragraphs[0]
    r = sp.add_run(); r.text = str(dn); r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
    r = sp.add_run(); r.text = '   ' + title; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

    # SINGLE content text box - fills entire remaining space
    pad = 0.06
    cty = cy + STRIP + 0.03
    cth = CH - STRIP - 0.06
    txb = slide.shapes.add_textbox(Inches(cx+pad), Inches(cty), Inches(CW-pad*2), Inches(cth))
    tf = txb.text_frame; tf.word_wrap = True
    bp2 = tf._txBody.find('{%s}bodyPr' % NS)
    if bp2 is not None:
        bp2.set('lIns', str(int(Inches(0.04))))
        bp2.set('rIns', str(int(Inches(0.04))))
        bp2.set('tIns', str(int(Inches(0.02))))
        bp2.set('bIns', str(int(Inches(0.02))))
        etree.SubElement(bp2, '{%s}normAutofit' % NS)

    first = True
    for text, ptype in paras:
        if first:
            pp = tf.paragraphs[0]; first = False
        else:
            pp = tf.add_paragraph()

        if ptype.startswith('h_'):
            pp.space_before = Pt(5)
            pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.bold = True; rr.font.size = Pt(9)
            if ptype == 'h_prob': rr.font.color.rgb = RED
            elif ptype == 'h_meth': rr.font.color.rgb = GREEN
            elif ptype == 'h_val': rr.font.color.rgb = dc['main']
            elif ptype == 'h_ref': rr.font.color.rgb = GRAY
            elif ptype == 'h_path': rr.font.color.rgb = GRAY
        elif ptype == 'val':
            pp.space_before = Pt(1); pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(9)
            rr.font.bold = True; rr.font.color.rgb = dc['main']
        elif ptype == 'ref':
            pp.space_before = Pt(0); pp.space_after = Pt(0)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(8)
            rr.font.bold = False; rr.font.color.rgb = GRAY
        else:  # body
            pp.space_before = Pt(1); pp.space_after = Pt(1)
            rr = pp.add_run(); rr.text = text; rr.font.size = Pt(9)
            rr.font.bold = False; rr.font.color.rgb = BODY


# ======== Direction 1: Multi-Perspective Spec Eval ========
add_cell(slide, 1, u'Multi-Perspective Spec Eval（多智能体多视角规格评估）', [
    (u'现状问题', 'h_prob'),
    (u'当前 Spec 质量检查依赖人工评审或单一模型实例评审，存在系统性盲区', 'body'),
    (u'Perspective-Based Reading (PBR) 研究表明，单一视角审查难以覆盖多维缺陷类别', 'body'),
    (u'方法', 'h_meth'),
    (u'构建多角色 Agent 审查组，在 spec 提交前执行多视角交叉验证：', 'body'),
    (u'• 完整性审查（遗漏场景与逻辑矛盾检测）', 'body'),
    (u'• 异常路径覆盖（错误流程与异常输入生成）', 'body'),
    (u'• 边界条件分析（边界值与竞态条件识别）', 'body'),
    (u'• 集成一致性检查（与已有 main specs 的隐性冲突检测）', 'body'),
    (u'多轮交叉审查后输出经多视角验证的高质量 spec', 'body'),
    (u'核心价值', 'h_val'),
    (u'✦ PBR 多视角团队比单一视角多发现 ~41% 缺陷（Shull, IEEE 2000）', 'val'),
    (u'✦ 可与方向 3 Self-Evolving 形成审查—学习闭环', 'val'),
    (u'学术依据', 'h_ref'),
    (u'Du et al. (ICML\'24) Multi-Agent Debate 提升推理准确性', 'ref'),
    (u'Chan et al. (ICLR\'24) ChatEval 多智能体辩论评估有效性', 'ref'),
    (u'Oriol et al. (RE@Next!\'25) MAD 在需求工程的初步探索', 'ref'),
    (u'实施路径', 'h_path'),
    (u'Agent 角色定义 → 审查协议 → 交叉审查引擎 → Spec 修复建议 → 与 Self-Evolving 闭环', 'ref'),
])

# ======== Direction 2: Spec-as-Runtime Guardrail ========
add_cell(slide, 2, u'Spec-as-Runtime Guardrail（规格即运行时护栏）', [
    (u'现状问题', 'h_prob'),
    (u'传统 Spec-Code 一致性在实现完成后验证，AI 高速生成代码场景下事后验证返工成本高', 'body'),
    (u'将 Spec 从"参考文档"变为"可执行的运行时约束"，缺乏强制约束机制违规无法拦截', 'body'),
    (u'方法', 'h_meth'),
    (u'Design by Contract (Meyer, 1988) 在 AI 辅助编程场景的延伸：', 'body'),
    (u'• Spec→Constraint 编译：结构化需求（EARS 格式）→ property-based predicates', 'body'),
    (u'• Agent 行为实时监控：AI 写代码时实时分析合规性', 'body'),
    (u'• 违规拦截：违反 spec 的代码在 save 前被拦截', 'body'),
    (u'• 自动生成 property-based tests', 'body'),
    (u'核心价值', 'h_val'),
    (u'✦ Kiro (AWS, GA 2025) 已验证 EARS→property test 可行性', 'val'),
    (u'✦ EU AI Act Art.72 要求高风险 AI 系统建立持续合规监控体系', 'val'),
    (u'✦ 安全关键行业（电信/汽车/能源）预防 > 追溯', 'val'),
    (u'学术依据', 'h_ref'),
    (u'Meyer (1988/1992) Design by Contract 理论基础', 'ref'),
    (u'Kiro (AWS, GA 2025) EARS→test 验证 | Tessl、spec-kit 也在探索此方向', 'ref'),
    (u'实施路径', 'h_path'),
    (u'Spec→Constraint 编译器 → Agent 行为监控 → 实时拦截引擎 → Property Test 生成 → CI/CD 集成', 'ref'),
])

# ======== Direction 3: Self-Evolving Spec Agents ========
add_cell(slide, 3, u'Self-Evolving Spec Agents（自我进化规格智能体）', [
    (u'现状问题', 'h_prob'),
    (u'现有 AI 记忆（Copilot Memory、Claude Code auto-memory）仅记录代码模式和用户偏好', 'body'),
    (u'缺乏从 spec 质量到 spec 编写策略的闭环反馈，同类项目反复犯相同错误', 'body'),
    (u'方法', 'h_meth'),
    (u'借鉴 experience-driven self-evolution 框架，引入 spec 质量闭环反馈：', 'body'),
    (u'• 每次 Apply→Verify 循环，记录哪些 spec 写法→高质量实现，哪些→返工', 'body'),
    (u'• 蒸馏领域级 Spec 最佳实践原则（如：风控 spec 必须包含 X/Y/Z）', 'body'),
    (u'• 创建 spec 时系统主动检索并注入相关原则', 'body'),
    (u'核心价值', 'h_val'),
    (u'✦ 经验积累效应：使用越多，spec 质量越高', 'val'),
    (u'✦ 经验沉淀为可复用的组织知识资产', 'val'),
    (u'✦ 与 Domain Packs 差异化互补：Pack = 行业通识，Self-Evolving = 项目经验', 'val'),
    (u'学术依据', 'h_ref'),
    (u'EvolveR (Wu et al., arXiv:2510.16079) LLM Agent 经验蒸馏自进化框架', 'ref'),
    (u'自进化 Agent 综述 (Gao et al., arXiv:2507.21046; Tao et al., arXiv:2508.07407)', 'ref'),
    (u'实施路径', 'h_path'),
    (u'Apply/Verify 日志采集 → 经验提炼引擎 → 原则库建设 → Spec 创建时自动注入 → 效果评估迭代', 'ref'),
])

# ======== Direction 4: Domain-Specific Packs ========
add_cell(slide, 4, u'Domain-Specific Packs（领域规格智能包）', [
    (u'现状问题', 'h_prob'),
    (u'通用大模型不懂行业：电信须符合 3GPP、汽车须通过 ISO 26262、能源须满足 IEC 62109', 'body'),
    (u'生成的 spec 缺乏行业深度与合规性，公司多年行业积累恰好填补空白', 'body'),
    (u'方法', 'h_meth'),
    (u'将行业 Know-How 封装为 Domain Packs（规格模板 + 规则集 + 验证逻辑）：', 'body'),
    (u'• Rules（行业规则）+ Spec 模板（必填字段与标准结构）', 'body'),
    (u'• Schema 约束 + 验证规则（Validators）', 'body'),
    (u'• AI 生成 spec/code 时自动注入行业先验知识', 'body'),
    (u'差异化价值', 'h_val'),
    (u'✦ 行业 Know-How 是稀缺资源，无法通过网页爬取或模型训练获得', 'val'),
    (u'✦ 先发优势 → 累积效应：越用越好 → 越好越多人用 → 竞争壁垒', 'val'),
    (u'落地路线', 'h_path'),
    (u'电信 Pack（Q1，3GPP 协议栈积累最深）→ 汽车（Q2，ASIL-D 认证经验）→ 能源（Q3，并网安全）', 'ref'),
    (u'行业落地', 'h_ref'),
    (u'ICT：3GPP 协议合规 + 网络功能 spec 模板 | 汽车：ASIL 等级模板 + AUTOSAR 约束', 'ref'),
    (u'能源：IEC 62109/62619 并网安全 + 极端工况场景库', 'ref'),
])

# Move slide to after slide 10
sldIdLst = prs.slides._sldIdLst
all_ids = list(sldIdLst)
new_id = all_ids[-1]; ref_id = all_ids[9]
sldIdLst.remove(new_id); ref_id.addnext(new_id)

prs.save('OpenSpec_技术研究汇报_精简版_update.pptx')
print('Done!')
