"""
追加量化项目深度对比分析 Slides 到现有 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (match existing) ──
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_RED_BG = RGBColor(0xFD, 0xE8, 0xE8)
SCORE_GOLD = RGBColor(0xC8, 0x9B, 0x2B)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.join(BASE_DIR, "OpenSpec_技术研究汇报.pptx")
prs = Presentation(pptx_path)
W = prs.slide_width
H = prs.slide_height


# ── Helpers (same as original) ──
def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        spPr = shape._element.spPr
        solidFill_elem = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill_elem is not None:
            srgbClr = solidFill_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr is not None:
                alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=DARK_GRAY, line_spacing=1.5, font_name="Microsoft YaHei",
                      alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, c, fs = item, False, color, font_size
        elif len(item) == 2:
            text, bold = item
            c, fs = color, font_size
        elif len(item) == 3:
            text, bold, c = item
            fs = font_size
        else:
            text, bold, c, fs = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_section_header(slide, text):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 text, font_size=28, color=BLACK, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(1.2), Inches(0.05), ACCENT_BLUE)


def add_page_number(slide, num, total):
    add_text_box(slide, W - Inches(1.2), H - Inches(0.5),
                 Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=ACCENT_BLUE, bg_color=WHITE, border_color=None):
    card = add_rect(slide, left, top, width, height, bg_color)
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    add_multiline_box(slide, left + Inches(0.2), top + Inches(0.6),
                      width - Inches(0.4), height - Inches(0.8),
                      body_lines, font_size=13, color=DARK_GRAY, line_spacing=1.4)
    return card


# ── Total slide count (original 15 + 6 new = 21) ──
EXISTING = 15
NEW_SLIDES = 6
TOTAL = EXISTING + NEW_SLIDES

# ════════════════════════════════════════════════════════════════════════
# APPENDIX DIVIDER (Slide 16)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)
add_rect(slide, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
             "附录：量化交易系统深度技术对比",
             font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.6),
             "A股有色金属多因子量化交易系统  |  同一需求 · 两种实现 · 逐项拆解",
             font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_multiline_box(slide, Inches(2.5), Inches(4.8), Inches(8), Inches(1.5), [
    ("OpenSpec 版：30+ 源码文件 · 16 个因子 · 100+ 测试用例 · CLI 工具", False, RGBColor(0x88, 0xBB, 0xEE), 15),
    ("无 OpenSpec 版：6 个源码文件 · 15+ 因子 · 0 测试 · Flask Web 仪表盘", False, RGBColor(0xEE, 0xBB, 0x88), 15),
], line_spacing=2.0, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 16, TOTAL)

# ════════════════════════════════════════════════════════════════════════
# Slide 17: 架构对比 — 一目了然
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "架构对比：企业级分层 vs 快速原型")

# ── Left: OpenSpec version ──
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.1), Inches(0.55), ACCENT_BLUE)
add_text_box(slide, Inches(0.4), Inches(1.47), Inches(6.1), Inches(0.4),
             "OpenSpec 版 — 6层架构 · 30+ 文件", font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

layers_openspec = [
    ("src/data/", "数据层", "多数据源适配器 (AKShare + BaoStock 自动切换)\nSQLite 增量更新 + 数据验证器", ACCENT_BLUE),
    ("src/factors/", "因子层", "注册器模式，16 个因子通过 @register 自注册\nMAD 去极值 + Z-Score 标准化", ACCENT_GREEN),
    ("src/strategy/", "策略层", "多因子评分 + IC 加权 + 约束优化分配\n商品动量择时（铜铝联动信号）", ACCENT_ORANGE),
    ("src/risk/", "风控层", "ATR 动态硬止损 + 追踪止损 + 分层回撤管理\n金属暴跌联动预警", ACCENT_RED),
    ("src/backtest/", "回测层", "事件驱动引擎 + A股规则模拟器\nT+1 / 涨跌停 / 停牌 / 整手 / 印花税", DARK_GRAY),
    ("src/report/", "报告层", "Plotly 交互图表 + HTML/PNG 双格式输出", MID_GRAY),
]

for i, (path, name, desc, color) in enumerate(layers_openspec):
    y = Inches(2.1) + Inches(0.85) * i
    add_rect(slide, Inches(0.4), y, Inches(1.35), Inches(0.72), color)
    add_text_box(slide, Inches(0.4), y + Inches(0.02), Inches(1.35), Inches(0.35),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.4), y + Inches(0.35), Inches(1.35), Inches(0.3),
                 path, font_size=9, color=RGBColor(0xDD, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER,
                 font_name="Consolas")
    add_multiline_box(slide, Inches(1.85), y + Inches(0.05), Inches(4.55), Inches(0.65),
                      desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.2)

# ── Right: No-OpenSpec version ──
add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(0.55), ACCENT_ORANGE)
add_text_box(slide, Inches(6.8), Inches(1.47), Inches(6.1), Inches(0.4),
             "无 OpenSpec 版 — 扁平结构 · 6 文件", font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

layers_no = [
    ("config.py", "配置", "Python 常量 · 硬编码股票池 · 无外部配置文件", ACCENT_ORANGE),
    ("data_fetcher.py", "数据", "单一 AKShare 源 · TTL 缓存 · 无数据验证", ACCENT_ORANGE),
    ("factor_engine.py", "因子", "所有因子挤在一个文件 · Z-Score 标准化\n无注册机制，添加因子需改多处代码", ACCENT_ORANGE),
    ("strategy.py", "策略", "阈值判断生成信号 · 等权分配\n无择时、无约束优化", ACCENT_ORANGE),
    ("backtester.py", "回测", "简化回测 · 因子计算与实时信号不一致\n无涨跌停/停牌/T+1 限制模拟", ACCENT_ORANGE),
    ("app.py", "界面", "Flask Web 仪表盘 · ECharts 图表\n暗色主题 · 参数在线调整  (亮点)", ACCENT_GREEN),
]

for i, (path, name, desc, color) in enumerate(layers_no):
    y = Inches(2.1) + Inches(0.85) * i
    add_rect(slide, Inches(6.8), y, Inches(1.35), Inches(0.72), color)
    add_text_box(slide, Inches(6.8), y + Inches(0.02), Inches(1.35), Inches(0.35),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(6.8), y + Inches(0.35), Inches(1.35), Inches(0.3),
                 path, font_size=9, color=RGBColor(0xDD, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER,
                 font_name="Consolas")
    add_multiline_box(slide, Inches(8.25), y + Inches(0.05), Inches(4.55), Inches(0.65),
                      desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.2)

# Bottom insight
add_rect(slide, Inches(0.4), Inches(7.3) - Inches(0.3), Inches(12.5), Inches(0.5), LIGHT_BLUE_BG)
add_text_box(slide, Inches(0.6), Inches(7.3) - Inches(0.25), Inches(12.1), Inches(0.4),
             "OpenSpec 的 specs 和 design 工件在写代码前就定义了模块边界，AI 被迫按契约分层实现，而非一股脑堆在一起",
             font_size=13, color=ACCENT_BLUE, bold=True)

add_page_number(slide, 17, TOTAL)

# ════════════════════════════════════════════════════════════════════════
# Slide 18: 风控系统对比 — 生死线
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "风控对比：量化交易的生死线")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "在量化交易中，风控能力决定了系统能否在真实市场中存活。这是两个版本差距最大的地方。",
             font_size=15, color=DARK_GRAY)

# ── Left: OpenSpec risk system ──
add_rect(slide, Inches(0.4), Inches(2.2), Inches(6.1), Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(0.4), Inches(2.25), Inches(6.1), Inches(0.4),
             "OpenSpec 版：4 层风控体系", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

risk_openspec = [
    ("第1层：ATR 动态止损", ACCENT_RED,
     "根据近期波动率自动调整止损位\n波动大 → 止损宽（避免误杀）\n波动小 → 止损窄（保护利润）"),
    ("第2层：追踪止损", ACCENT_ORANGE,
     "股价上涨 10% 后激活\n从最高点回撤 8% 则卖出\n锁定盈利，让利润奔跑"),
    ("第3层：组合回撤管理", ACCENT_BLUE,
     "回撤 15% → 仓位砍半（减少暴露）\n回撤 20% → 全部清仓（保命）\n分层响应，不会一刀切"),
    ("第4层：金属暴跌联动", DARK_GRAY,
     "监控铜/铝期货日跌幅 > 3%\n自动识别受影响的持仓股\n触发黄金对冲机制"),
]

for i, (title, color, desc) in enumerate(risk_openspec):
    y = Inches(2.85) + Inches(1.1) * i
    add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.9), color)
    add_text_box(slide, Inches(0.6), y, Inches(2.5), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_multiline_box(slide, Inches(0.6), y + Inches(0.35), Inches(5.7), Inches(0.55),
                      desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.1)

# ── Right: No-OpenSpec risk system ──
add_rect(slide, Inches(6.8), Inches(2.2), Inches(6.1), Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(6.8), Inches(2.25), Inches(6.1), Inches(0.4),
             "无 OpenSpec 版：基础止损", font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

risk_no = [
    ("固定止损 -8%", ACCENT_ORANGE,
     "不管市场波动大小，统一 -8% 触发\n牛市中容易被正常回调误杀\n熊市中 -8% 可能仍然太慢"),
    ("回撤预警 -15%", MID_GRAY,
     "仅发出预警，没有自动减仓动作\n需要人工介入处理\n晚间/周末无法及时响应"),
]

for i, (title, color, desc) in enumerate(risk_no):
    y = Inches(2.85) + Inches(1.1) * i
    add_rect(slide, Inches(6.8), y, Inches(0.08), Inches(0.9), color)
    add_text_box(slide, Inches(7.0), y, Inches(2.5), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_multiline_box(slide, Inches(7.0), y + Inches(0.35), Inches(5.7), Inches(0.55),
                      desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.1)

# Gap visualization
add_rect(slide, Inches(6.8), Inches(5.05), Inches(6.1), Inches(1.8), LIGHT_RED_BG)
add_text_box(slide, Inches(7.0), Inches(5.15), Inches(5.7), Inches(0.35),
             "缺失的风控能力", font_size=14, color=ACCENT_RED, bold=True)
add_multiline_box(slide, Inches(7.0), Inches(5.55), Inches(5.7), Inches(1.2), [
    ("  无动态止损 — 无法适应不同市场环境", False, DARK_GRAY, 11),
    ("  无追踪止损 — 赚到的利润无法锁定", False, DARK_GRAY, 11),
    ("  无自动减仓 — 回撤时只能眼睁睁看着", False, DARK_GRAY, 11),
    ("  无联动预警 — 上游商品崩盘时毫无防备", False, DARK_GRAY, 11),
], line_spacing=1.3)

# Bottom insight
add_rect(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.45), LIGHT_BLUE_BG)
add_text_box(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35),
             "OpenSpec 的 usecases.md 和 specs 工件要求明确定义异常场景，迫使 AI 实现完整的风控逻辑，而非只处理正常情况",
             font_size=12, color=ACCENT_BLUE, bold=True)

add_page_number(slide, 18, TOTAL)

# ════════════════════════════════════════════════════════════════════════
# Slide 19: 回测引擎对比 — 结果可信度
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "回测对比：你的回测结果能信多少？")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "回测是量化策略上线前的最后一道关卡。回测越贴近真实，上线后的意外就越少。",
             font_size=15, color=DARK_GRAY)

# Table header
th_y = Inches(2.2)
add_rect(slide, Inches(0.6), th_y, Inches(3.8), Inches(0.55), MID_GRAY)
add_text_box(slide, Inches(0.6), th_y + Inches(0.08), Inches(3.8), Inches(0.4),
             "A股交易规则", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.4), th_y, Inches(4.3), Inches(0.55), ACCENT_BLUE)
add_text_box(slide, Inches(4.4), th_y + Inches(0.08), Inches(4.3), Inches(0.4),
             "OpenSpec 版", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(8.7), th_y, Inches(4.2), Inches(0.55), ACCENT_ORANGE)
add_text_box(slide, Inches(8.7), th_y + Inches(0.08), Inches(4.2), Inches(0.4),
             "无 OpenSpec 版", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ("T+1 交割制度",
     "完整模拟：今日买入明日才能卖",
     "未实现：买卖当日即可完成"),
    ("涨跌停板限制",
     "涨停 +10% 不可买入，跌停不可卖出",
     "未实现：任意价格均可成交"),
    ("停牌处理",
     "成交量=0 自动冻结，不可交易",
     "未实现：停牌期间仍可交易"),
    ("整手交易（100股）",
     "强制 100 股倍数买入",
     "未实现：可买任意股数"),
    ("交易费用建模",
     "印花税 0.05% + 佣金 0.03%\n+ 滑点 0.15%（分别建模）",
     "合并费率 0.025%\n（欠精确）"),
    ("因子计算一致性",
     "回测与实盘使用同一因子引擎",
     "回测中重新实现了简化版因子\n与实时信号逻辑不一致"),
    ("仓位分配",
     "得分加权 + 单股 10% + 行业 25%\n上限约束 + 择时调仓比例",
     "等权分配\n无约束优化"),
]

for i, (rule, openspec_val, no_val) in enumerate(rows):
    y = th_y + Inches(0.55) + Inches(0.60) * i
    bgc = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.6), y, Inches(3.8), Inches(0.60), bgc)
    add_rect(slide, Inches(4.4), y, Inches(4.3), Inches(0.60), bgc)
    add_rect(slide, Inches(8.7), y, Inches(4.2), Inches(0.60), bgc)

    add_text_box(slide, Inches(0.75), y + Inches(0.08), Inches(3.5), Inches(0.44),
                 rule, font_size=11, color=BLACK, bold=True)
    add_multiline_box(slide, Inches(4.55), y + Inches(0.04), Inches(4.0), Inches(0.52),
                      openspec_val.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.1)
    add_multiline_box(slide, Inches(8.85), y + Inches(0.04), Inches(3.9), Inches(0.52),
                      no_val.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.1)

# Bottom insight
add_rect(slide, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.5), LIGHT_ORANGE_BG)
add_multiline_box(slide, Inches(0.6), Inches(6.93), Inches(12.1), Inches(0.4), [
    ("无 OpenSpec 版的回测忽略了大量 A 股特有规则，回测收益率会虚高，上线后会因为 T+1 无法卖出、涨停买不进等问题导致实际表现大幅偏离", False, ACCENT_ORANGE, 12),
], line_spacing=1.0)

add_page_number(slide, 19, TOTAL)

# ════════════════════════════════════════════════════════════════════════
# Slide 20: 综合评分雷达图（文字模拟）
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "综合评分：8 个维度逐项打分")

# Score table
dimensions = [
    ("架构设计",    9, 7, "分层清晰，模块间通过接口通信", "扁平直接，够用但难以扩展"),
    ("风控完备性",  9, 5, "4 层风控 + 联动预警", "仅固定止损 + 预警"),
    ("回测真实性",  9, 6, "完整模拟 A 股交易规则", "简化模拟，结果可信度低"),
    ("测试覆盖",    8, 0, "100+ 测试用例，pytest 框架", "零测试"),
    ("可扩展性",    9, 5, "注册器模式，新增因子只需一个类", "添加因子需改多处代码"),
    ("配置管理",    9, 6, "外部 YAML，12 类参数可调", "Python 常量硬编码"),
    ("用户体验",    6, 8, "CLI + HTML 报告", "Web 仪表盘 + 在线调参"),
    ("文档质量",    8, 7, "README + 类型注解 + BDD 场景", "中文注释 + docstrings"),
]

# Header
header_y = Inches(1.5)
col_starts = [Inches(0.5), Inches(2.6), Inches(4.0), Inches(5.4), Inches(7.0), Inches(10.7)]
col_widths = [Inches(2.0), Inches(1.3), Inches(1.3), Inches(1.5), Inches(3.6), Inches(2.2)]
headers = ["维度", "OpenSpec", "无OpenSpec", "差距", "OpenSpec 亮点", "无OpenSpec 亮点"]

add_rect(slide, Inches(0.5), header_y, Inches(12.4), Inches(0.5), DARK_BLUE)
for j, h in enumerate(headers):
    add_text_box(slide, col_starts[j] + Inches(0.05), header_y + Inches(0.06),
                 col_widths[j] - Inches(0.1), Inches(0.38),
                 h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Rows
for i, (dim, s1, s2, desc1, desc2) in enumerate(dimensions):
    y = header_y + Inches(0.5) + Inches(0.58) * i
    bgc = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.5), y, Inches(12.4), Inches(0.58), bgc)

    add_text_box(slide, col_starts[0] + Inches(0.1), y + Inches(0.1),
                 col_widths[0] - Inches(0.2), Inches(0.38),
                 dim, font_size=12, color=BLACK, bold=True)

    # OpenSpec score (colored)
    s1_color = ACCENT_GREEN if s1 >= 8 else (ACCENT_ORANGE if s1 >= 6 else ACCENT_RED)
    add_text_box(slide, col_starts[1], y + Inches(0.1),
                 col_widths[1], Inches(0.38),
                 f"{s1}/10", font_size=13, color=s1_color, bold=True, alignment=PP_ALIGN.CENTER)

    # No-OpenSpec score (colored)
    s2_color = ACCENT_GREEN if s2 >= 8 else (ACCENT_ORANGE if s2 >= 6 else ACCENT_RED)
    add_text_box(slide, col_starts[2], y + Inches(0.1),
                 col_widths[2], Inches(0.38),
                 f"{s2}/10", font_size=13, color=s2_color, bold=True, alignment=PP_ALIGN.CENTER)

    # Gap
    gap = s1 - s2
    gap_text = f"+{gap}" if gap > 0 else str(gap)
    gap_color = ACCENT_GREEN if gap > 0 else (ACCENT_RED if gap < 0 else MID_GRAY)
    add_text_box(slide, col_starts[3], y + Inches(0.1),
                 col_widths[3], Inches(0.38),
                 gap_text, font_size=13, color=gap_color, bold=True, alignment=PP_ALIGN.CENTER)

    # Descriptions
    add_text_box(slide, col_starts[4] + Inches(0.05), y + Inches(0.1),
                 col_widths[4] - Inches(0.1), Inches(0.38),
                 desc1, font_size=10, color=DARK_GRAY)
    add_text_box(slide, col_starts[5] + Inches(0.05), y + Inches(0.1),
                 col_widths[5] - Inches(0.1), Inches(0.38),
                 desc2, font_size=10, color=DARK_GRAY)

# Total scores
total_y = header_y + Inches(0.5) + Inches(0.58) * 8 + Inches(0.1)
add_rect(slide, Inches(0.5), total_y, Inches(12.4), Inches(0.55), DARK_BLUE)
add_text_box(slide, col_starts[0] + Inches(0.1), total_y + Inches(0.08),
             col_widths[0], Inches(0.38),
             "总分", font_size=14, color=WHITE, bold=True)

total_1 = sum(s1 for _, s1, _, _, _ in dimensions)
total_2 = sum(s2 for _, _, s2, _, _ in dimensions)
add_text_box(slide, col_starts[1], total_y + Inches(0.08),
             col_widths[1], Inches(0.38),
             f"{total_1}/80", font_size=14, color=RGBColor(0x66, 0xFF, 0x66), bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, col_starts[2], total_y + Inches(0.08),
             col_widths[2], Inches(0.38),
             f"{total_2}/80", font_size=14, color=RGBColor(0xFF, 0xCC, 0x66), bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, col_starts[3], total_y + Inches(0.08),
             col_widths[3], Inches(0.38),
             f"+{total_1 - total_2}", font_size=14, color=RGBColor(0x66, 0xFF, 0x66), bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.0), total_y + Inches(0.08),
             Inches(5.8), Inches(0.38),
             f"OpenSpec 版总分领先 {total_1 - total_2} 分，唯一弱项是缺少 Web 界面",
             font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

add_page_number(slide, 20, TOTAL)

# ════════════════════════════════════════════════════════════════════════
# Slide 21: 结论 — 为什么 OpenSpec 版更好
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "结论：OpenSpec 让 AI 写出了工程级别的量化系统")

# Main verdict
add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0), LIGHT_BLUE_BG)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.35),
             "OpenSpec 版在 8 个维度中 7 个领先，总分 67/80 vs 44/80",
             font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.35),
             "不是代码写得多就好 — 而是规格约束让 AI 在正确的方向上写出高质量代码",
             font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 3 key takeaways
takeaways = [
    ("规格前置 = 架构质量",
     ACCENT_BLUE, LIGHT_BLUE_BG,
     [
         "OpenSpec 强制 AI 在写代码前定义模块",
         "边界、接口契约和数据流。",
         "",
         "结果：6 层清晰架构，每层职责单一，",
         "模块间通过 Protocol 接口通信。",
         "",
         "对比：无 OpenSpec 时 AI 把所有逻辑",
         "堆在 6 个文件里，改一处要查全局。",
     ]),
    ("异常场景 = 风控完备",
     ACCENT_RED, LIGHT_RED_BG,
     [
         "OpenSpec 的 usecases.md 要求定义异常",
         "流程，specs 要求 Given/When/Then 验证。",
         "",
         "结果：4 层风控体系，覆盖个股止损、",
         "组合回撤、品种联动等完整场景。",
         "",
         "对比：无 OpenSpec 时 AI 只实现正常流程",
         "的 happy path，风控形同虚设。",
     ]),
    ("任务清单 = 测试护城河",
     ACCENT_GREEN, LIGHT_GREEN_BG,
     [
         "OpenSpec 的 tasks.md 把测试作为显式",
         "任务项，AI 必须逐项完成。",
         "",
         "结果：100+ 测试用例，覆盖因子计算、",
         "策略逻辑、风控触发、回测指标等。",
         "",
         "对比：无 OpenSpec 时 AI 认为「功能写完",
         "就是完成」，测试数量为 0。",
     ]),
]

for i, (title, tc, bgc, lines) in enumerate(takeaways):
    x = Inches(0.4) + Inches(4.2) * i
    y = Inches(2.8)
    card = add_rect(slide, x, y, Inches(4.0), Inches(4.3), bgc)
    card.line.color.rgb = tc
    card.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.6), Inches(0.45),
                 title, font_size=16, color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(1.0), y + Inches(0.7), Inches(2.0), Inches(0.03), tc)
    add_multiline_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.6), Inches(3.2),
                      lines, font_size=11, color=DARK_GRAY, line_spacing=1.2)

add_page_number(slide, 21, TOTAL)

# ── Save ──
prs.save(pptx_path)
print(f"PPT updated: added {NEW_SLIDES} slides (total {TOTAL})")
print(f"Saved to: {pptx_path}")
