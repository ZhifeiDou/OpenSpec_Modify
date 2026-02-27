# -*- coding: utf-8 -*-
"""Insert ONE combined slide with all 3 experiment evaluation tables."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报_精简版.pptx"
OUTPUT_PATH = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报_精简版_v2.pptx"

FONT = "Microsoft YaHei"
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_COLOR = RGBColor(0x44, 0x44, 0x44)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_BAR_COLOR = RGBColor(0x00, 0x78, 0xD4)
HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
ROW_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TOTAL_BG = RGBColor(0xE3, 0xF2, 0xFD)
SUMMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = 12190730


def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def add_top_bar(slide):
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W + 965, Emu(54864))
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_accent_underline(slide):
    shape = slide.shapes.add_shape(
        1, Emu(731520), Emu(910000), Emu(1097280), Emu(45720)
    )
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def insert_slide_after(prs, index):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_list = prs.slides._sldIdLst
    slide_elem = slide_list[-1]
    slide_list.remove(slide_elem)
    slide_list.insert(index + 1, slide_elem)
    return slide


def score_color(val, is_total=False):
    """Return color based on score value."""
    if isinstance(val, str):
        try:
            val = float(val.split('/')[0])
        except:
            return BODY_COLOR
    if val >= 4:
        return GREEN
    elif val >= 3:
        return ORANGE
    elif val >= 2:
        return RGBColor(0xE8, 0x6C, 0x00)
    else:
        return RED


def draw_table(slide, x_start, y_start, table_title, table_subtitle,
               col_headers, col_widths, rows_data, total_row,
               title_bg_color):
    """
    Draw a styled table with colored cells.
    col_headers: list of (text, bg_color)
    col_widths: list of EMU widths
    rows_data: list of [dim_label, score1, score2, ...] or [dim_label, score1, score2, score3]
    total_row: [label, total1, total2, ...]
    """
    total_w = sum(col_widths)
    row_h = Emu(295000)
    header_h = Emu(310000)
    total_h = Emu(320000)

    # Table title bar
    title_h = Emu(290000)
    add_rect(slide, x_start, y_start, total_w, title_h, title_bg_color)
    add_textbox(slide, x_start, y_start + Emu(10000), total_w, title_h,
                table_title, Pt(12), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    subtitle_h = Emu(230000)
    y = y_start + title_h
    add_rect(slide, x_start, y, total_w, subtitle_h, RGBColor(0xF0, 0xF0, 0xF0))
    add_textbox(slide, x_start, y + Emu(5000), total_w, subtitle_h,
                table_subtitle, Pt(8), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # Column headers
    y = y + subtitle_h
    cx = x_start
    for i, (hdr_text, hdr_bg) in enumerate(col_headers):
        add_rect(slide, cx, y, col_widths[i], header_h, hdr_bg)
        add_textbox(slide, cx, y + Emu(15000), col_widths[i], header_h,
                    hdr_text, Pt(9), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        cx += col_widths[i]

    # Data rows
    y = y + header_h
    for ri, row in enumerate(rows_data):
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE
        cx = x_start
        for ci, cell_text in enumerate(row):
            add_rect(slide, cx, y, col_widths[ci], row_h, bg)
            if ci == 0:
                # Dimension label - left aligned
                add_textbox(slide, cx + Emu(40000), y + Emu(8000),
                            col_widths[ci] - Emu(50000), row_h,
                            cell_text, Pt(9), TITLE_COLOR, bold=False,
                            alignment=PP_ALIGN.LEFT)
            else:
                # Score value - centered, colored
                txt_color = score_color(cell_text)
                add_textbox(slide, cx, y + Emu(5000), col_widths[ci], row_h,
                            str(cell_text), Pt(13), txt_color, bold=True,
                            alignment=PP_ALIGN.CENTER)
            cx += col_widths[ci]
        y += row_h

    # Total row
    cx = x_start
    for ci, cell_text in enumerate(total_row):
        add_rect(slide, cx, y, col_widths[ci], total_h, TOTAL_BG)
        if ci == 0:
            add_textbox(slide, cx + Emu(40000), y + Emu(10000),
                        col_widths[ci] - Emu(50000), total_h,
                        cell_text, Pt(10), TITLE_COLOR, bold=True,
                        alignment=PP_ALIGN.LEFT)
        else:
            txt_color = score_color(cell_text, is_total=True)
            add_textbox(slide, cx, y + Emu(5000), col_widths[ci], total_h,
                        str(cell_text), Pt(14), txt_color, bold=True,
                        alignment=PP_ALIGN.CENTER)
        cx += col_widths[ci]

    return y + total_h  # return bottom y


def build_combined_slide(prs, insert_after_idx):
    slide = insert_slide_after(prs, insert_after_idx)
    add_top_bar(slide)

    # Title
    add_textbox(slide, Emu(731520), Emu(280000), Emu(10500000), Emu(460000),
                u"三组实验综合评分：五维度 SE 质量对比",
                Pt(26), TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Subtitle
    add_textbox(slide, Emu(731520), Emu(960000), Emu(10500000), Emu(250000),
                u"A股量化初始构建（3方） | A股量化需求变更（2方） | Dashboard 应用（3方），独立评估五个 SE 维度",
                Pt(11), SUBTITLE_COLOR)

    # =====================================================================
    # Table layout: 3 tables side by side
    # =====================================================================
    gap = Emu(100000)
    margin_left = Emu(250000)
    # Total available: ~11700000
    # Exp2 has 3 cols (narrower), Exp1 and Dashboard have 4 cols (wider)
    # Exp1: 3900000, Exp2: 3400000, Dashboard: 3900000
    # Check: 3900 + 100 + 3400 + 100 + 3900 = 11400. Plus margin 250*2 = 500 -> fits in 11900

    # --- Table 1: Experiment 1 (initial build, 3 projects) ---
    t1_x = margin_left
    t1_w_dim = Emu(1350000)
    t1_w_score = Emu(850000)
    t1_total_w = t1_w_dim + t1_w_score * 3  # 3900000
    t1_col_widths = [t1_w_dim, t1_w_score, t1_w_score, t1_w_score]

    t1_headers = [
        (u"评估维度", HEADER_BG),
        (u"Vanilla", RGBColor(0xC6, 0x28, 0x28)),
        (u"OpenSpec", RGBColor(0xE8, 0x6C, 0x00)),
        (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32)),
    ]
    t1_rows = [
        [u"1. 需求与规格", "1", "4", "5"],
        [u"2. 设计", "1", "4", "5"],
        [u"3. 实现", "2", "4", "4"],
        [u"4. 验证与确认", "1", "3", "4"],
        [u"5. 维护与演进", "1", "4", "5"],
    ]
    t1_total = [u"总分", "6/25", "19/25", "23/25"]

    # --- Table 2: Experiment 2 (requirement change, 2 projects) ---
    t2_x = t1_x + t1_total_w + gap
    t2_w_dim = Emu(1350000)
    t2_w_score = Emu(1075000)
    t2_total_w = t2_w_dim + t2_w_score * 2  # 3500000
    t2_col_widths = [t2_w_dim, t2_w_score, t2_w_score]

    t2_headers = [
        (u"评估维度", HEADER_BG),
        (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32)),
        (u"Vanilla", RGBColor(0xC6, 0x28, 0x28)),
    ]
    t2_rows = [
        [u"1. 需求与规格", "5", "1"],
        [u"2. 设计", "5", "1"],
        [u"3. 实现", "4", "2"],
        [u"4. 验证与确认", "4", "1"],
        [u"5. 维护与演进", "5", "1"],
    ]
    t2_total = [u"总分", "23/25", "6/25"]

    # --- Table 3: Dashboard (3 projects) ---
    t3_x = t2_x + t2_total_w + gap
    t3_w_dim = Emu(1350000)
    t3_w_score = Emu(850000)
    t3_total_w = t3_w_dim + t3_w_score * 3  # 3900000
    t3_col_widths = [t3_w_dim, t3_w_score, t3_w_score, t3_w_score]

    t3_headers = [
        (u"评估维度", HEADER_BG),
        (u"Vanilla", RGBColor(0xC6, 0x28, 0x28)),
        (u"OpenSpec", RGBColor(0xE8, 0x6C, 0x00)),
        (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32)),
    ]
    t3_rows = [
        [u"1. 需求与规格", "1", "4", "5"],
        [u"2. 设计", "2", "4", "5"],
        [u"3. 实现", "2", "4", "4"],
        [u"4. 验证与确认", "1", "3", "5"],
        [u"5. 维护与演进", "2", "4", "5"],
    ]
    t3_total = [u"总分", "8/25", "19/25", "24/25"]

    # Draw all three tables
    table_y = Emu(1280000)

    draw_table(slide, t1_x, table_y,
               u"实验1：A股量化 — 初始构建",
               u"同一需求，3 种方法从零构建",
               t1_headers, t1_col_widths, t1_rows, t1_total,
               RGBColor(0x33, 0x33, 0x66))

    draw_table(slide, t2_x, table_y,
               u"实验2：A股量化 — 需求变更",
               u"同一基线，3 项需求变更迭代",
               t2_headers, t2_col_widths, t2_rows, t2_total,
               RGBColor(0x33, 0x66, 0x33))

    draw_table(slide, t3_x, table_y,
               u"实验3：Dashboard — 初始构建",
               u"Web Dashboard，3 种方法对比",
               t3_headers, t3_col_widths, t3_rows, t3_total,
               RGBColor(0x66, 0x33, 0x33))

    # =====================================================================
    # Below tables: score highlight boxes (3 big numbers)
    # =====================================================================
    # Table bottom is approximately at:
    # 1280000 + 290000(title) + 230000(subtitle) + 310000(header) + 5*295000(rows) + 320000(total)
    # = 1280000 + 290000 + 230000 + 310000 + 1475000 + 320000 = 3905000
    box_y = Emu(4000000)
    box_h = Emu(1050000)

    # --- Experiment 1 highlight ---
    add_rect(slide, t1_x, box_y, t1_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    # Three score circles/numbers
    score_data_1 = [
        ("6", RED, u"Vanilla"),
        ("19", ORANGE, u"OpenSpec"),
        ("23", GREEN, u"OS-HW"),
    ]
    sx = t1_x
    sw = Emu(int(t1_total_w / 3))
    for score_val, score_col, score_label in score_data_1:
        add_textbox(slide, sx, box_y + Emu(30000), sw, Emu(580000),
                    score_val, Pt(32), score_col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(630000), sw, Emu(300000),
                    u"/ 25 " + score_label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # --- Experiment 2 highlight ---
    add_rect(slide, t2_x, box_y, t2_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    score_data_2 = [
        ("23", GREEN, u"OS-HW"),
        ("6", RED, u"Vanilla"),
    ]
    sx = t2_x
    sw = Emu(int(t2_total_w / 2))
    for score_val, score_col, score_label in score_data_2:
        add_textbox(slide, sx, box_y + Emu(30000), sw, Emu(580000),
                    score_val, Pt(32), score_col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(630000), sw, Emu(300000),
                    u"/ 25 " + score_label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # --- Experiment 3 highlight ---
    add_rect(slide, t3_x, box_y, t3_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    score_data_3 = [
        ("8", RED, u"Vanilla"),
        ("19", ORANGE, u"OpenSpec"),
        ("24", GREEN, u"OS-HW"),
    ]
    sx = t3_x
    sw = Emu(int(t3_total_w / 3))
    for score_val, score_col, score_label in score_data_3:
        add_textbox(slide, sx, box_y + Emu(30000), sw, Emu(580000),
                    score_val, Pt(32), score_col, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(630000), sw, Emu(300000),
                    u"/ 25 " + score_label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # =====================================================================
    # Bottom summary bar
    # =====================================================================
    summary_y = Emu(5200000)
    summary_h = Emu(750000)
    add_rect(slide, Emu(250000), summary_y, Emu(11500000), summary_h,
             RGBColor(0xE8, 0xF5, 0xE9))

    # Line 1: main conclusion
    txBox = slide.shapes.add_textbox(Emu(400000), summary_y + Emu(30000),
                                     Emu(11200000), Emu(350000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = u"核心结论：三组独立实验一致证明，OpenSpec-HW 的规格驱动工作流在所有 SE 维度系统性领先 Vanilla Claude Code。"
    run.font.size = Pt(12)
    run.font.color.rgb = SUMMARY_COLOR
    run.font.bold = True
    run.font.name = FONT

    # Line 2: supporting detail
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    segs = [
        (u"差距最大维度：", Pt(11), SUMMARY_COLOR, True),
        (u"需求规格 & 维护演进（5:1）", Pt(11), BLUE, True),
        (u"   差距最小维度：", Pt(11), SUMMARY_COLOR, True),
        (u"实现（4:2）", Pt(11), ORANGE, True),
        (u"   Vanilla Claude 三组实验平均 ", Pt(11), SUBTITLE_COLOR, False),
        (u"7.3/25", Pt(11), RED, True),
        (u"，OpenSpec-HW 平均 ", Pt(11), SUBTITLE_COLOR, False),
        (u"23.3/25", Pt(11), GREEN, True),
    ]
    for text, size, color, bold in segs:
        run = p2.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT

    # Line 3
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = u"结论适用于不同项目类型（量化交易 vs Web Dashboard）和不同开发阶段（初始构建 vs 需求变更迭代），具有可复现性。"
    run3.font.size = Pt(10)
    run3.font.color.rgb = SUBTITLE_COLOR
    run3.font.bold = False
    run3.font.name = FONT

    return slide


def main():
    prs = Presentation(PPTX_PATH)
    original = len(prs.slides)
    print(f"Original slide count: {original}")

    # Insert after slide 11 (index 10) - after all exp2 detail slides
    build_combined_slide(prs, 10)

    print(f"New slide count: {len(prs.slides)}")
    prs.save(OUTPUT_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
