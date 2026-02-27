# -*- coding: utf-8 -*-
"""Insert SE evaluation slides into the PPTX after slide 4."""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报_精简版.pptx"
OUTPUT_PATH = PPTX_PATH  # overwrite

# --- Style constants (extracted from existing slides) ---
FONT = "Microsoft YaHei"
TITLE_SIZE = Pt(28)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_SIZE = Pt(16)
SUBTITLE_COLOR = RGBColor(0x44, 0x44, 0x44)
BODY_SIZE = Pt(13)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)
LABEL_SIZE = Pt(11)
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BAR_COLOR = RGBColor(0x00, 0x78, 0xD4)
HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
ROW_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
HIGHLIGHT_RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
PAGE_COLOR = RGBColor(0x88, 0x88, 0x88)
BOTTOM_BG = RGBColor(0x1A, 0x1A, 0x2E)
SUMMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = 12190730
SLIDE_H = 6858000


def set_shape_fill(shape, color):
    """Set solid fill on a shape."""
    spPr = shape._element.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr'
    )
    if spPr is None:
        spPr = shape._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr'
        )
    if spPr is None:
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        spPr = etree.SubElement(shape._element,
                                '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size, color, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name=FONT):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if isinstance(line_text, tuple):
            # (text, color, bold)
            run = p.add_run()
            run.text = line_text[0]
            run.font.size = font_size
            run.font.color.rgb = line_text[2] if len(line_text) > 2 else color
            run.font.bold = line_text[1] if len(line_text) > 1 else bold
            run.font.name = font_name
        else:
            run = p.add_run()
            run.text = line_text
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font_name
    return txBox


def add_top_bar(slide):
    """Add the thin colored bar at the very top."""
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W + 965, Emu(54864)
    )
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_accent_underline(slide):
    """Add accent underline below title."""
    shape = slide.shapes.add_shape(
        1, Emu(731520), Emu(960120), Emu(1097280), Emu(45720)
    )
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_page_number(slide, text):
    """Add page number in bottom-right."""
    add_textbox(slide, Emu(11094415), Emu(6400800), Emu(914400), Emu(365760),
                text, Pt(10), PAGE_COLOR, alignment=PP_ALIGN.RIGHT)


def add_rect_cell(slide, left, top, width, height, fill_color):
    """Add a rectangle that acts as a table cell background."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def insert_slide_after(prs, index):
    """Insert a blank slide after the given index, return the new slide."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    # Move the new slide to the correct position
    slide_list = prs.slides._sldIdLst
    slide_elem = slide_list[-1]  # last added
    slide_list.remove(slide_elem)
    slide_list.insert(index + 1, slide_elem)
    return slide


def build_slide_1(prs, insert_after_idx):
    """Slide: experiment 1 overview with score matrix."""
    slide = insert_slide_after(prs, insert_after_idx)

    add_top_bar(slide)
    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验1 - 初始构建实验：三方软件工程质量五维度评估",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Subtitle
    add_textbox(slide, Emu(731520), Emu(1080000), Emu(10058400), Emu(340000),
                u"同一量化交易需求，分别用 Vanilla Claude Code、原生 OpenSpec、OpenSpec-HW 构建，评估五个软件工程维度",
                Pt(15), SUBTITLE_COLOR)

    # === Score matrix as grid of rectangles ===
    # Column positions
    col_starts = [Emu(457200), Emu(2600000), Emu(5100000), Emu(7600000), Emu(10100000)]
    col_widths = [Emu(2100000), Emu(2460000), Emu(2460000), Emu(2460000), Emu(1500000)]
    row_h = Emu(370000)
    row_start = Emu(1550000)

    headers = [u"评估维度", u"A: Vanilla Claude", u"B: 原生 OpenSpec", u"C: OpenSpec-HW", u"差距"]

    # Header row
    for ci, hdr in enumerate(headers):
        add_rect_cell(slide, col_starts[ci], row_start, col_widths[ci], row_h, HEADER_BG)
        add_textbox(slide, col_starts[ci], row_start + Emu(15000),
                    col_widths[ci], row_h,
                    hdr, Pt(13), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Data rows
    rows_data = [
        [u"1. 需求与规格", "1", "4", "5", u"A:C = 1:5"],
        [u"2. 设计", "1", "4", "5", u"A:C = 1:5"],
        [u"3. 实现", "2", "4", "4", u"A:C = 1:2"],
        [u"4. 验证与确认", "1", "3", "4", u"A:C = 1:4"],
        [u"5. 维护与演进", "1", "4", "5", u"A:C = 1:5"],
        [u"总分 (满分25)", "6", "19", "23", u""],
        [u"平均分 (满分5)", "1.2", "3.8", "4.6", u"A:C = 1:3.8"],
    ]

    for ri, row_data in enumerate(rows_data):
        y = row_start + row_h * (ri + 1)
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE
        is_total = ri >= 5
        for ci, cell_text in enumerate(row_data):
            fill = bg
            if is_total:
                fill = RGBColor(0xE3, 0xF2, 0xFD)  # light blue for total
            add_rect_cell(slide, col_starts[ci], y, col_widths[ci], row_h, fill)

            # Color coding for score columns
            text_color = LABEL_COLOR
            text_bold = False
            if ci == 1:  # No OpenSpec - red
                text_color = RED
                text_bold = is_total
            elif ci == 2:  # Native OpenSpec
                text_color = ORANGE
                text_bold = is_total
            elif ci == 3:  # OpenSpec-HW - green
                text_color = GREEN
                text_bold = is_total
            elif ci == 4:  # Gap
                text_color = BLUE
                text_bold = True
            elif ci == 0:
                text_bold = is_total

            add_textbox(slide, col_starts[ci], y + Emu(15000),
                        col_widths[ci], row_h,
                        cell_text, Pt(12) if not is_total else Pt(13),
                        text_color, bold=text_bold,
                        alignment=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # === Bottom summary box ===
    box_y = Emu(4650000)
    add_rect_cell(slide, Emu(365760), box_y, Emu(11430000), Emu(580000), RGBColor(0xE8, 0xF5, 0xE9))

    add_multiline_textbox(
        slide, Emu(548640), box_y + Emu(30000), Emu(11064240), Emu(520000),
        [
            (u"核心结论：OpenSpec-HW 的规格驱动工作流在所有维度全面领先 Vanilla Claude Code（23 vs 6），差距最大在需求/设计/维护（5:1）", True, SUMMARY_COLOR),
        ],
        Pt(13), SUMMARY_COLOR
    )

    # === Three project boxes at bottom ===
    box_w = Emu(3580000)
    box_h = Emu(1250000)
    box_y2 = Emu(5400000)
    box_gap = Emu(100000)
    box_x = [Emu(365760), Emu(365760) + box_w + box_gap, Emu(365760) + (box_w + box_gap) * 2]

    proj_data = [
        (u"Project A: Vanilla Claude Code", u"6 个 Python 文件，0 测试\n0 规格文档，0 设计文档\n总分：6/25", RED, RGBColor(0xFF, 0xEB, 0xEE)),
        (u"Project B: 原生 OpenSpec", u"28 个 Python 文件，58 个测试\n16 个规格文档（600行）\n总分：19/25", ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
        (u"Project C: OpenSpec-HW", u"27 个 Python 文件，56 个测试\n24 个规格文档（1634行）\n总分：23/25", GREEN, HIGHLIGHT_GREEN_BG),
    ]

    for i, (title, desc, title_color, bg_color) in enumerate(proj_data):
        add_rect_cell(slide, box_x[i], box_y2, box_w, box_h, bg_color)
        add_textbox(slide, box_x[i] + Emu(80000), box_y2 + Emu(50000),
                    box_w - Emu(160000), Emu(280000),
                    title, Pt(12), title_color, bold=True)
        add_multiline_textbox(
            slide, box_x[i] + Emu(80000), box_y2 + Emu(350000),
            box_w - Emu(160000), Emu(800000),
            desc.split('\n'), Pt(10), BODY_COLOR
        )

    return slide


def build_slide_2(prs, insert_after_idx):
    """Slide: dimension-by-dimension comparison grid (like slide 6 style)."""
    slide = insert_slide_after(prs, insert_after_idx)

    add_top_bar(slide)
    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验1 - 初始构建结果：五维度逐项对比",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Grid layout like slide 6
    col_starts = [Emu(457200), Emu(3017520), Emu(6858000), Emu(10698480)]
    col_widths = [Emu(2560320), Emu(3840480), Emu(3840480), Emu(1097280)]
    row_h = Emu(420000)
    row_start = Emu(1180000)

    # Header row
    header_data = [
        (u"对比维度", WHITE, HEADER_BG),
        (u"OpenSpec-HW（C）", WHITE, RGBColor(0x2E, 0x7D, 0x32)),
        (u"Vanilla Claude Code（A）", WHITE, RGBColor(0xC6, 0x28, 0x28)),
        (u"评分", WHITE, HEADER_BG),
    ]
    for ci, (text, text_color, bg_color) in enumerate(header_data):
        add_rect_cell(slide, col_starts[ci], row_start, col_widths[ci], row_h, bg_color)
        add_textbox(slide, col_starts[ci], row_start + Emu(30000),
                    col_widths[ci], row_h,
                    text, Pt(14), text_color, bold=True, alignment=PP_ALIGN.CENTER)

    # Comparison rows
    rows = [
        {
            "label": u"需求与规格",
            "good": [u"8 个规格文档 + 用例文档", u"38 需求 + 56 场景 (WHEN/THEN)"],
            "bad": [u"0 个规格文档", u"无需求、无用例、无 API 规格"],
            "score": u"5 vs 1",
        },
        {
            "label": u"设计",
            "good": [u"254 行设计文档 + 8 个架构决策", u"Protocol/ABC/Adapter 设计模式"],
            "bad": [u"0 行设计文档", u"无 OOP、无抽象、逻辑重复"],
            "score": u"5 vs 1",
        },
        {
            "label": u"实现",
            "good": [u"3909 行 / 7 个包 / type hints", u"YAML 配置 + CLI + 优雅降级"],
            "bad": [u"1200 行 / 6 个平铺文件", u"无 type hints、print 调试、代码重复"],
            "score": u"4 vs 2",
        },
        {
            "label": u"验证确认",
            "good": [u"56 个测试 + 集成测试（779 行）", u"边界/异常覆盖 + conftest 复用"],
            "bad": [u"0 个测试", u"无 pytest / CI / linting"],
            "score": u"4 vs 1",
        },
        {
            "label": u"维护演进",
            "good": [u"README + 变更归档 + 活文档", u"插件架构 + Protocol 扩展点"],
            "bad": [u"无 README / 无文档", u"无扩展点、硬编码配置、无日志"],
            "score": u"5 vs 1",
        },
    ]

    for ri, row in enumerate(rows):
        y = row_start + row_h * (ri + 1)
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE

        for ci in range(4):
            add_rect_cell(slide, col_starts[ci], y, col_widths[ci], row_h, bg)

        # Column 0: label
        add_textbox(slide, col_starts[0] + Emu(80000), y + Emu(10000),
                    col_widths[0] - Emu(80000), row_h,
                    row["label"], Pt(11), LABEL_COLOR, bold=True)

        # Column 1: OpenSpec-HW (good)
        add_multiline_textbox(
            slide, col_starts[1] + Emu(80000), y + Emu(5000),
            col_widths[1] - Emu(100000), row_h,
            row["good"], Pt(10), GREEN
        )

        # Column 2: No OpenSpec (bad)
        add_multiline_textbox(
            slide, col_starts[2] + Emu(80000), y + Emu(5000),
            col_widths[2] - Emu(100000), row_h,
            row["bad"], Pt(10), RED
        )

        # Column 3: score
        add_textbox(slide, col_starts[3], y + Emu(30000),
                    col_widths[3], row_h,
                    row["score"], Pt(14), BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Summary box at bottom
    summary_y = row_start + row_h * 6 + Emu(80000)
    add_rect_cell(slide, Emu(365760), summary_y, Emu(11430000), Emu(600000),
                  RGBColor(0xE8, 0xF5, 0xE9))

    add_multiline_textbox(
        slide, Emu(548640), summary_y + Emu(20000), Emu(11064240), Emu(550000),
        [
            (u"关键发现：OpenSpec 的规格驱动流程使 AI 在编码前先建立需求共识，从而在设计、测试、维护各环节产生系统性提升。", True, SUMMARY_COLOR),
            (u"无 OpenSpec 时，AI 跳过了需求分析直接写代码，导致零测试、零文档、架构混乱的「能跑但不可维护」的结果。", False, SUBTITLE_COLOR),
        ],
        Pt(13), SUMMARY_COLOR
    )

    return slide


def build_slide_3(prs, insert_after_idx):
    """Slide: key findings with visual callouts."""
    slide = insert_slide_after(prs, insert_after_idx)

    add_top_bar(slide)
    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验1 - 关键发现与分析",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Five finding boxes arranged in a grid (2 cols x 3 rows)
    findings = [
        {
            "num": "01",
            "title": u"规格驱动 = 全维度提升",
            "desc": u"OpenSpec 项目在所有 5 个维度均显著领先。结构化的 proposal→specs→design→tasks 工作流确保需求在编码前被明确捕获。",
            "color": BLUE,
        },
        {
            "num": "02",
            "title": u"有 Spec 才有 Test",
            "desc": u"OpenSpec-HW 的 WHEN/THEN 场景直接映射为测试用例（56 个），实现需求→测试的可追溯性。无 Spec 的项目测试为 0。",
            "color": GREEN,
        },
        {
            "num": "03",
            "title": u"设计文档驱动架构质量",
            "desc": u"有设计文档的项目采用 Protocol/ABC/Adapter 等模式，模块化清晰。无设计文档的项目出现逻辑重复、紧耦合。",
            "color": ORANGE,
        },
        {
            "num": "04",
            "title": u"维护差距最为显著",
            "desc": u"OpenSpec-HW 的变更归档、活文档、插件架构使系统可持续演进。无 OpenSpec 的项目任何修改都面临回归风险。",
            "color": RGBColor(0x9C, 0x27, 0xB0),
        },
        {
            "num": "05",
            "title": u"实现质量差距相对最小",
            "desc": u"三个项目都能产出可运行的量化交易系统，说明 AI 编码能力本身已足够。差距在于「写对的代码」而非「能写代码」。",
            "color": RED,
        },
    ]

    # Layout: 3 boxes on top row, 2 boxes on bottom row
    box_w = Emu(3500000)
    box_h = Emu(1700000)
    gap = Emu(130000)
    start_x = Emu(457200)
    row1_y = Emu(1200000)
    row2_y = row1_y + box_h + gap

    positions = [
        (start_x, row1_y),
        (start_x + box_w + gap, row1_y),
        (start_x + (box_w + gap) * 2, row1_y),
        (start_x + Emu(900000), row2_y),
        (start_x + Emu(900000) + box_w + gap, row2_y),
    ]

    for i, finding in enumerate(findings):
        x, y = positions[i]
        add_rect_cell(slide, x, y, box_w, box_h, RGBColor(0xFA, 0xFA, 0xFA))

        # Number badge
        badge_size = Emu(400000)
        add_rect_cell(slide, x + Emu(80000), y + Emu(80000), badge_size, badge_size, finding["color"])
        add_textbox(slide, x + Emu(80000), y + Emu(100000), badge_size, badge_size,
                    finding["num"], Pt(18), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, x + Emu(550000), y + Emu(100000), box_w - Emu(630000), Emu(350000),
                    finding["title"], Pt(14), finding["color"], bold=True)

        # Description
        add_textbox(slide, x + Emu(100000), y + Emu(550000), box_w - Emu(200000), Emu(1100000),
                    finding["desc"], Pt(11), BODY_COLOR)

    # Bottom summary
    add_rect_cell(slide, Emu(365760), Emu(5100000), Emu(11430000), Emu(500000),
                  RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, Emu(548640), Emu(5130000), Emu(11064240), Emu(440000),
                u"总结：OpenSpec 解决的不是「AI 能不能写代码」的问题，而是「AI 写的代码是否符合工程标准」的问题。从 6/25 到 23/25 的提升证明了规格驱动开发的系统性价值。",
                Pt(13), SUMMARY_COLOR, bold=True)

    return slide


def main():
    prs = Presentation(PPTX_PATH)
    print(f"Original slide count: {len(prs.slides)}")

    # Insert 3 slides after slide 4 (index 3)
    # Insert in reverse order so indices stay correct
    s3 = build_slide_3(prs, 3)
    s2 = build_slide_2(prs, 3)
    s1 = build_slide_1(prs, 3)

    print(f"New slide count: {len(prs.slides)}")
    prs.save(OUTPUT_PATH)
    print(f"Saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
