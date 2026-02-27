# -*- coding: utf-8 -*-
"""Replace slide 12 with enhanced version that includes experiment descriptions."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

PPTX_PATH = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报_精简版.pptx"
OUTPUT_PATH = PPTX_PATH

FONT = "Microsoft YaHei"
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_COLOR = RGBColor(0x44, 0x44, 0x44)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_BAR_COLOR = RGBColor(0x00, 0x78, 0xD4)
HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
ROW_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TOTAL_BG = RGBColor(0xE3, 0xF2, 0xFD)
SUMMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
DESC_COLOR = RGBColor(0x55, 0x55, 0x55)

SLIDE_W = 12190730


def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size, color, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if isinstance(line_text, tuple):
            txt, sz, clr, bld = line_text
            run = p.add_run()
            run.text = txt
            run.font.size = sz
            run.font.color.rgb = clr
            run.font.bold = bld
            run.font.name = font_name
        else:
            run = p.add_run()
            run.text = line_text
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def add_top_bar(slide):
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W + 965, Emu(54864))
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_accent_underline(slide):
    slide.shapes.add_shape(
        1, Emu(731520), Emu(910000), Emu(1097280), Emu(45720)
    )
    set_shape_fill(slide.shapes[-1], ACCENT_BAR_COLOR)
    slide.shapes[-1].line.fill.background()


def score_color(val):
    if isinstance(val, str):
        try:
            val = float(val.split('/')[0])
        except:
            return BODY_COLOR
    if val >= 4:
        return GREEN
    elif val >= 3:
        return ORANGE
    else:
        return RED


def delete_slide(prs, index):
    """Delete slide at given index."""
    rId = prs.slides._sldIdLst[index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[index]
    prs.slides._sldIdLst.remove(sldId)


def insert_slide_after(prs, index):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_list = prs.slides._sldIdLst
    slide_elem = slide_list[-1]
    slide_list.remove(slide_elem)
    slide_list.insert(index + 1, slide_elem)
    return slide


def draw_table(slide, x_start, y_start, table_title, desc_lines,
               col_headers, col_widths, rows_data, total_row,
               title_bg_color):
    """Draw a table with title, description, headers, data rows, and total."""
    total_w = sum(col_widths)
    row_h = Emu(270000)
    header_h = Emu(285000)
    total_h = Emu(295000)

    # --- Table title bar ---
    title_h = Emu(270000)
    add_rect(slide, x_start, y_start, total_w, title_h, title_bg_color)
    add_textbox(slide, x_start, y_start + Emu(5000), total_w, title_h,
                table_title, Pt(11), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # --- Description area ---
    desc_h = Emu(380000)
    y = y_start + title_h
    add_rect(slide, x_start, y, total_w, desc_h, RGBColor(0xF8, 0xF8, 0xF0))
    add_multiline_textbox(
        slide, x_start + Emu(50000), y + Emu(10000),
        total_w - Emu(100000), desc_h,
        desc_lines, Pt(8), DESC_COLOR, alignment=PP_ALIGN.LEFT
    )

    # --- Column headers ---
    y = y + desc_h
    cx = x_start
    for i, (hdr_text, hdr_bg) in enumerate(col_headers):
        add_rect(slide, cx, y, col_widths[i], header_h, hdr_bg)
        add_textbox(slide, cx, y + Emu(12000), col_widths[i], header_h,
                    hdr_text, Pt(9), WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        cx += col_widths[i]

    # --- Data rows ---
    y = y + header_h
    for ri, row in enumerate(rows_data):
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE
        cx = x_start
        for ci, cell_text in enumerate(row):
            add_rect(slide, cx, y, col_widths[ci], row_h, bg)
            if ci == 0:
                add_textbox(slide, cx + Emu(30000), y + Emu(5000),
                            col_widths[ci] - Emu(40000), row_h,
                            cell_text, Pt(8), TITLE_COLOR,
                            alignment=PP_ALIGN.LEFT)
            else:
                add_textbox(slide, cx, y + Emu(2000), col_widths[ci], row_h,
                            str(cell_text), Pt(12), score_color(cell_text),
                            bold=True, alignment=PP_ALIGN.CENTER)
            cx += col_widths[ci]
        y += row_h

    # --- Total row ---
    cx = x_start
    for ci, cell_text in enumerate(total_row):
        add_rect(slide, cx, y, col_widths[ci], total_h, TOTAL_BG)
        if ci == 0:
            add_textbox(slide, cx + Emu(30000), y + Emu(8000),
                        col_widths[ci] - Emu(40000), total_h,
                        cell_text, Pt(9), TITLE_COLOR, bold=True,
                        alignment=PP_ALIGN.LEFT)
        else:
            add_textbox(slide, cx, y + Emu(2000), col_widths[ci], total_h,
                        str(cell_text), Pt(13), score_color(cell_text),
                        bold=True, alignment=PP_ALIGN.CENTER)
        cx += col_widths[ci]

    return y + total_h


def build_combined_slide(prs, insert_after_idx):
    slide = insert_slide_after(prs, insert_after_idx)
    add_top_bar(slide)

    # Title
    add_textbox(slide, Emu(731520), Emu(260000), Emu(10500000), Emu(440000),
                u"三组实验综合评分：五维度 SE 质量对比",
                Pt(26), TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Subtitle
    add_textbox(slide, Emu(731520), Emu(960000), Emu(10500000), Emu(230000),
                u"两个项目类型（量化交易 / Web Dashboard）× 两个开发阶段（初始构建 / 需求变更）× 五个 SE 维度",
                Pt(11), SUBTITLE_COLOR)

    # =====================================================================
    # Table dimensions
    # =====================================================================
    gap = Emu(90000)
    margin_left = Emu(230000)

    t1_w_dim = Emu(1300000)
    t1_w_sc = Emu(830000)
    t1_total_w = t1_w_dim + t1_w_sc * 3  # 3790000

    t2_w_dim = Emu(1300000)
    t2_w_sc = Emu(1050000)
    t2_total_w = t2_w_dim + t2_w_sc * 2  # 3400000

    t3_w_dim = Emu(1300000)
    t3_w_sc = Emu(830000)
    t3_total_w = t3_w_dim + t3_w_sc * 3  # 3790000

    t1_x = margin_left
    t2_x = t1_x + t1_total_w + gap
    t3_x = t2_x + t2_total_w + gap

    table_y = Emu(1250000)

    # =====================================================================
    # Table 1: Experiment 1 - A-share quant initial build
    # =====================================================================
    t1_desc = [
        u"实验内容：从零构建 A 股有色金属量化交易机器人",
        u"含数据获取、多因子引擎、策略评分、风控止损、",
        u"事件驱动回测、报告生成 6 大模块（~27 个源文件）",
    ]
    t1_bottom = draw_table(
        slide, t1_x, table_y,
        u"实验1：A股量化 — 初始构建",
        t1_desc,
        [(u"评估维度", HEADER_BG), (u"Claude Code", RGBColor(0xC6, 0x28, 0x28)),
         (u"OpenSpec", RGBColor(0xE8, 0x6C, 0x00)), (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32))],
        [t1_w_dim, t1_w_sc, t1_w_sc, t1_w_sc],
        [
            [u"1. 需求与规格", "1", "4", "5"],
            [u"2. 设计", "1", "4", "5"],
            [u"3. 实现", "2", "4", "4"],
            [u"4. 验证与确认", "1", "3", "4"],
            [u"5. 维护与演进", "1", "4", "5"],
        ],
        [u"总分", "6/25", "19/25", "23/25"],
        RGBColor(0x33, 0x33, 0x66),
    )

    # =====================================================================
    # Table 2: Experiment 2 - A-share quant requirement change
    # =====================================================================
    t2_desc = [
        u"实验内容：在实验1基线上执行 3 项需求变更",
        u"① 波动率自适应止损（ATR 动态乘数）",
        u"② 动态再平衡频率 ③ 极端行情熔断机制（新功能）",
    ]
    t2_bottom = draw_table(
        slide, t2_x, table_y,
        u"实验2：A股量化 — 需求变更",
        t2_desc,
        [(u"评估维度", HEADER_BG), (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32)),
         (u"Claude Code", RGBColor(0xC6, 0x28, 0x28))],
        [t2_w_dim, t2_w_sc, t2_w_sc],
        [
            [u"1. 需求与规格", "5", "1"],
            [u"2. 设计", "5", "1"],
            [u"3. 实现", "4", "2"],
            [u"4. 验证与确认", "4", "1"],
            [u"5. 维护与演进", "5", "1"],
        ],
        [u"总分", "23/25", "6/25"],
        RGBColor(0x33, 0x66, 0x33),
    )

    # =====================================================================
    # Table 3: Dashboard experiment
    # =====================================================================
    t3_desc = [
        u"实验内容：从零构建 Web Dashboard 应用",
        u"含前端交互界面、后端 REST API、数据可视化、",
        u"用户认证等功能模块，验证结论跨项目类型通用性",
    ]
    t3_bottom = draw_table(
        slide, t3_x, table_y,
        u"实验3：Dashboard — 初始构建",
        t3_desc,
        [(u"评估维度", HEADER_BG), (u"Claude Code", RGBColor(0xC6, 0x28, 0x28)),
         (u"OpenSpec", RGBColor(0xE8, 0x6C, 0x00)), (u"OS-HW", RGBColor(0x2E, 0x7D, 0x32))],
        [t3_w_dim, t3_w_sc, t3_w_sc, t3_w_sc],
        [
            [u"1. 需求与规格", "1", "4", "5"],
            [u"2. 设计", "2", "4", "5"],
            [u"3. 实现", "2", "4", "4"],
            [u"4. 验证与确认", "1", "3", "5"],
            [u"5. 维护与演进", "2", "4", "5"],
        ],
        [u"总分", "8/25", "19/25", "24/25"],
        RGBColor(0x66, 0x33, 0x33),
    )

    # =====================================================================
    # Score highlight boxes
    # =====================================================================
    box_y = max(t1_bottom, t2_bottom, t3_bottom) + Emu(60000)
    box_h = Emu(900000)

    # --- Exp 1 ---
    add_rect(slide, t1_x, box_y, t1_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    scores_1 = [("6", RED, u"Claude Code"), ("19", ORANGE, u"OpenSpec"), ("23", GREEN, u"OS-HW")]
    sw = Emu(int(t1_total_w / 3))
    sx = t1_x
    for val, clr, label in scores_1:
        add_textbox(slide, sx, box_y + Emu(20000), sw, Emu(500000),
                    val, Pt(30), clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(530000), sw, Emu(280000),
                    u"/ 25  " + label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # --- Exp 2 ---
    add_rect(slide, t2_x, box_y, t2_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    scores_2 = [("23", GREEN, u"OS-HW"), ("6", RED, u"Claude Code")]
    sw = Emu(int(t2_total_w / 2))
    sx = t2_x
    for val, clr, label in scores_2:
        add_textbox(slide, sx, box_y + Emu(20000), sw, Emu(500000),
                    val, Pt(30), clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(530000), sw, Emu(280000),
                    u"/ 25  " + label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # --- Exp 3 ---
    add_rect(slide, t3_x, box_y, t3_total_w, box_h, RGBColor(0xF5, 0xF5, 0xF5))
    scores_3 = [("8", RED, u"Claude Code"), ("19", ORANGE, u"OpenSpec"), ("24", GREEN, u"OS-HW")]
    sw = Emu(int(t3_total_w / 3))
    sx = t3_x
    for val, clr, label in scores_3:
        add_textbox(slide, sx, box_y + Emu(20000), sw, Emu(500000),
                    val, Pt(30), clr, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, box_y + Emu(530000), sw, Emu(280000),
                    u"/ 25  " + label, Pt(9), SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
        sx += sw

    # =====================================================================
    # Bottom summary bar
    # =====================================================================
    summary_y = box_y + box_h + Emu(60000)
    summary_h = Emu(650000)
    add_rect(slide, Emu(230000), summary_y, Emu(11500000), summary_h,
             RGBColor(0xE8, 0xF5, 0xE9))

    txBox = slide.shapes.add_textbox(Emu(380000), summary_y + Emu(25000),
                                     Emu(11200000), Emu(620000))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Line 1
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = u"核心结论：三组独立实验一致证明，OpenSpec-HW 在所有 SE 维度系统性领先 Claude Code（平均 23.3 vs 6.7 / 25）。"
    run.font.size = Pt(11)
    run.font.color.rgb = SUMMARY_COLOR
    run.font.bold = True
    run.font.name = FONT

    # Line 2
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    segs = [
        (u"差距最大：", Pt(10), SUMMARY_COLOR, True),
        (u"需求规格 & 维护演进（5:1）", Pt(10), BLUE, True),
        (u"   差距最小：", Pt(10), SUMMARY_COLOR, True),
        (u"实现（4:2）", Pt(10), ORANGE, True),
        (u"   结论跨项目类型（量化交易 / Dashboard）与开发阶段（初始构建 / 需求变更）均成立，具有可复现性。", Pt(10), SUBTITLE_COLOR, False),
    ]
    for text, size, color, bold in segs:
        run = p2.add_run()
        run.text = text
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT

    return slide


def main():
    prs = Presentation(PPTX_PATH)
    original = len(prs.slides)
    print(f"Original slide count: {original}")

    # Insert after slide 11 (index 10) -> becomes slide 12
    build_combined_slide(prs, 10)
    print(f"After insert: {len(prs.slides)}")

    prs.save(OUTPUT_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
