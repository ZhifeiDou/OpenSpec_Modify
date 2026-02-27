# -*- coding: utf-8 -*-
"""Insert Experiment 2 (requirements change) SE evaluation slides after slide 8."""

from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报_精简版.pptx"
OUTPUT_PATH = PPTX_PATH  # overwrite

# --- Style constants (match existing slides exactly) ---
FONT = "Microsoft YaHei"
TITLE_SIZE = Pt(28)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_SIZE = Pt(16)
SUBTITLE_COLOR = RGBColor(0x44, 0x44, 0x44)
BODY_SIZE = Pt(13)
BODY_COLOR = RGBColor(0x44, 0x44, 0x44)
LABEL_SIZE = Pt(11)
LABEL_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BAR_COLOR = RGBColor(0x00, 0x78, 0xD4)
HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
ROW_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
PAGE_COLOR = RGBColor(0x88, 0x88, 0x88)
SUMMARY_COLOR = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = 12190730
SLIDE_H = 6858000


def set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size, color, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name=FONT,
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if line_spacing is not None:
            p.space_after = Pt(line_spacing)
        if isinstance(line_text, tuple):
            run = p.add_run()
            run.text = line_text[0]
            run.font.size = font_size
            run.font.color.rgb = line_text[2] if len(line_text) > 2 else color
            run.font.bold = line_text[1] if len(line_text) > 1 else bold
            run.font.name = font_name
        else:
            run = p.add_run()
            run.text = line_text
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font_name
    return txBox


def add_rich_paragraph(tf, segments, alignment=PP_ALIGN.LEFT, space_after=None):
    """Add a paragraph with multiple styled runs.
    segments: list of (text, font_size, color, bold) tuples.
    """
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_after is not None:
        p.space_after = Pt(space_after)
    for text, font_size, color, bold in segments:
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT
    return p


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W + 965, Emu(54864)
    )
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_accent_underline(slide):
    shape = slide.shapes.add_shape(
        1, Emu(731520), Emu(960120), Emu(1097280), Emu(45720)
    )
    set_shape_fill(shape, ACCENT_BAR_COLOR)
    shape.line.fill.background()


def add_rect_cell(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    set_shape_fill(shape, fill_color)
    shape.line.fill.background()
    return shape


def insert_slide_after(prs, index):
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    slide_list = prs.slides._sldIdLst
    slide_elem = slide_list[-1]
    slide_list.remove(slide_elem)
    slide_list.insert(index + 1, slide_elem)
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Score matrix - 实验2 评估结果总览
# ═══════════════════════════════════════════════════════════════
def build_exp2_slide_1(prs, insert_after_idx):
    slide = insert_slide_after(prs, insert_after_idx)
    add_top_bar(slide)

    # Title
    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验2 - 需求变更实验：五维度软件工程质量评估",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Subtitle
    add_textbox(slide, Emu(731520), Emu(1050000), Emu(10500000), Emu(340000),
                u"同一基线系统 + 同一需求变更，分别用 OpenSpec-HW 与 Vanilla Claude Code 实现，独立评估五个 SE 维度",
                Pt(14), SUBTITLE_COLOR)

    # === Score matrix ===
    col_starts = [Emu(731520), Emu(3800000), Emu(6400000), Emu(9000000)]
    col_widths = [Emu(3030000), Emu(2560000), Emu(2560000), Emu(2560000)]
    row_h = Emu(380000)
    row_start = Emu(1500000)

    headers = [u"评估维度", u"OpenSpec-HW", u"Vanilla Claude", u"差距分析"]

    # Header row
    for ci, hdr in enumerate(headers):
        add_rect_cell(slide, col_starts[ci], row_start, col_widths[ci], row_h, HEADER_BG)
        add_textbox(slide, col_starts[ci], row_start + Emu(18000),
                    col_widths[ci], row_h,
                    hdr, Pt(13), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    rows_data = [
        [u"1. 需求与规格", "5", "1", u"规格文档 636 行 vs 0 行"],
        [u"2. 设计", "5", "1", u"模块化 vs 单文件膨胀 62%"],
        [u"3. 实现", "4", "2", u"类型安全 + 日志 vs 无类型/print"],
        [u"4. 验证与确认", "4", "1", u"31 个测试 vs 0 个测试"],
        [u"5. 维护与演进", "5", "1", u"YAML配置+归档 vs 硬编码"],
        [u"总分（满分 25）", "23", "6", u"OpenSpec-HW 领先 +17"],
        [u"平均分（满分 5）", "4.6", "1.2", u"差距 3.8 倍"],
    ]

    for ri, row_data in enumerate(rows_data):
        y = row_start + row_h * (ri + 1)
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE
        is_total = ri >= 5
        for ci, cell_text in enumerate(row_data):
            fill = bg
            if is_total:
                fill = RGBColor(0xE3, 0xF2, 0xFD)
            add_rect_cell(slide, col_starts[ci], y, col_widths[ci], row_h, fill)

            text_color = LABEL_COLOR
            text_bold = False
            if ci == 1:  # OpenSpec-HW green
                text_color = GREEN
                text_bold = is_total
            elif ci == 2:  # Vanilla Claude red
                text_color = RED
                text_bold = is_total
            elif ci == 3:  # Gap blue
                text_color = BLUE
                text_bold = is_total
            elif ci == 0:
                text_bold = is_total

            add_textbox(slide, col_starts[ci], y + Emu(18000),
                        col_widths[ci], row_h,
                        cell_text, Pt(12) if not is_total else Pt(13),
                        text_color, bold=text_bold,
                        alignment=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # === Bottom: two project summary boxes ===
    box_y = Emu(5050000)
    box_w = Emu(5400000)
    box_h = Emu(1450000)
    gap = Emu(150000)

    # OpenSpec-HW box
    add_rect_cell(slide, Emu(548640), box_y, box_w, box_h, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, Emu(700000), box_y + Emu(60000), box_w - Emu(200000), Emu(280000),
                u"OpenSpec-HW 方法", Pt(14), GREEN, bold=True)
    add_multiline_textbox(
        slide, Emu(700000), box_y + Emu(380000), box_w - Emu(250000), Emu(1000000),
        [
            u"6 个 artifact 文档（636行）+ 10 个边界场景",
            u"3 个新模块 + 31 个新测试 + 19 个配置参数",
            u"变更分布在 7 个文件，职责隔离清晰",
            u"T+1 冲突、渐进恢复、频率切换全覆盖",
        ],
        Pt(11), BODY_COLOR
    )

    # Vanilla Claude box
    x2 = Emu(548640) + box_w + gap
    add_rect_cell(slide, x2, box_y, box_w, box_h, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, x2 + Emu(150000), box_y + Emu(60000), box_w - Emu(200000), Emu(280000),
                u"Vanilla Claude Code 方法", Pt(14), RED, bold=True)
    add_multiline_textbox(
        slide, x2 + Emu(150000), box_y + Emu(380000), box_w - Emu(250000), Emu(1000000),
        [
            u"0 个规格文档 + 仅 2 个边界场景",
            u"0 个新模块、0 个测试，仅改 2 个文件",
            u"backtester.py 膨胀 62%（257→416 行）",
            u"遗漏 T+1 冲突、渐进恢复、频率切换",
        ],
        Pt(11), BODY_COLOR
    )

    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Dimension-by-dimension comparison
# ═══════════════════════════════════════════════════════════════
def build_exp2_slide_2(prs, insert_after_idx):
    slide = insert_slide_after(prs, insert_after_idx)
    add_top_bar(slide)

    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验2 - 需求变更结果：五维度逐项对比",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # Column layout matching slide 6 style
    col_starts = [Emu(457200), Emu(2900000), Emu(6700000), Emu(10500000)]
    col_widths = [Emu(2400000), Emu(3760000), Emu(3760000), Emu(1097280)]
    row_h = Emu(440000)
    row_start = Emu(1150000)

    header_data = [
        (u"对比维度", WHITE, HEADER_BG),
        (u"OpenSpec-HW", WHITE, RGBColor(0x2E, 0x7D, 0x32)),
        (u"Vanilla Claude Code", WHITE, RGBColor(0xC6, 0x28, 0x28)),
        (u"评分", WHITE, HEADER_BG),
    ]
    for ci, (text, text_color, bg_color) in enumerate(header_data):
        add_rect_cell(slide, col_starts[ci], row_start, col_widths[ci], row_h, bg_color)
        add_textbox(slide, col_starts[ci], row_start + Emu(30000),
                    col_widths[ci], row_h,
                    text, Pt(13), text_color, bold=True, alignment=PP_ALIGN.CENTER)

    rows = [
        {
            "label": u"需求与规格",
            "good": [u"proposal+usecases+specs+design+tasks", u"10 个边界场景 + 636 行 artifact"],
            "bad": [u"0 个规格文档、0 个用例", u"仅靠 requirement_change.md 直接编码"],
            "score": u"5 vs 1",
        },
        {
            "label": u"设计",
            "good": [u"3 新模块隔离 + CircuitBreakerState 状态机", u"向后兼容（volatility_ratio 可选参数）"],
            "bad": [u"全部逻辑塞入 run_backtest() 单函数", u"内联 if/elif、无状态机、硬编码阈值"],
            "score": u"5 vs 1",
        },
        {
            "label": u"实现",
            "good": [u"type hints + logging + 防御编程", u"dataclass 结构化 + YAML 配置 19 参数"],
            "bad": [u"无 type hints、print 调试、1e-9 硬编码", u"止损用百分比而非 ATR（偏离需求）"],
            "score": u"4 vs 2",
        },
        {
            "label": u"验证确认",
            "good": [u"31 新测试（边界值 + 状态转换 + T+1）", u"测试方法名映射用例（test_UC1_1 等）"],
            "bad": [u"0 个测试、0 测试基础设施", u"三个关键边界场景完全遗漏"],
            "score": u"4 vs 1",
        },
        {
            "label": u"维护演进",
            "good": [u"变更归档 + delta specs + 可追溯链", u"新增 4 维度只需改 YAML + 新模块"],
            "bad": [u"0 文档、0 设计记录、无变更历史", u"任何修改需读懂 416 行单函数"],
            "score": u"5 vs 1",
        },
    ]

    for ri, row in enumerate(rows):
        y = row_start + row_h * (ri + 1)
        bg = ROW_LIGHT if ri % 2 == 0 else ROW_WHITE

        for ci in range(4):
            add_rect_cell(slide, col_starts[ci], y, col_widths[ci], row_h, bg)

        # Column 0: label
        add_textbox(slide, col_starts[0] + Emu(80000), y + Emu(10000),
                    col_widths[0] - Emu(80000), row_h,
                    row["label"], Pt(12), LABEL_COLOR, bold=True)

        # Column 1: OpenSpec-HW (good)
        add_multiline_textbox(
            slide, col_starts[1] + Emu(80000), y + Emu(5000),
            col_widths[1] - Emu(100000), row_h,
            row["good"], Pt(10), GREEN
        )

        # Column 2: Vanilla Claude (bad)
        add_multiline_textbox(
            slide, col_starts[2] + Emu(80000), y + Emu(5000),
            col_widths[2] - Emu(100000), row_h,
            row["bad"], Pt(10), RED
        )

        # Column 3: score
        add_textbox(slide, col_starts[3], y + Emu(30000),
                    col_widths[3], row_h,
                    row["score"], Pt(14), BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Summary box at bottom
    summary_y = row_start + row_h * 6 + Emu(60000)
    add_rect_cell(slide, Emu(365760), summary_y, Emu(11430000), Emu(600000),
                  RGBColor(0xE8, 0xF5, 0xE9))

    add_multiline_textbox(
        slide, Emu(548640), summary_y + Emu(20000), Emu(11064240), Emu(550000),
        [
            (u"核心结论：需求变更场景下差距比初始构建更大。OpenSpec 的 usecases 阶段强制识别边界场景（10 vs 2），是质量差异的根源。", True, SUMMARY_COLOR),
            (u"Vanilla Claude 跳过需求分析直接编码，导致 3 个关键边界遗漏（T+1 冲突、渐进恢复、频率切换），且产生单文件膨胀的技术债。", False, SUBTITLE_COLOR),
        ],
        Pt(12), SUMMARY_COLOR
    )

    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Key findings & boundary case analysis
# ═══════════════════════════════════════════════════════════════
def build_exp2_slide_3(prs, insert_after_idx):
    slide = insert_slide_after(prs, insert_after_idx)
    add_top_bar(slide)

    add_textbox(slide, Emu(731520), Emu(365760), Emu(10500000), Emu(521970),
                u"实验2 - 关键发现：需求变更中的工程质量差异",
                TITLE_SIZE, TITLE_COLOR, bold=True)
    add_accent_underline(slide)

    # --- Left column: 4 finding boxes stacked ---
    findings = [
        {
            "num": "01",
            "title": u"边界场景发现：10 vs 2",
            "desc": u"OpenSpec 的 usecases 阶段系统性枚举了 10 个边界场景（含 T+1 冲突、熔断恢复中断、频率切换等），"
                    u"Vanilla Claude 仅识别 2 个，遗漏 3 个关键场景。",
            "color": BLUE,
        },
        {
            "num": "02",
            "title": u"模块化 vs 单文件膨胀",
            "desc": u"OpenSpec 创建 3 个新模块（volatility/circuit_breaker/stop_loss），变更分布 7 个文件。"
                    u"Vanilla Claude 将全部逻辑塞入 backtester.py，单文件膨胀 62%。",
            "color": GREEN,
        },
        {
            "num": "03",
            "title": u"需求偏离：止损实现不符规格",
            "desc": u"需求要求 ATR-based 自适应止损（1.5x/2.0x/2.5x ATR 乘数），Vanilla Claude 实现为简单百分比止损"
                    u"（-6%/-8%/-10%），偏离了原始需求规格。",
            "color": ORANGE,
        },
        {
            "num": "04",
            "title": u"零测试 = 零信心",
            "desc": u"OpenSpec 产出 31 个测试覆盖状态转换、边界值和 T+1。Vanilla Claude 零测试——"
                    u"量化交易系统的数值正确性直接影响资金安全，零测试是致命缺陷。",
            "color": RED,
        },
    ]

    box_w = Emu(5200000)
    box_h = Emu(1080000)
    gap = Emu(80000)
    start_x = Emu(380000)
    start_y = Emu(1200000)

    for i, finding in enumerate(findings):
        y = start_y + (box_h + gap) * i
        add_rect_cell(slide, start_x, y, box_w, box_h, RGBColor(0xFA, 0xFA, 0xFA))

        # Number badge
        badge_w = Emu(360000)
        badge_h = Emu(360000)
        add_rect_cell(slide, start_x + Emu(60000), y + Emu(60000), badge_w, badge_h, finding["color"])
        add_textbox(slide, start_x + Emu(60000), y + Emu(75000), badge_w, badge_h,
                    finding["num"], Pt(16), WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, start_x + Emu(480000), y + Emu(70000), box_w - Emu(560000), Emu(300000),
                    finding["title"], Pt(13), finding["color"], bold=True)

        # Description
        add_textbox(slide, start_x + Emu(80000), y + Emu(440000), box_w - Emu(160000), Emu(600000),
                    finding["desc"], Pt(10), BODY_COLOR)

    # --- Right column: missed boundary cases highlight ---
    right_x = Emu(5800000)
    right_w = Emu(5700000)

    # Header
    add_textbox(slide, right_x, Emu(1200000), right_w, Emu(350000),
                u"Vanilla Claude 遗漏的 3 个关键边界场景", Pt(14), RED, bold=True)

    # Case 1: T+1
    case_y = Emu(1600000)
    case_h = Emu(1100000)
    add_rect_cell(slide, right_x, case_y, right_w, case_h, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, right_x + Emu(100000), case_y + Emu(40000), right_w - Emu(200000), Emu(250000),
                u"T+1 结算冲突", Pt(12), RED, bold=True)
    add_multiline_textbox(
        slide, right_x + Emu(100000), case_y + Emu(340000), right_w - Emu(200000), Emu(750000),
        [
            u"熔断触发全清仓时，当日买入的股票因 T+1 规则无法卖出。",
            u"OpenSpec 在 usecases 阶段即识别此冲突，设计了延迟卖出队列。",
            u"Vanilla Claude 完全未考虑，熔断日可能产生非法交易指令。",
        ],
        Pt(10), BODY_COLOR
    )

    # Case 2: Gradual recovery
    case_y2 = case_y + case_h + Emu(70000)
    add_rect_cell(slide, right_x, case_y2, right_w, case_h, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, right_x + Emu(100000), case_y2 + Emu(40000), right_w - Emu(200000), Emu(250000),
                u"渐进恢复机制", Pt(12), RED, bold=True)
    add_multiline_textbox(
        slide, right_x + Emu(100000), case_y2 + Emu(340000), right_w - Emu(200000), Emu(750000),
        [
            u"L2 熔断后恢复应分阶段（L2→recovering→normal），而非一步跳回满仓。",
            u"OpenSpec 设计了 4 状态状态机，含恢复中断（新下跌重置计数器）。",
            u"Vanilla Claude 直接从清仓跳回满仓，可能在市场未稳定时重新满仓。",
        ],
        Pt(10), BODY_COLOR
    )

    # Case 3: Frequency switching
    case_y3 = case_y2 + case_h + Emu(70000)
    add_rect_cell(slide, right_x, case_y3, right_w, case_h, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, right_x + Emu(100000), case_y3 + Emu(40000), right_w - Emu(200000), Emu(250000),
                u"再平衡频率切换", Pt(12), RED, bold=True)
    add_multiline_textbox(
        slide, right_x + Emu(100000), case_y3 + Emu(340000), right_w - Emu(200000), Emu(750000),
        [
            u"月度→双周→周度切换时，需处理月中状态变更和计时器重置。",
            u"OpenSpec 用例明确定义了频率切换的过渡规则。",
            u"Vanilla Claude 仅用简单天数计数器，频率切换时可能跳过再平衡。",
        ],
        Pt(10), BODY_COLOR
    )

    # Bottom summary
    add_rect_cell(slide, Emu(365760), Emu(5850000), Emu(11430000), Emu(500000),
                  RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, Emu(548640), Emu(5880000), Emu(11064240), Emu(440000),
                u"结论：需求变更场景放大了规格驱动开发的优势。边界场景覆盖率（10:2）是质量差异的核心驱动力——有 spec 才能发现边界，有边界才能写对代码。",
                Pt(13), SUMMARY_COLOR, bold=True)

    return slide


def main():
    prs = Presentation(PPTX_PATH)
    print(f"Original slide count: {len(prs.slides)}")

    # Insert 3 slides after slide 8 (index 7)
    # Insert in reverse order so indices stay correct
    s3 = build_exp2_slide_3(prs, 7)
    s2 = build_exp2_slide_2(prs, 7)
    s1 = build_exp2_slide_1(prs, 7)

    print(f"New slide count: {len(prs.slides)}")
    prs.save(OUTPUT_PATH)
    print(f"Saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
