"""
Render Slide 4 as a PNG image for visual verification.
Uses PIL to draw shapes based on PPTX data.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from PIL import Image, ImageDraw, ImageFont
import os

prs = Presentation('OpenSpec_技术研究汇报_精简版.pptx')
slide = prs.slides[3]

# Scale: EMU to pixels
# Slide: 12190730 x 6858000 EMU
# Target image: ~2400 x 1350 px
SCALE = 2400 / 12190730
IMG_W = int(12190730 * SCALE)
IMG_H = int(6858000 * SCALE)

img = Image.new('RGB', (IMG_W, IMG_H), (255, 255, 255))
draw = ImageDraw.Draw(img)

def emu_to_px(val):
    return int(val * SCALE)

# Try to load a font
try:
    font_small = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 11)
    font_medium = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 14)
    font_large = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 18)
    font_xlarge = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 28)
except:
    font_small = ImageFont.load_default()
    font_medium = font_small
    font_large = font_small
    font_xlarge = font_small

def get_fill_color(shape):
    """Extract fill color from shape."""
    try:
        if shape.fill.type is not None:
            rgb = shape.fill.fore_color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except:
        pass
    return None

def get_text_info(shape):
    """Extract text and basic formatting."""
    if not shape.has_text_frame:
        return []
    results = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = run.text
            color = (0, 0, 0)
            try:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    color = (rgb[0], rgb[1], rgb[2])
            except:
                pass
            size = 10
            try:
                if run.font.size:
                    size = run.font.size.pt
            except:
                pass
            results.append((text, color, size, run.font.bold))
    return results

# Draw all shapes
for shape in slide.shapes:
    x = emu_to_px(shape.left)
    y = emu_to_px(shape.top)
    w = emu_to_px(shape.width)
    h = emu_to_px(shape.height)

    # Draw fill
    fill = get_fill_color(shape)
    if fill and w > 2 and h > 2:
        draw.rectangle([x, y, x + w, y + h], fill=fill)

    # Draw text
    text_info = get_text_info(shape)
    if text_info:
        combined_text = ''.join(t[0] for t in text_info)
        if combined_text.strip():
            # Use first run's color
            color = text_info[0][1]
            size = text_info[0][2]

            # Choose font based on size
            if size >= 24:
                font = font_xlarge
            elif size >= 14:
                font = font_large
            elif size >= 10:
                font = font_medium
            else:
                font = font_small

            # Truncate text to fit
            max_chars = max(1, w // 7)
            display_text = combined_text[:max_chars]

            # Draw text with padding
            tx = x + 3
            ty = y + 2
            try:
                draw.text((tx, ty), display_text, fill=color, font=font)
            except:
                draw.text((tx, ty), display_text, fill=color)

# Save
output_path = os.path.join(os.getcwd(), 'slide4_preview.png')
img.save(output_path, 'PNG')
print(f"Saved: {output_path}")
print(f"Image size: {IMG_W}x{IMG_H}")
