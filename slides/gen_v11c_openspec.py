# -*- coding: utf-8 -*-
"""
Add OpenSpec application section to slide 3 of v11c.
Places a new card in the space below the middle and right columns.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import glob, os

DARK  = RGBColor(0x2D, 0x2D, 0x2D)
GRAY  = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = '\u5fae\u8f6f\u96c5\u9ed1'

# Teal theme for the OpenSpec section
TEAL_HDR = RGBColor(0x00, 0x89, 0x7B)
TEAL_BG  = RGBColor(0xE0, 0xF2, 0xF1)

base = r'C:\Users\Zhifei Dou\Desktop\Openspec_modified\slides'
v11c = [f for f in glob.glob(os.path.join(base, '*v11c.pptx')) if '~$' not in f][0]
print(f"Reading: {v11c}")
prs = Presentation(v11c)

slide = prs.slides[2]  # Slide 3

# --- Card position: below middle+right columns, above summary bar ---
card_x = 4.52
card_y = 3.85
card_w = 8.32   # spans middle + right column width
card_h = 2.65

# Background
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(card_x), Inches(card_y), Inches(card_w), Inches(card_h))
bg.fill.solid()
bg.fill.fore_color.rgb = TEAL_BG
bg.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
bg.line.width = Pt(0.5)

# Header bar
hdr = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(card_x), Inches(card_y), Inches(card_w), Inches(0.26))
hdr.fill.solid()
hdr.fill.fore_color.rgb = TEAL_HDR
hdr.line.fill.background()

# Left accent bar
acc = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(card_x), Inches(card_y), Inches(0.04), Inches(card_h))
acc.fill.solid()
acc.fill.fore_color.rgb = TEAL_HDR
acc.line.fill.background()

# Header text
htb = slide.shapes.add_textbox(
    Inches(card_x + 0.12), Inches(card_y + 0.02),
    Inches(card_w - 0.2), Inches(0.26))
htf = htb.text_frame
htf.word_wrap = True
p = htf.paragraphs[0]
p.space_before = Pt(0)
p.space_after = Pt(0)
r = p.add_run()
r.text = "\u4ee5OpenSpec\u4e3a\u4f8b: \u77e5\u8bc6\u6c89\u6dc0\u4e0e\u590d\u7528\u5982\u4f55\u670d\u52a1Spec\u9a71\u52a8\u7684\u7cfb\u7edf\u5de5\u7a0b\u6846\u67b6"
# 以OpenSpec为例: 知识沉淀与复用如何服务Spec驱动的系统工程框架
r.font.size = Pt(10)
r.font.bold = True
r.font.color.rgb = WHITE
r.font.name = FONT


# Body content
body_paras = [
    # ① Spec质量反馈闭环
    ("\u2460 Spec\u8d28\u91cf\u53cd\u9988\u95ed\u73af: \u4ece\u6d4b\u8bd5\u7ed3\u679c\u4e2d\u6c89\u6dc0\u7f16\u5199\u7ecf\u9a8c",
     True, 8.5, DARK),
    # ① Spec质量反馈闭环: 从测试结果中沉淀编写经验

    ("OpenSpec\u7684\u6838\u5fc3\u77e5\u8bc6\u662fSpec(\u7ed3\u6784\u5316\u884c\u4e3a\u89c4\u683c, \u5305\u542bRequirement + Scenario + GIVEN/WHEN/THEN)\u3002",
     False, 8, DARK),
    # OpenSpec的核心知识是Spec(结构化行为规格, 包含Requirement + Scenario + GIVEN/WHEN/THEN)。

    ("Agent\u6839\u636eSpec\u5b9e\u73b0\u4ee3\u7801\u5e76\u6d4b\u8bd5, \u4f46\u6d4b\u8bd5\u901a\u8fc7\u6216\u5931\u8d25\u7684\u7ed3\u679c\u76ee\u524d\u6ca1\u6709\u53cd\u9988\u56deSpec\u8d28\u91cf\u5c42\u9762\u3002",
     False, 8, DARK),
    # Agent根据Spec实现代码并测试, 但测试通过或失败的结果目前没有反馈回Spec质量层面。

    ("\u673a\u4f1a: \u8ffd\u8e2aSpec\u2192\u5b9e\u73b0\u2192\u6d4b\u8bd5\u7684\u5b8c\u6574\u94fe\u8def, \u8bb0\u5f55\u6bcf\u4e2aSpec\u7684\u5b9e\u73b0\u6210\u529f\u7387, \u63d0\u70bc\"\u4ec0\u4e48\u5199\u6cd5\u7684Spec\u66f4\u5bb9\u6613\u88abAgent\u6b63\u786e\u5b9e\u73b0\"\u7684\u89c4\u5f8b\u3002",
     False, 8, GRAY),
    # 机会: 追踪Spec→实现→测试的完整链路, 记录每个Spec的实现成功率, 提炼"什么写法的Spec更容易被Agent正确实现"的规律。

    ("", False, 2, DARK),

    # ② 领域知识条件检索注入
    ("\u2461 \u9886\u57df\u77e5\u8bc6\u7684\u6761\u4ef6\u68c0\u7d22\u4e0e\u6ce8\u5165",
     True, 8.5, DARK),
    # ② 领域知识的条件检索与注入

    ("\u9879\u76ee\u6d89\u53ca\u7279\u5b9a\u884c\u4e1a(\u5149\u901a\u4fe1\u3001\u667a\u9a7e\u7b49), \u6bcf\u4e2a\u884c\u4e1a\u6709\u81ea\u5df1\u7684\u6807\u51c6\u548c\u7ea6\u675f(\u59823GPP\u3001ISO 26262)\u3002",
     False, 8, DARK),
    # 项目涉及特定行业(光通信、智驾等), 每个行业有自己的标准和约束(如3GPP、ISO 26262)。

    ("\u73b0\u72b6: \u8fd9\u4e9b\u9886\u57df\u77e5\u8bc6\u5b8c\u5168\u4f9d\u8d56\u7528\u6237\u624b\u52a8\u63d0\u4f9b, Agent\u65e0\u6cd5\u4e3b\u52a8\u83b7\u53d6\u3002",
     False, 8, DARK),
    # 现状: 这些领域知识完全依赖用户手动提供, Agent无法主动获取。

    ("\u673a\u4f1a: Agent\u8bc6\u522b\u5f53\u524dSpec\u6240\u5c5e\u9886\u57df\u540e, \u81ea\u52a8\u4ece\u9886\u57df\u77e5\u8bc6\u5e93\u68c0\u7d22\u76f8\u5173\u6807\u51c6\u6761\u76ee, \u6ce8\u5165Agent\u4e0a\u4e0b\u6587\u8f85\u52a9Spec\u7f16\u5199\u548c\u4ee3\u7801\u5b9e\u73b0\u3002",
     False, 8, GRAY),
    # 机会: Agent识别当前Spec所属领域后, 自动从领域知识库检索相关标准条目, 注入Agent上下文辅助Spec编写和代码实现。

    ("", False, 2, DARK),

    # ③ 跨项目Spec模式复用
    ("\u2462 \u8de8\u9879\u76eeSpec\u6a21\u5f0f\u590d\u7528",
     True, 8.5, DARK),
    # ③ 跨项目Spec模式复用

    ("\u4e0d\u540c\u9879\u76ee\u4e2d\u5b58\u5728\u5927\u91cf\u76f8\u4f3c\u7684Spec\u6a21\u5f0f(\u5982CRUD\u63a5\u53e3\u89c4\u683c\u3001\u8ba4\u8bc1\u6388\u6743\u89c4\u683c\u3001\u6d88\u606f\u961f\u5217\u89c4\u683c), \u6bcf\u6b21\u90fd\u4ece\u96f6\u7f16\u5199\u3002",
     False, 8, DARK),
    # 不同项目中存在大量相似的Spec模式(如CRUD接口规格、认证授权规格、消息队列规格), 每次都从零编写。

    ("\u673a\u4f1a: \u4ece\u5df2\u5f52\u6863\u7684\u5386\u53f2changes\u4e2d, \u7531LLM\u81ea\u52a8\u5f52\u7eb3\u5e38\u89c1Spec\u6a21\u677f\u3002\u65b0\u9700\u6c42\u65f6\u6839\u636e\u63cf\u8ff0\u5339\u914d\u63a8\u8350\u6700\u76f8\u5173\u7684\u5386\u53f2Spec\u4f5c\u4e3a\u8d77\u70b9\u3002",
     False, 8, GRAY),
    # 机会: 从已归档的历史changes中, 由LLM自动归纳常见Spec模板。新需求时根据描述匹配推荐最相关的历史Spec作为起点。
]

# Add body textbox
btb = slide.shapes.add_textbox(
    Inches(card_x + 0.12), Inches(card_y + 0.30),
    Inches(card_w - 0.24), Inches(card_h - 0.36))
btf = btb.text_frame
btf.word_wrap = True

first = True
for text, bold, sz, clr in body_paras:
    if first:
        p = btf.paragraphs[0]
        first = False
    else:
        p = btf.add_paragraph()
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FONT

prs.save(v11c)
print(f"Done! OpenSpec section added to slide 3. Saved to: {v11c}")
