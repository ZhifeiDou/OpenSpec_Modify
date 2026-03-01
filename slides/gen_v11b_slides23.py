# -*- coding: utf-8 -*-
"""
Add 2 slides to v11b.pptx condensing the research report.
Style matches reference PPT: vivid header bars, tinted backgrounds, bottom summary bar.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import glob, os

DARK  = RGBColor(0x2D, 0x2D, 0x2D)
GRAY  = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAR   = RGBColor(0x1A, 0x2B, 0x4A)
RED   = RGBColor(0xC0, 0x39, 0x2B)
FONT  = '\u5fae\u8f6f\u96c5\u9ed1'

# Card theme colors (header bar, body bg, left accent)
BLUE_HDR   = RGBColor(0x1A, 0x73, 0xE8)
BLUE_BG    = RGBColor(0xEB, 0xF3, 0xFE)
PURPLE_HDR = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_BG  = RGBColor(0xF5, 0xF0, 0xFF)
RED_HDR    = RGBColor(0xD1, 0x34, 0x38)
RED_BG     = RGBColor(0xFF, 0xF0, 0xF0)
GREEN_HDR  = RGBColor(0x0D, 0x7C, 0x3F)
GREEN_BG   = RGBColor(0xE6, 0xF4, 0xEA)
DARK_BAR   = RGBColor(0x1A, 0x1A, 0x2E)

base = r'C:\Users\Zhifei Dou\Desktop\Openspec_modified\slides'
v11b = [f for f in glob.glob(os.path.join(base, '*v11b.pptx')) if '~$' not in f][0]
print(f"Reading: {v11b}")
prs = Presentation(v11b)


def add_tb(slide, left, top, width, height, paras):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for text, bold, sz, clr in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = clr
        r.font.name = FONT


def add_card(slide, x, y, w, h, hdr_color, bg_color, hdr_text, body_paras):
    """Add a card: bg rect + colored header bar + left accent + text."""
    # Body background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    bg.line.width = Pt(0.5)

    # Header bar
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.26))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = hdr_color
    hdr.line.fill.background()

    # Left accent bar
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(0.04), Inches(h))
    acc.fill.solid()
    acc.fill.fore_color.rgb = hdr_color
    acc.line.fill.background()

    # Header text
    add_tb(slide, x + 0.12, y + 0.02, w - 0.2, 0.26,
           [(hdr_text, True, 10, WHITE)])

    # Body text
    add_tb(slide, x + 0.12, y + 0.32, w - 0.24, h - 0.38, body_paras)


def make_slide_3col(prs, title, subtitle, col_specs, summary_paras):
    """
    3-column layout + bottom summary bar.
    col_specs: list of 3 tuples (hdr_color, bg_color, hdr_text, body_paras)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Top/bottom bars
    for yy in [0.0, 7.45]:
        b = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(yy),
            Inches(13.33), Inches(0.05))
        b.fill.solid()
        b.fill.fore_color.rgb = BAR
        b.line.fill.background()

    # Title
    add_tb(slide, 0.49, 0.15, 10.0, 0.48,
           [(title, True, 18, DARK)])

    # Subtitle
    if subtitle:
        add_tb(slide, 0.49, 0.60, 10.0, 0.30,
               [(subtitle, False, 9, GRAY)])

    # Accent line
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.49), Inches(0.95),
        Inches(1.2), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RED
    ln.line.fill.background()

    # 3 columns
    col_w = [4.10, 4.10, 4.10]
    col_x = [0.30, 4.52, 8.74]
    col_y = 1.10
    col_h = 4.80

    for i, (hc, bc, ht, bp) in enumerate(col_specs):
        add_card(slide, col_x[i], col_y, col_w[i], col_h,
                 hc, bc, ht, bp)

    # Bottom summary bar
    sbar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.30), Inches(6.05),
        Inches(12.54), Inches(0.75))
    sbar.fill.solid()
    sbar.fill.fore_color.rgb = DARK_BAR
    sbar.line.fill.background()

    add_tb(slide, 0.50, 6.10, 12.14, 0.65, summary_paras)


# =====================================================================
# SLIDE 2: 知识沉淀的形态、范式与生命周期管理
# =====================================================================

s2_col1 = (BLUE_HDR, BLUE_BG,
    "\u4e03\u79cd\u77e5\u8bc6\u6c89\u6dc0\u5f62\u6001 (2023\u20132026)",  # 七种知识沉淀形态 (2023-2026)
    [
    ("Agent\u7ecf\u9a8c\u79ef\u7d2f\u7814\u7a76\u4e2d, \u77e5\u8bc6\u6c89\u6dc0\u7684\u5f62\u6001\u5448\u73b0\u51fa\u4ece\u4f4e\u62bd\u8c61\u5230\u9ad8\u62bd\u8c61\u7684\u6e05\u6670\u8c31\u7cfb:", False, 8, DARK),
    # Agent经验积累研究中, 知识沉淀的形态呈现出从低抽象到高抽象的清晰谱系:
    ("", False, 2, DARK),

    ("\u2460 \u539f\u59cb\u8f68\u8ff9 (ExpeL, JARVIS-1)", True, 8, DARK),
    # ① 原始轨迹
    ("Agent\u4e0e\u73af\u5883\u4ea4\u4e92\u7684\u5b8c\u6574\u884c\u52a8\u5e8f\u5217, \u5305\u542b\u6210\u529f\u548c\u5931\u8d25\u4e24\u7c7b\u8f68\u8ff9\u3002", False, 8, DARK),
    # Agent与环境交互的完整行动序列, 包含成功和失败两类轨迹。
    ("JARVIS-1\u6269\u5c55\u4e3a\u591a\u6a21\u6001\u8f68\u8ff9, \u540c\u65f6\u5b58\u50a8\u89c6\u89c9\u622a\u56fe\u548c\u6587\u672c\u8ba1\u5212\u3002", False, 8, DARK),
    # JARVIS-1扩展为多模态轨迹, 同时存储视觉截图和文本计划。
    ("\u4e0a\u4e0b\u6587\u6700\u4e30\u5bcc, \u4f46\u68c0\u7d22\u6548\u7387\u4f4e, \u76f4\u63a5\u6ce8\u5165\u6210\u672c\u9ad8\u3002", False, 8, GRAY),
    # 上下文最丰富, 但检索效率低, 直接注入成本高。

    ("\u2461 \u81ea\u7136\u8bed\u8a00\u53cd\u601d (Reflexion, NeurIPS 2023)", True, 8, DARK),
    # ② 自然语言反思
    ("Agent\u5728\u4efb\u52a1\u5931\u8d25\u540e\u751f\u6210\"\u54ea\u91cc\u51fa\u4e86\u9519\"\u4ee5\u53ca\"\u4e0b\u6b21\u5982\u4f55\u6539\u8fdb\"", False, 8, DARK),
    # Agent在任务失败后生成"哪里出了错"以及"下次如何改进"
    ("\u7684\u7b80\u77ed\u6587\u672c\u3002\u6781\u5176\u7075\u6d3b, LLM\u539f\u751f\u53ef\u7406\u89e3, \u4f46\u7f3a\u4e4f\u7ed3\u6784\u5316\u7ea6\u675f\u3002", False, 8, DARK),
    # 的简短文本。极其灵活, LLM原生可理解, 但缺乏结构化约束。

    ("\u2462 \u84b8\u998f\u89c4\u5219/\u6d1e\u5bdf (ExpeL/EvolveR)", True, 8, DARK),
    # ③ 蒸馏规则/洞察
    ("\u901a\u8fc7\u5bf9\u6bd4\u6210\u529f\u548c\u5931\u8d25\u8f68\u8ff9, \u63d0\u53d6\u8de8\u4efb\u52a1\u7684\u9ad8\u5c42\u7b56\u7565\u89c4\u5219\u3002", False, 8, DARK),
    # 通过对比成功和失败轨迹, 提取跨任务的高层策略规则。
    ("EvolveR\u5c06\u89c4\u5219\u5347\u7ea7\u4e3a\u5e26\u8d28\u91cf\u8bc4\u5206\u7684\u6218\u7565\u539f\u5219, \u6bcf\u6761\u9644\u6709metric score\u3002", False, 8, DARK),
    # EvolveR将规则升级为带质量评分的战略原则, 每条附有metric score。

    ("\u2463 \u6761\u4ef6\u5316\u6307\u5357 (AutoGuide, NeurIPS 2024)", True, 8, DARK),
    # ④ 条件化指南
    ("\u6bcf\u6761\u6307\u5357\u662f{context, guideline}\u952e\u503c\u5bf9, \u660e\u786e\u63cf\u8ff0\u9002\u7528\u7684\u4e0a\u4e0b\u6587\u6761\u4ef6\u3002", False, 8, DARK),
    # 每条指南是{context, guideline}键值对, 明确描述适用的上下文条件。
    ("\u534a\u7ed3\u6784\u5316\u6bd4\u7eaf\u81ea\u7136\u8bed\u8a00\u66f4\u7cbe\u786e, \u652f\u6301\u57fa\u4e8e\u4e0a\u4e0b\u6587\u5339\u914d\u7684\u9009\u62e9\u6027\u6ce8\u5165\u3002", False, 8, DARK),
    # 半结构化比纯自然语言更精确, 支持基于上下文匹配的选择性注入。

    ("\u2464 \u5de5\u4f5c\u6d41\u6a21\u677f (AWM, ICML 2025)", True, 8, DARK),
    # ⑤ 工作流模板
    ("\u4ece\u6210\u529f\u8f68\u8ff9\u5f52\u7eb3\u53ef\u590d\u7528\u7684\u62bd\u8c61\u64cd\u4f5c\u5e8f\u5217, \u652f\u6301\u5c42\u7ea7\u7ec4\u5408\u3002", False, 8, DARK),
    # 从成功轨迹归纳可复用的抽象操作序列, 支持层级组合。
    ("WebArena\u4e0a\u5e26\u676551.1%\u7684\u76f8\u5bf9\u6210\u529f\u7387\u63d0\u5347\u3002", False, 8, DARK),
    # WebArena上带来51.1%的相对成功率提升。

    ("\u2465 \u53ef\u6267\u884c\u6280\u80fd (Voyager/LILO)", True, 8, DARK),
    # ⑥ 可执行技能
    ("\u6bcf\u4e2a\u6280\u80fd\u662f\u5b8c\u6574\u7684\u4ee3\u7801\u51fd\u6570+\u81ea\u7136\u8bed\u8a00\u6587\u6863\u3002LILO\u901a\u8fc7Stitch\u7b97\u6cd5", False, 8, DARK),
    # 每个技能是完整的代码函数+自然语言文档。LILO通过Stitch算法
    ("\u81ea\u52a8\u8bc6\u522b\u8de8\u4ee3\u7801\u7684\u5171\u6027\u6a21\u5f0f\u5e76\u5408\u5e76\u4e3a\u53ef\u590d\u7528\u62bd\u8c61\u3002\u590d\u7528\u6027\u6700\u9ad8\u3002", False, 8, DARK),
    # 自动识别跨代码的共性模式并合并为可复用抽象。复用性最高。

    ("\u2466 \u7ed3\u6784\u5316\u7ecf\u9a8c\u5361\u7247 (MemGovern, 2026)", True, 8, DARK),
    # ⑦ 结构化经验卡片
    ("\u4e13\u95e8\u9762\u5411\u4ee3\u7801Agent, \u5c06Issue/PR\u8f6c\u5316\u4e3a\u53cc\u5c42\u7ed3\u6784:", False, 8, DARK),
    # 专门面向代码Agent, 将Issue/PR转化为双层结构:
    ("\u7d22\u5f15\u5c42(\u95ee\u9898\u6458\u8981, \u9519\u8bef\u7b7e\u540d) + \u51b3\u8bae\u5c42(\u6839\u56e0, \u4fee\u590d, \u9a8c\u8bc1)\u3002", False, 8, DARK),
    # 索引层(问题摘要, 错误签名) + 决议层(根因, 修复, 验证)。
])

s2_col2 = (PURPLE_HDR, PURPLE_BG,
    "\u5b58\u50a8\u7ed3\u6784\u4e0e\u751f\u547d\u5468\u671f\u7ba1\u7406",  # 存储结构与生命周期管理
    [
    ("\u5b58\u50a8\u7ed3\u6784\u6f14\u8fdb", True, 9, DARK),
    # 存储结构演进
    ("\u7eaf\u81ea\u7136\u8bed\u8a00 (Reflexion, ExpeL\u7684insights)", False, 8, DARK),
    # 纯自然语言 (Reflexion, ExpeL的insights)
    ("\u6700\u7075\u6d3b\u4f46\u6700\u4e0d\u7cbe\u786e\u3002", False, 8, GRAY),
    # 最灵活但最不精确。
    ("\u534a\u7ed3\u6784\u5316\u952e\u503c\u5bf9 (AutoGuide\u7684context-guideline\u5b57\u5178)", False, 8, DARK),
    # 半结构化键值对 (AutoGuide的context-guideline字典)
    ("\u5728\u7075\u6d3b\u6027\u548c\u7cbe\u786e\u6027\u95f4\u53d6\u5f97\u5e73\u8861\u3002", False, 8, GRAY),
    # 在灵活性和精确性间取得平衡。
    ("\u53ef\u6267\u884c\u4ee3\u7801 (Voyager\u7684JavaScript\u51fd\u6570, LILO\u7684\u03bb-abstraction)", False, 8, DARK),
    # 可执行代码 (Voyager的JavaScript函数, LILO的λ-abstraction)
    ("\u7cbe\u786e\u5ea6\u6700\u9ad8\u4f46\u6cdb\u5316\u53d7\u9650\u3002", False, 8, GRAY),
    # 精确度最高但泛化受限。
    ("\u7ed3\u6784\u5316\u7b14\u8bb0\u7f51\u7edc (A-MEM, Zettelkasten\u5f0f)", False, 8, DARK),
    # 结构化笔记网络 (A-MEM, Zettelkasten式)
    ("\u901a\u8fc7keywords, tags, embedding\u548c\u7b14\u8bb0\u95f4\u94fe\u63a5\u5b9e\u73b0\u591a\u5c5e\u6027\u5173\u8054\u3002", False, 8, GRAY),
    # 通过keywords, tags, embedding和笔记间链接实现多属性关联。
    ("\u77e5\u8bc6\u56fe\u8c31\u4e09\u5143\u7ec4 (Mem0, \u6709\u5411\u6807\u7b7e\u56fe)", False, 8, DARK),
    # 知识图谱三元组 (Mem0, 有向标签图)
    ("\u652f\u6301\u591a\u8df3\u63a8\u7406\u548c\u65f6\u5e8f\u63a8\u7406\u3002", False, 8, GRAY),
    # 支持多跳推理和时序推理。
    ("", False, 3, DARK),

    ("\u77e5\u8bc6\u751f\u547d\u5468\u671f\u7ba1\u7406", True, 9, DARK),
    # 知识生命周期管理
    ("", False, 2, DARK),
    ("\u65b0\u589e:", True, 8, DARK),
    # 新增:
    ("\u81ea\u52a8\u751f\u6210 \u2014 Voyager\u901a\u8fc7\u4ee3\u7801\u751f\u6210+self-verify\u540e\u5165\u5e93", False, 8, DARK),
    # 自动生成 — Voyager通过代码生成+self-verify后入库
    ("\u538b\u7f29\u5f52\u7eb3 \u2014 LILO\u7684Stitch\u7b97\u6cd5, AWM\u7684workflow\u5f52\u7eb3", False, 8, DARK),
    # 压缩归纳 — LILO的Stitch算法, AWM的workflow归纳
    ("\u52a8\u6001\u63d0\u53d6 \u2014 Mem0\u4ece\u4eba-Agent\u5bf9\u8bdd\u4e2d\u81ea\u52a8\u63d0\u53d6\u5173\u952e\u4e8b\u5b9e", False, 8, DARK),
    # 动态提取 — Mem0从人-Agent对话中自动提取关键事实
    ("", False, 2, DARK),
    ("\u66f4\u65b0\u4e0e\u6f14\u5316:", True, 8, DARK),
    # 更新与演化:
    ("A-MEM: \u65b0\u8bb0\u5fc6\u6574\u5408\u65f6\u53ef\u89e6\u53d1\u5bf9\u5df2\u6709\u5386\u53f2\u8bb0\u5fc6\u7684\u4e0a\u4e0b\u6587\u8868\u793a\u548c\u5c5e\u6027\u66f4\u65b0", False, 8, DARK),
    # A-MEM: 新记忆整合时可触发对已有历史记忆的上下文表示和属性更新
    ("Mem0: \u51b2\u7a81\u68c0\u6d4b\u5668+\u66f4\u65b0\u89e3\u6790\u5668, \u652f\u6301add/merge/invalidate/skip", False, 8, DARK),
    # Mem0: 冲突检测器+更新解析器, 支持add/merge/invalidate/skip
    ("EvolveR: \u5b8c\u6574\u7ef4\u62a4\u7ba1\u7ebf(\u53bb\u91cd, \u6574\u5408, \u8d28\u91cf\u8bc4\u4f30)", False, 8, DARK),
    # EvolveR: 完整维护管线(去重, 整合, 质量评估)
    ("", False, 2, DARK),
    ("\u5220\u9664\u4e0e\u8fc7\u671f:", True, 8, DARK),
    # 删除与过期:
    ("SAGE: \u57fa\u4e8e\u827e\u5bbe\u6d69\u65af\u9057\u5fd8\u66f2\u7ebf, \u672a\u88ab\u53cd\u590d\u5f3a\u5316\u7684\u8bb0\u5fc6\u81ea\u7136\u8870\u51cf", False, 8, DARK),
    # SAGE: 基于艾宾浩斯遗忘曲线, 未被反复强化的记忆自然衰减
    ("Mem0: \u652f\u6301TTL\u548cdecay\u673a\u5236", False, 8, DARK),
    # Mem0: 支持TTL和decay机制
    ("Copilot Memory: 28\u5929\u672a\u4f7f\u7528\u81ea\u52a8\u5220\u9664+\u5f15\u7528\u6709\u6548\u6027\u9a8c\u8bc1", False, 8, DARK),
    # Copilot Memory: 28天未使用自动删除+引用有效性验证
    ("\u5927\u591a\u6570\u7cfb\u7edf\u53ea\u589e\u4e0d\u5220, \u7ecf\u9a8c\u5e93\u65e0\u9650\u589e\u957f\u4ecd\u672a\u89e3\u51b3", False, 8, GRAY),
    # 大多数系统只增不删, 经验库无限增长仍未解决
])

s2_col3 = (GREEN_HDR, GREEN_BG,
    "\u5546\u4e1a\u4ea7\u54c1\u7684\u77e5\u8bc6\u6c89\u6dc0\u5b9e\u8df5",  # 商业产品的知识沉淀实践
    [
    ("GitHub Copilot Memory (2025.12\u516c\u6d4b)", True, 8, DARK),
    # GitHub Copilot Memory (2025.12公测)
    ("Copilot\u5728\u4f7f\u7528\u8fc7\u7a0b\u4e2d\u81ea\u52a8\u53d1\u73b0\u4ee3\u7801\u5e93\u4e2d\u7684\u5173\u952e\u6d1e\u5bdf", False, 8, DARK),
    # Copilot在使用过程中自动发现代码库中的关键洞察
    ("(\u5982\u6570\u636e\u5e93\u8fde\u63a5\u6a21\u5f0f, \u6587\u4ef6\u540c\u6b65\u89c4\u5219), \u5b58\u50a8\u65f6\u9644\u5e26\u4ee3\u7801", False, 8, DARK),
    # (如数据库连接模式, 文件同步规则), 存储时附带代码
    ("\u4f4d\u7f6e\u5f15\u7528(citations), \u4f7f\u7528\u524d\u9a8c\u8bc1\u5f15\u7528\u662f\u5426\u4ecd\u6709\u6548\u3002", False, 8, DARK),
    # 位置引用(citations), 使用前验证引用是否仍有效。
    ("\u8bb0\u5fc6\u53ef\u8de8\u529f\u80fd\u5171\u4eab: coding agent\u53d1\u73b0\u7684\u8bb0\u5fc6\u53ef\u88abcode review\u4f7f\u7528\u3002", False, 8, DARK),
    # 记忆可跨功能共享: coding agent发现的记忆可被code review使用。
    ("", False, 3, DARK),

    ("Claude Code + Auto Memory (Anthropic)", True, 8, DARK),
    ("\u56db\u5c42Markdown\u8bb0\u5fc6\u5c42\u6b21:", False, 8, DARK),
    # 四层Markdown记忆层次:
    ("Enterprise\u7ea7(\u7ba1\u7406\u5458MDM\u90e8\u7f72) / User\u7ea7(~/.claude/CLAUDE.md) /", False, 8, DARK),
    ("Project\u7ea7(./CLAUDE.md, Git\u7248\u672c\u63a7\u5236) / Sub-directory\u7ea7\u3002", False, 8, DARK),
    ("Auto Memory\u5728\u5de5\u4f5c\u8fc7\u7a0b\u4e2d\u81ea\u52a8\u8bb0\u5f55\u5b66\u4e60\u548c\u6a21\u5f0f\u3002", False, 8, DARK),
    # Auto Memory在工作过程中自动记录学习和模式。
    ("\u547d\u4ee4\u5f0f\u6307\u4ee4\u7684\u9075\u5faa\u7387(94%)\u8fdc\u9ad8\u4e8e\u63cf\u8ff0\u5f0f(73%)\u3002", False, 8, DARK),
    # 命令式指令的遵循率(94%)远高于描述式(73%)。
    ("", False, 3, DARK),

    ("Windsurf Cascade (Codeium)", True, 8, DARK),
    ("\u5b9e\u65f6\u52a8\u4f5c\u611f\u77e5: \u8ffd\u8e2a\u7528\u6237\u6240\u6709\u64cd\u4f5c(\u6587\u4ef6\u7f16\u8f91, \u7ec8\u7aef\u547d\u4ee4,", False, 8, DARK),
    # 实时动作感知: 追踪用户所有操作(文件编辑, 终端命令,
    ("\u526a\u8d34\u677f, \u5bf9\u8bdd\u5386\u53f2), \u63a8\u65ad\u610f\u56fe\u5e76\u5b9e\u65f6\u9002\u5e94\u3002", False, 8, DARK),
    # 剪贴板, 对话历史), 推断意图并实时适应。
    ("\u56db\u79cd\u6fc0\u6d3b\u6a21\u5f0f: Manual, Always On, Auto Attached, Model Decision\u3002", False, 8, DARK),
    # 四种激活模式: Manual, Always On, Auto Attached, Model Decision。
    ("", False, 3, DARK),

    ("Cursor + \u793e\u533aMemory Bank", True, 8, DARK),
    # Cursor + 社区Memory Bank
    ("Cursor\u672c\u8eab\u4e0d\u5177\u5907\u8de8\u4f1a\u8bdd\u81ea\u52a8\u8bb0\u5fc6\u80fd\u529b, \u4f46\u901a\u8fc7\u793e\u533a", False, 8, DARK),
    # Cursor本身不具备跨会话自动记忆能力, 但通过社区
    ("\u521b\u5efa\u7684Memory Bank\u6a21\u5f0f\u5f25\u8865\u3002\u5178\u578b\u7ed3\u6784\u5305\u542b:", False, 8, DARK),
    # 创建的Memory Bank模式弥补。典型结构包含:
    ("projectbrief.md, systemPatterns.md, techContext.md\u7b49\u6587\u4ef6\u3002", False, 8, DARK),
    ("", False, 3, DARK),
    ("AGENTS.md\u4f5c\u4e3a\u8de8AI\u5de5\u5177\u7684\u901a\u7528\u6307\u4ee4\u683c\u5f0f\u6b63\u5728\u5f62\u6210\u4e8b\u5b9e\u6807\u51c6,", False, 8, GRAY),
    # AGENTS.md作为跨AI工具的通用指令格式正在形成事实标准,
    ("Claude Code, Cursor, Copilot\u548cWindsurf\u5747\u5df2\u652f\u6301\u3002", False, 8, GRAY),
    # Claude Code, Cursor, Copilot和Windsurf均已支持。
])

s2_summary = [
    ("\u5173\u952e\u89c2\u5bdf: ", True, 9, WHITE),
    # 关键观察:
    ("\u4e03\u79cd\u5f62\u6001\u6784\u6210\u4ece\u4f4e\u5230\u9ad8\u7684\u62bd\u8c61\u8c31\u7cfb, \u4e0a\u4e0b\u6587\u4fdd\u7559\u9012\u51cf\u3001\u53ef\u590d\u7528\u6027\u9012\u589e\u3001\u521b\u5efa\u6210\u672c\u9012\u589e\u3002\u5546\u4e1a\u4ea7\u54c1\u5df2\u9a8c\u8bc1\"\u7ecf\u9a8c\u5230\u8bb0\u5fc6\u5230\u590d\u7528\"\u95ed\u73af\u53ef\u884c, \u4f46\u8bb0\u5fc6\u7c92\u5ea6", False, 9, WHITE),
    # 七种形态构成从低到高的抽象谱系, 上下文保留递减、可复用性递增、创建成本递增。商业产品已验证"经验到记忆到复用"闭环可行, 但记忆粒度
    ("\u505c\u7559\u5728\u7f16\u7801\u504f\u597d\u5c42\u9762\u3002\u751f\u547d\u5468\u671f\u7ba1\u7406\u662f\u6700\u8584\u5f31\u73af\u8282: \u5927\u591a\u6570\u7cfb\u7edf\u53ea\u589e\u4e0d\u5220, \u76ee\u524d\u6ca1\u6709\u4efb\u4f55\u5b66\u672f\u5de5\u4f5c\u5efa\u7acb\u4e86\u7cfb\u7edf\u7684\u77e5\u8bc6\u8d28\u91cf\u8bc4\u4f30\u548c\u6dd8\u6c70\u6846\u67b6\u3002", False, 9, WHITE),
    # 停留在编码偏好层面。生命周期管理是最薄弱环节: 大多数系统只增不删, 目前没有任何学术工作建立了系统的知识质量评估和淘汰框架。
]

make_slide_3col(prs,
    "\u77e5\u8bc6\u6c89\u6dc0\u7684\u5f62\u6001\u3001\u8303\u5f0f\u4e0e\u751f\u547d\u5468\u671f\u7ba1\u7406",
    # 知识沉淀的形态、范式与生命周期管理
    "\u57fa\u4e8e46\u7bc7\u5173\u952e\u8bba\u6587\u548c6\u4e2a\u5546\u4e1a\u4ea7\u54c1\u7684\u7cfb\u7edf\u68b3\u7406 (2023\u20132026)",
    # 基于46篇关键论文和6个商业产品的系统梳理 (2023-2026)
    [s2_col1, s2_col2, s2_col3],
    s2_summary)
print("Slide 2 added.")


# =====================================================================
# SLIDE 3: 沉淀知识的使用机制
# =====================================================================

s3_col1 = (BLUE_HDR, BLUE_BG,
    "\u68c0\u7d22\u5339\u914d\u673a\u5236",  # 检索匹配机制
    [
    ("\u77e5\u8bc6\u68c0\u7d22\u673a\u5236\u7ecf\u5386\u4e86\u663e\u8457\u7684\u590d\u6742\u5316\u6f14\u8fdb:", False, 8, DARK),
    # 知识检索机制经历了显著的复杂化演进:
    ("", False, 2, DARK),

    ("\u8bed\u4e49\u5d4c\u5165\u68c0\u7d22 (Voyager, A-MEM)", True, 8, DARK),
    # 语义嵌入检索 (Voyager, A-MEM)
    ("\u6700\u57fa\u7840\u4e5f\u6700\u666e\u904d\u7684\u65b9\u6cd5\u3002Voyager\u4f7f\u7528GPT-3.5\u751f\u6210\u6280\u80fd", False, 8, DARK),
    # 最基础也最普遍的方法。Voyager使用GPT-3.5生成技能
    ("\u63cf\u8ff0\u7684embedding\u8fdb\u884c\u76f8\u4f3c\u5ea6\u5339\u914d; A-MEM\u901a\u8fc7ChromaDB\u68c0\u7d22\u3002", False, 8, DARK),
    # 描述的embedding进行相似度匹配; A-MEM通过ChromaDB检索。
    ("", False, 2, DARK),

    ("\u6761\u4ef6\u89e6\u53d1\u68c0\u7d22 (Self-RAG, DRAGIN)", True, 8, DARK),
    # 条件触发检索 (Self-RAG, DRAGIN)
    ("Self-RAG (ICLR 2024 Oral): \u8bad\u7ec3LLM\u4f7f\u7528\u56db\u79cd\u53cd\u5c04token", False, 8, DARK),
    # Self-RAG (ICLR 2024 Oral): 训练LLM使用四种反射token
    ("([Retrieve], [IsRel], [IsSup], [IsUse])\u81ea\u9002\u5e94\u5224\u65ad\u4f55\u65f6", False, 8, DARK),
    # ([Retrieve], [IsRel], [IsSup], [IsUse])自适应判断何时
    ("\u9700\u8981\u68c0\u7d22\u4ee5\u53ca\u68c0\u7d22\u8d28\u91cf\u5982\u4f55\u3002DRAGIN\u57fa\u4e8etoken\u7ea7\u71b5\u4fe1\u53f7\u89e6\u53d1\u3002", False, 8, DARK),
    # 需要检索以及检索质量如何。DRAGIN基于token级熵信号触发。
    ("", False, 2, DARK),

    ("\u56fe\u68c0\u7d22 (GraphRAG, Think-on-Graph)", True, 8, DARK),
    # 图检索 (GraphRAG, Think-on-Graph)
    ("GraphRAG (Microsoft): \u4ece\u6587\u672c\u63d0\u53d6\u77e5\u8bc6\u56fe\u8c31\u5e76\u5206\u5c42\u7ed3\u6784\u5316\u3002", False, 8, DARK),
    # GraphRAG (Microsoft): 从文本提取知识图谱并分层结构化。
    ("Think-on-Graph: LLM\u4f5c\u4e3aAgent\u5728\u77e5\u8bc6\u56fe\u8c31\u4e0a\u8fed\u4ee3beam search\u3002", False, 8, DARK),
    # Think-on-Graph: LLM作为Agent在知识图谱上迭代beam search。
    ("", False, 2, DARK),

    ("\u4e0a\u4e0b\u6587\u5339\u914d (AutoGuide)", True, 8, DARK),
    # 上下文匹配 (AutoGuide)
    ("\u4e0d\u4f7f\u7528\u7b80\u5355\u8bed\u4e49\u76f8\u4f3c\u5ea6, \u800c\u662f\u5148\u7528Context Identification", False, 8, DARK),
    # 不使用简单语义相似度, 而是先用Context Identification
    ("Module\u8bc6\u522b\u5f53\u524d\u8f68\u8ff9\u7684\u4e0a\u4e0b\u6587\u7c7b\u578b, \u518d\u4ece\u5b57\u5178\u4e2d\u7cbe\u786e\u5339\u914d\u3002", False, 8, DARK),
    # Module识别当前轨迹的上下文类型, 再从字典中精确匹配。
    ("", False, 2, DARK),

    ("\u6df7\u5408\u68c0\u7d22 (Mem0, M-RAG)", True, 8, DARK),
    # 混合检索 (Mem0, M-RAG)
    ("Mem0\u7ed3\u5408\u5411\u91cf\u76f8\u4f3c\u5ea6\u68c0\u7d22\u548c\u56fe\u5b50\u56fe\u68c0\u7d22, \u5e76\u901a\u8fc7\u91cd\u6392\u5e8f\u4f18\u5316\u3002", False, 8, DARK),
    # Mem0结合向量相似度检索和图子图检索, 并通过重排序优化。
    ("M-RAG\u5c06\u77e5\u8bc6\u805a\u7c7b\u4e3a\u8bed\u4e49\u5206\u533a, \u7531\u53ccAgent\u9009\u62e9\u548c\u7cbe\u7ec6\u5185\u5bb9\u3002", False, 8, DARK),
    # M-RAG将知识聚类为语义分区, 由双Agent选择和精炼内容。
    ("", False, 2, DARK),

    ("\u51c6\u786e\u6027\u4fdd\u969c:", True, 8, DARK),
    # 准确性保障:
    ("Self-RAG\u7684\u53cd\u5c04token\u673a\u5236\u6700\u4e3a\u7cfb\u7edf\u5316; CRAG\u5f15\u5165\u68c0\u7d22\u8bc4\u4f30\u5668,", False, 8, DARK),
    # Self-RAG的反射token机制最为系统化; CRAG引入检索评估器,
    ("\u5bf9\u6bcf\u4e2aquery-document\u5bf9\u9884\u6d4b\u4e09\u7ea7\u7f6e\u4fe1\u5ea6(Correct/Incorrect/", False, 8, DARK),
    ("Ambiguous), \u5206\u522b\u89e6\u53d1\u7cbe\u7ec6, Web\u641c\u7d22\u6216\u4e24\u8005\u7ed3\u5408\u3002", False, 8, DARK),
    # Ambiguous), 分别触发精炼, Web搜索或两者结合。
])

s3_col2 = (PURPLE_HDR, PURPLE_BG,
    "\u77e5\u8bc6\u6ce8\u5165\u4e0e\u51b2\u7a81\u89e3\u51b3",  # 知识注入与冲突解决
    [
    ("\u6ce8\u5165\u65b9\u5f0f", True, 9, DARK),
    # 注入方式
    ("", False, 2, DARK),

    ("In-context\u793a\u4f8b\u6ce8\u5165 (\u6700\u4e3b\u6d41)", True, 8, DARK),
    # In-context示例注入 (最主流)
    ("\u51e0\u4e4e\u6240\u6709\u7cfb\u7edf(Voyager, ExpeL, AWM, Reflexion, A-MEM)", False, 8, DARK),
    # 几乎所有系统(Voyager, ExpeL, AWM, Reflexion, A-MEM)
    ("\u5747\u91c7\u7528: \u5c06\u68c0\u7d22\u5230\u7684\u8bb0\u5fc6\u76f4\u63a5\u653e\u5165prompt\u4f5c\u4e3afew-shot\u793a\u4f8b\u3002", False, 8, DARK),
    # 均采用: 将检索到的记忆直接放入prompt作为few-shot示例。
    ("\u7b80\u5355\u76f4\u63a5, LLM\u539f\u751f\u652f\u6301, \u4f46\u53d7\u9650\u4e8e\u4e0a\u4e0b\u6587\u7a97\u53e3\u5927\u5c0f\u3002", False, 8, GRAY),
    # 简单直接, LLM原生支持, 但受限于上下文窗口大小。
    ("", False, 2, DARK),

    ("\u7cfb\u7edf\u63d0\u793a\u6ce8\u5165 (\u59cb\u7ec8\u5728\u573a\u7684\u6838\u5fc3\u77e5\u8bc6)", True, 8, DARK),
    # 系统提示注入 (始终在场的核心知识)
    ("MemGPT: Core Memory\u59cb\u7ec8\u653e\u5728system prompt\u4e2d\u3002", False, 8, DARK),
    # MemGPT: Core Memory始终放在system prompt中。
    ("Claude Code: MEMORY.md\u524d200\u884c\u81ea\u52a8\u52a0\u8f7d\u5230\u6bcf\u6b21\u4f1a\u8bdd\u3002", False, 8, DARK),
    # Claude Code: MEMORY.md前200行自动加载到每次会话。
    ("", False, 2, DARK),

    ("\u5de5\u5177\u8c03\u7528\u6ce8\u5165 (\u6309\u9700\u7075\u6d3b\u8bbf\u95ee)", True, 8, DARK),
    # 工具调用注入 (按需灵活访问)
    ("MemGPT: LLM\u901a\u8fc7function calling\u5bf9archival/recall storage\u5206\u9875\u67e5\u8be2\u3002", False, 8, DARK),
    # MemGPT: LLM通过function calling对archival/recall storage分页查询。
    ("AWM: workflow\u6269\u5c55\u4e3aAgent\u52a8\u4f5c\u7a7a\u95f4\u4e2d\u7684\u53ef\u8c03\u7528\u51fd\u6570\u3002", False, 8, DARK),
    # AWM: workflow扩展为Agent动作空间中的可调用函数。
    ("MCP (Anthropic 2024): AI\u4e0e\u5916\u90e8\u6570\u636e\u6e90\u7684\u901a\u7528\u8fde\u63a5\u6807\u51c6\u3002", False, 8, DARK),
    # MCP (Anthropic 2024): AI与外部数据源的通用连接标准。
    ("", False, 2, DARK),

    ("\u53c2\u6570\u5316\u6ce8\u5165 (\u65b0\u5174\u65b9\u5411)", True, 8, DARK),
    # 参数化注入 (新兴方向)
    ("Parametric RAG (SIGIR 2025): \u5c06\u5916\u90e8\u77e5\u8bc6\u76f4\u63a5\u6ce8\u5165LLM\u7684", False, 8, DARK),
    # Parametric RAG (SIGIR 2025): 将外部知识直接注入LLM的
    ("\u524d\u9988\u7f51\u7edc(FFN), \u907f\u514d\u4e0a\u4e0b\u6587\u7a97\u53e3\u9650\u5236\u3002", False, 8, DARK),
    # 前馈网络(FFN), 避免上下文窗口限制。
    ("", False, 3, DARK),

    ("\u77e5\u8bc6\u51b2\u7a81\u68c0\u6d4b\u4e0e\u89e3\u51b3", True, 9, DARK),
    # 知识冲突检测与解决
    ("\u4e09\u7c7b\u51b2\u7a81: \u4e0a\u4e0b\u6587-\u8bb0\u5fc6\u51b2\u7a81, \u4e0a\u4e0b\u6587\u95f4\u51b2\u7a81, \u8bb0\u5fc6\u5185\u51b2\u7a81\u3002", False, 8, DARK),
    # 三类冲突: 上下文-记忆冲突, 上下文间冲突, 记忆内冲突。
    ("ConflictBank (NeurIPS 2024): 286\u4e07\u6761\u58f0\u660e\u7684\u57fa\u51c6\u6570\u636e\u96c6,", False, 8, DARK),
    # ConflictBank (NeurIPS 2024): 286万条声明的基准数据集,
    ("\u53d1\u73b0LLM\u9ad8\u5ea6\u63a5\u53d7\u5916\u90e8\u8bc1\u636e, \u5b58\u5728\u88ab\u8bef\u5bfc\u7684\u98ce\u9669\u3002", False, 8, DARK),
    # 发现LLM高度接受外部证据, 存在被误导的风险。
    ("KARMA (NeurIPS 2025): \u4e13\u95e8\u7684\u51b2\u7a81\u89e3\u51b3Agent, LLM\u8fa9\u8bba\u673a\u5236\u3002", False, 8, DARK),
    # KARMA (NeurIPS 2025): 专门的冲突解决Agent, LLM辩论机制。
])

s3_col3 = (RED_HDR, RED_BG,
    "\u4eba-Code Agent\u4ea4\u4e92\u4e2d\u7684\u7ecf\u9a8c\u4f7f\u7528",  # 人-Code Agent交互中的经验使用
    [
    ("\u53cd\u9988\u95ed\u73af\u9a71\u52a8\u7684\u7ecf\u9a8c\u79ef\u7d2f\u5448\u4e09\u7ea7\u8fdb\u5316", True, 9, DARK),
    # 反馈闭环驱动的经验积累呈三级进化
    ("", False, 2, DARK),

    ("\u7b2c\u4e00\u9636\u6bb5: \u5355\u8f6e\u9a8c\u8bc1\u53cd\u9988", True, 8, DARK),
    # 第一阶段: 单轮验证反馈
    ("SpecGen (ICSE 2025): LLM\u751f\u6210JML\u89c4\u683c, \u5f62\u5f0f\u5316\u9a8c\u8bc1\u5668\u68c0\u67e5,", False, 8, DARK),
    # SpecGen (ICSE 2025): LLM生成JML规格, 形式化验证器检查,
    ("\u5931\u8d25\u4fe1\u606f\u53cd\u9988\u5230LLM\u91cd\u65b0\u751f\u6210\u3002385\u7a0b\u5e8f\u4e2d279\u4e2a\u6210\u529f\u3002", False, 8, DARK),
    # 失败信息反馈到LLM重新生成。385程序中279个成功。
    ("Voyager: \u4ee3\u7801\u751f\u6210, \u6267\u884c, \u9519\u8bef\u53cd\u9988, \u4fee\u590d, \u6210\u529f\u540e\u5165\u5e93\u3002", False, 8, DARK),
    # Voyager: 代码生成, 执行, 错误反馈, 修复, 成功后入库。
    ("\u5b8c\u5168\u81ea\u52a8\u4f46\u4e0d\u8de8\u4f1a\u8bdd, \u7ecf\u9a8c\u4e0d\u6301\u4e45\u3002", False, 8, GRAY),
    # 完全自动但不跨会话, 经验不持久。
    ("", False, 2, DARK),

    ("\u7b2c\u4e8c\u9636\u6bb5: \u8de8\u4efb\u52a1\u7ecf\u9a8c\u63d0\u70bc", True, 8, DARK),
    # 第二阶段: 跨任务经验提炼
    ("ExpeL (AAAI 2024): \u4ece\u6210\u529f/\u5931\u8d25\u8f68\u8ff9\u5bf9\u4e2d\u63d0\u53d6\u901a\u7528\u6d1e\u5bdf,", False, 8, DARK),
    # ExpeL (AAAI 2024): 从成功/失败轨迹对中提取通用洞察,
    ("\u5b9e\u73b0\u8de8\u4efb\u52a1\u8fc1\u79fb\u5b66\u4e60(\u5982\u4eceHotpotQA\u5230FEVER\u7684\u8fc1\u79fb)\u3002", False, 8, DARK),
    # 实现跨任务迁移学习(如从HotpotQA到FEVER的迁移)。
    ("EvolveR: \u5b8c\u6574\"\u6536\u96c6, \u84b8\u998f, \u68c0\u7d22, \u8fdb\u5316\"\u95ed\u73af,", False, 8, DARK),
    # EvolveR: 完整"收集, 蒸馏, 检索, 进化"闭环,
    ("\u901a\u8fc7GRPO\u5f3a\u5316\u5b66\u4e60\u66f4\u65b0Agent\u7b56\u7565\u53c2\u6570\u3002", False, 8, DARK),
    # 通过GRPO强化学习更新Agent策略参数。
    ("", False, 2, DARK),

    ("\u7b2c\u4e09\u9636\u6bb5: \u7ec4\u7ec7\u7ea7\u7ecf\u9a8c\u6cbb\u7406", True, 8, DARK),
    # 第三阶段: 组织级经验治理
    ("MemGovern (\u5f00\u521b\u65b9\u5411): \u5c06GitHub\u4e0a\u975e\u7ed3\u6784\u5316\u7684Issue/PR", False, 8, DARK),
    # MemGovern (开创方向): 将GitHub上非结构化的Issue/PR
    ("\u6570\u636e\u8f6c\u5316\u4e3a\u7ed3\u6784\u5316\u7ecf\u9a8c\u5361\u7247, \u901a\u8fc7\u5c42\u6b21\u5316\u7b5b\u9009\u8fc7\u6ee4\u4f4e\u8d28\u91cf\u5b9e\u4f8b,", False, 8, DARK),
    # 数据转化为结构化经验卡片, 通过层次化筛选过滤低质量实例,
    ("\u5f15\u5165Agentic Experience Search(Search-then-Browse\u6a21\u5f0f)\u3002", False, 8, DARK),
    # 引入Agentic Experience Search(Search-then-Browse模式)。
    ("SWE-bench Verified\u4e0a\u5e73\u5747\u63d0\u53474.65%\u89e3\u51b3\u7387\u3002", False, 8, DARK),
    # SWE-bench Verified上平均提升4.65%解决率。
    ("\u4f46\u7ecf\u9a8c\u6765\u6e90\u662f\u5386\u53f2\u7684\u3001\u5916\u90e8\u7684\u6570\u636e, \u800c\u975e\u5f53\u524d\u7528\u6237\u4ea4\u4e92\u3002", False, 8, GRAY),
    # 但经验来源是历史的、外部的数据, 而非当前用户交互。
    ("", False, 2, DARK),

    ("\u7528\u6237\u53c2\u4e0e\u77e5\u8bc6\u7ef4\u62a4\u7684\u4e09\u79cd\u6a21\u5f0f:", True, 8, DARK),
    # 用户参与知识维护的三种模式:
    ("\u4e3b\u52a8\u544a\u77e5: \u7528\u6237\u76f4\u63a5\u8bf4\"remember that we use pnpm\"", False, 8, DARK),
    # 主动告知: 用户直接说"remember that we use pnpm"
    ("\u5ba1\u6838\u5220\u9664: \u7528\u6237\u4f5c\u4e3a\"\u7b56\u5c55\u4eba\"\u800c\u975e\"\u521b\u4f5c\u8005\"\u53c2\u4e0e", False, 8, DARK),
    # 审核删除: 用户作为"策展人"而非"创作者"参与
    ("\u534f\u4f5c\u7cbe\u7ec6: HULA (Atlassian)\u591a\u9636\u6bb5\u4eba\u5de5\u7cbe\u7ec6\u548c\u5f15\u5bfc", False, 8, DARK),
    # 协作精炼: HULA (Atlassian)多阶段人工精炼和引导
])

s3_summary = [
    ("\u6838\u5fc3\u7f3a\u5931: ", True, 9, WHITE),
    # 核心缺失:
    ("\u76ee\u524d\u6ca1\u6709\u4efb\u4f55\u7cfb\u7edf\u5b9e\u73b0\u4e86\"\u7528\u6237\u9a8c\u8bc1\u8d28\u91cf \u2192 \u53cd\u9988\u5199\u5165 \u2192 \u81ea\u52a8\u8c03\u6574\u672a\u6765\u7b56\u7565\"\u7684\u95ed\u73af\u3002\u7528\u6237\u53cd\u9988\u4ec5\u505c\u7559\u5728\u7f16\u7801\u504f\u597d\u5c42\u9762, \u672a\u6df1\u5165\u5230", False, 9, WHITE),
    # 目前没有任何系统实现了"用户验证质量 → 反馈写入 → 自动调整未来策略"的闭环。用户反馈仅停留在编码偏好层面, 未深入到
    ("\u529f\u80fd\u610f\u56fe\u548c\u8bbe\u8ba1\u51b3\u7b56\u7684\u7ecf\u9a8c\u79ef\u7d2f\u3002\u77e5\u8bc6\u51b2\u7a81\u68c0\u6d4b\u4e0e\u89e3\u51b3\u662f\u4e00\u4e2a\u5df2\u8bc6\u522b\u4f46\u8fdc\u672a\u89e3\u51b3\u7684\u96be\u9898, \u65f6\u6548\u6027\u4fdd\u8bc1\u4f9d\u8d56\u52a8\u6001\u66f4\u65b0\u548c\u8fc7\u671f\u673a\u5236\u3002", False, 9, WHITE),
    # 功能意图和设计决策的经验积累。知识冲突检测与解决是一个已识别但远未解决的难题, 时效性保证依赖动态更新和过期机制。
]

make_slide_3col(prs,
    "\u6c89\u6dc0\u77e5\u8bc6\u7684\u4f7f\u7528\u673a\u5236: \u68c0\u7d22\u3001\u6ce8\u5165\u4e0e\u4eba\u673a\u534f\u540c",
    # 沉淀知识的使用机制: 检索、注入与人机协同
    "\u5982\u4f55\u66f4\u51c6\u786e\u3001\u66f4\u9ad8\u6548\u5730\u5728\u4eba-Code Agent\u4ea4\u4e92\u4e2d\u4f7f\u7528\u6c89\u6dc0\u7684\u7ecf\u9a8c",
    # 如何更准确、更高效地在人-Code Agent交互中使用沉淀的经验
    [s3_col1, s3_col2, s3_col3],
    s3_summary)
print("Slide 3 added.")

out = v11b.replace('v11b.pptx', 'v11c.pptx')
prs.save(out)
print(f"Done! Saved to: {out}")
