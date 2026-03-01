# -*- coding: utf-8 -*-
"""
Complete the summary bar text on slide 3 of v11c.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import glob, os

DARK  = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = '\u5fae\u8f6f\u96c5\u9ed1'

base = r'C:\Users\Zhifei Dou\Desktop\Openspec_modified\slides'
v11c = [f for f in glob.glob(os.path.join(base, '*v11c.pptx')) if '~$' not in f][0]
print(f"Reading: {v11c}")
prs = Presentation(v11c)

slide = prs.slides[2]  # Slide 3

# Find the summary text shape on slide 3
# It's the textbox at approximately (0.36, 6.71) containing "总结"
summary_tb = None
for shape in slide.shapes:
    if shape.has_text_frame:
        t = round(shape.top / 914400, 2)
        if t > 6.5 and shape.text_frame.text.strip().startswith('\u603b\u7ed3'):
            summary_tb = shape
            break

if not summary_tb:
    print("ERROR: Summary textbox not found on slide 3")
    sys.exit(1)

print(f"Found summary at ({round(summary_tb.left/914400,2)}, {round(summary_tb.top/914400,2)})")
print(f"Current text: {summary_tb.text_frame.text[:80]}...")

# Clear existing paragraphs
tf = summary_tb.text_frame
while len(tf.paragraphs) > 1:
    tf._txBody.remove(tf.paragraphs[-1]._p)
for r in list(tf.paragraphs[0].runs):
    tf.paragraphs[0]._p.remove(r._r)

# Also resize the textbox height to fit more text
summary_tb.height = Inches(0.60)

# Write new summary content
# 总结：知识的沉淀和复用潜在可以为人-Code Agent协同的系统工程框架...
summary_lines = [
    ("\u603b\u7ed3: ", True, 9, WHITE),
    # 总结:

    ("\u77e5\u8bc6\u7684\u6c89\u6dc0\u4e0e\u590d\u7528\u6f5c\u5728\u53ef\u4ee5\u4e3a\u4eba-Code Agent\u534f\u540c\u7684\u7cfb\u7edf\u5de5\u7a0b\u6846\u67b6\u63d0\u4f9b\u4e09\u65b9\u9762\u652f\u6491: "
     "\u68c0\u7d22\u4e0e\u6ce8\u5165\u673a\u5236\u4f7f\u9886\u57df\u6807\u51c6\u3001\u5386\u53f2Spec\u6a21\u5f0f\u7b49\u77e5\u8bc6\u80fd\u5728\u7f16\u5199\u548c\u5b9e\u73b0\u9636\u6bb5\u88ab\u81ea\u52a8\u5339\u914d\u548c\u590d\u7528; "
     "\u53cd\u9988\u95ed\u73af\u4f7fSpec\u2192\u5b9e\u73b0\u2192\u6d4b\u8bd5\u7684\u7ed3\u679c\u56de\u6d41\u4e3aSpec\u7f16\u5199\u8d28\u91cf\u7684\u7ecf\u9a8c\u79ef\u7d2f; "
     "\u751f\u547d\u5468\u671f\u7ba1\u7406\u4f7f\u8fc7\u65f6\u7684\u7ecf\u9a8c\u88ab\u53ca\u65f6\u6dd8\u6c70\u3001\u51b2\u7a81\u7684\u77e5\u8bc6\u88ab\u68c0\u6d4b\u548c\u89e3\u51b3\u3002",
     False, 9, WHITE),
    # 知识的沉淀与复用潜在可以为人-Code Agent协同的系统工程框架提供三方面支撑:
    # 检索与注入机制使领域标准、历史Spec模式等知识能在编写和实现阶段被自动匹配和复用;
    # 反馈闭环使Spec→实现→测试的结果回流为Spec编写质量的经验积累;
    # 生命周期管理使过时的经验被及时淘汰、冲突的知识被检测和解决。

    ("\u4f46\u73b0\u6709\u68c0\u7d22\u548c\u6ce8\u5165\u673a\u5236\u4e3b\u8981\u9762\u5411\u4ee3\u7801\u7247\u6bb5, "
     "\u5bf9\u89c4\u683c\u9700\u6c42\u3001\u8bbe\u8ba1\u51b3\u7b56\u3001\u9886\u57df\u6807\u51c6\u7b49\u975e\u4ee3\u7801\u77e5\u8bc6\u7684\u652f\u6301\u4ecd\u662f\u7a7a\u767d, "
     "\u8fd9\u662f\u5c06\u77e5\u8bc6\u6c89\u6dc0\u878d\u5165\u7cfb\u7edf\u5de5\u7a0b\u6846\u67b6\u9700\u8981\u9996\u5148\u89e3\u51b3\u7684\u95ee\u9898\u3002",
     False, 9, WHITE),
    # 但现有检索和注入机制主要面向代码片段,
    # 对规格需求、设计决策、领域标准等非代码知识的支持仍是空白,
    # 这是将知识沉淀融入系统工程框架需要首先解决的问题。
]

first = True
for text, bold, sz, clr in summary_lines:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FONT

out = v11c.replace('v11c.pptx', 'v11d.pptx')
prs.save(out)
print(f"Done! Summary on slide 3 updated. Saved to: {out}")
