"""
Add Experiment 4 (Testing Workflow) to Slide 4 of the PPT.
Replaces the bottom summary score boxes with the testing experiment block.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation('OpenSpec_技术研究汇报_精简版.pptx')
slide = prs.slides[3]  # Slide 4 (0-indexed)

# === Step 1: Remove summary score boxes (the big-number area at bottom) ===
shapes_to_remove = {
    # Exp1 summary box
    'Rectangle 172', 'TextBox 173', 'TextBox 174', 'TextBox 175',
    'TextBox 176', 'TextBox 177', 'TextBox 178',
    # Exp2 summary box
    'Rectangle 179', 'TextBox 180', 'TextBox 181',
    'TextBox 182', 'TextBox 183',
    # Exp3 summary box
    'Rectangle 184', 'TextBox 185', 'TextBox 186', 'TextBox 187',
    'TextBox 188', 'TextBox 189', 'TextBox 190',
}

sp_tree = slide.shapes._spTree
for shape in list(slide.shapes):
    if shape.name in shapes_to_remove:
        sp_tree.remove(shape._element)

# === Step 2: Update conclusion text (三组 -> 四组) ===
for shape in slide.shapes:
    if shape.name == 'TextBox 192':
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if '三组' in run.text:
                    run.text = run.text.replace('三组', '四组')

# === Color constants (matching existing slide style) ===
DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF0, 0xF0, 0xF0)
VERY_LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
MID_GRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
BLUE_ACCENT = RGBColor(0x00, 0x78, 0xD4)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

def no_line(shape):
    """Remove border line from shape."""
    shape.line.fill.background()

def make_rect(slide, left, top, width, height, fill_rgb):
    """Add a rectangle with solid fill and no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    no_line(shape)
    return shape

def make_rounded_rect(slide, left, top, width, height, fill_rgb):
    """Add a rounded rectangle with solid fill and no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    no_line(shape)
    return shape

def make_textbox(slide, left, top, width, height, text, font_size_pt, color,
                 bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box with a single run."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txbox

def make_multi_run_textbox(slide, left, top, width, height, runs_data, align=PP_ALIGN.LEFT):
    """Add a text box with multiple styled runs.
    runs_data: list of (text, font_size_pt, color, bold, font_name)
    """
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, color, bold, name in runs_data:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = name or "Microsoft YaHei"
    return txbox

# === Step 3: Add Experiment 4 block ===
# Position: same y as removed summary boxes
EXP4_LEFT = 230000
EXP4_TOP = 4401810
EXP4_FULL_WIDTH = 11500000

# --- 3a: Header bar ---
HEADER_H = 270000
make_rect(slide, EXP4_LEFT, EXP4_TOP, EXP4_FULL_WIDTH, HEADER_H, DARK_BLUE)
make_multi_run_textbox(
    slide, EXP4_LEFT, EXP4_TOP + 5000, EXP4_FULL_WIDTH, 260000,
    [
        ("实验4：Dashboard — 规格驱动测试", 11, WHITE, True, "Microsoft YaHei"),
        ("  (OS-HW Testing Workflow: gen-tests → run-tests)", 9, RGBColor(0xBB, 0xBB, 0xBB), False, "Microsoft YaHei"),
    ],
    align=PP_ALIGN.CENTER
)

# --- 3b: Description bar ---
DESC_TOP = EXP4_TOP + HEADER_H + 10000
DESC_H = 180000
make_rect(slide, EXP4_LEFT, DESC_TOP, EXP4_FULL_WIDTH, DESC_H, LIGHT_GRAY_BG)
make_textbox(
    slide, EXP4_LEFT + 50000, DESC_TOP + 10000, EXP4_FULL_WIDTH - 100000, 160000,
    "实验内容：使用 OS-HW 测试工作流从 Use Case 规格自动生成测试用例（gen-tests），执行测试并追踪 Happy / Alt / Exception 三类路径覆盖率（run-tests），全流程自动化验证规格一致性",
    8, DARK_GRAY
)

# --- 3c: Left side - Coverage table (manual shapes, matching existing style) ---
TABLE_TOP = DESC_TOP + DESC_H + 15000
TABLE_LEFT = EXP4_LEFT

# Table dimensions
COL_WIDTHS = [1800000, 700000, 700000, 800000, 700000]
ROW_H = 220000
HEADER_ROW_H = 240000
TABLE_W = sum(COL_WIDTHS)  # 4700000

# Header row cells
headers = ["用例", "Happy", "Alt", "Exception", "覆盖率"]
x = TABLE_LEFT
for i, (hdr, cw) in enumerate(zip(headers, COL_WIDTHS)):
    make_rect(slide, x, TABLE_TOP, cw, HEADER_ROW_H, DARK_BLUE)
    make_textbox(slide, x, TABLE_TOP + 5000, cw, HEADER_ROW_H - 10000,
                 hdr, 8, WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += cw

# Data rows
data_rows = [
    ("UC-01: Drag Widget",    "✅ 1/1", "✅ 2/2", "❌ 0/1",  "75%"),
    ("UC-02: Snap Widget",    "✅ 1/1", "✅ 1/1", "✅ 1/1", "100%"),
    ("UC-03: Resize Widget",  "✅ 1/1", "✅ 1/1", "✅ 1/1", "100%"),
    ("UC-04: Add Widget",     "✅ 1/1", "✅ 1/1", "✅ 1/1", "100%"),
    ("UC-05: Remove Widget",  "✅ 1/1", "✅ 1/1", "✅ 1/1", "100%"),
]

for r_idx, row in enumerate(data_rows):
    row_top = TABLE_TOP + HEADER_ROW_H + r_idx * ROW_H
    bg_color = ROW_GRAY if r_idx % 2 == 0 else ROW_WHITE
    x = TABLE_LEFT
    for c_idx, (val, cw) in enumerate(zip(row, COL_WIDTHS)):
        make_rect(slide, x, row_top, cw, ROW_H, bg_color)
        # Determine text color
        if "❌" in val:
            txt_color = RED
        elif val == "75%":
            txt_color = ORANGE
        elif "✅" in val or val == "100%":
            txt_color = GREEN
        else:
            txt_color = MID_GRAY
        align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
        left_pad = 30000 if c_idx == 0 else 0
        make_textbox(slide, x + left_pad, row_top + 3000, cw - left_pad, ROW_H - 6000,
                     val, 7.5, txt_color, bold=(c_idx > 0), align=align)
        x += cw

# Total row
total_top = TABLE_TOP + HEADER_ROW_H + len(data_rows) * ROW_H
total_vals = ["总计 (15/17)", "✅ 5/5", "✅ 6/6", "4/5", "88%"]
x = TABLE_LEFT
for c_idx, (val, cw) in enumerate(zip(total_vals, COL_WIDTHS)):
    make_rect(slide, x, total_top, cw, HEADER_ROW_H, VERY_LIGHT_GREEN)
    if val == "88%":
        txt_color = GREEN
    elif "4/5" in val:
        txt_color = ORANGE
    elif "✅" in val:
        txt_color = GREEN
    else:
        txt_color = DARK_BLUE
    align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
    left_pad = 30000 if c_idx == 0 else 0
    make_textbox(slide, x + left_pad, total_top + 5000, cw - left_pad, HEADER_ROW_H - 10000,
                 val, 8, txt_color, bold=True, align=align)
    x += cw

# --- 3d: Right side - Key metrics panel ---
METRIC_LEFT = TABLE_LEFT + TABLE_W + 200000
METRIC_W = EXP4_FULL_WIDTH - TABLE_W - 200000  # remaining width
METRIC_TOP = TABLE_TOP
METRIC_TOTAL_H = HEADER_ROW_H + len(data_rows) * ROW_H + HEADER_ROW_H

# Background
make_rounded_rect(slide, METRIC_LEFT, METRIC_TOP, METRIC_W, METRIC_TOTAL_H, RGBColor(0xFA, 0xFA, 0xFA))

# Big metric: 88%
make_textbox(slide, METRIC_LEFT, METRIC_TOP + 30000, METRIC_W, 350000,
             "88%", 36, GREEN, bold=True, align=PP_ALIGN.CENTER)
make_textbox(slide, METRIC_LEFT, METRIC_TOP + 330000, METRIC_W, 150000,
             "规格路径覆盖率", 10, DARK_GRAY, align=PP_ALIGN.CENTER)
make_textbox(slide, METRIC_LEFT, METRIC_TOP + 460000, METRIC_W, 100000,
             "(15 / 17 paths covered)", 8, MID_GRAY, align=PP_ALIGN.CENTER)

# Divider line (thin rectangle)
divider_y = METRIC_TOP + 580000
make_rect(slide, METRIC_LEFT + 200000, divider_y, METRIC_W - 400000, 15000, RGBColor(0xE0, 0xE0, 0xE0))

# Second row metrics
m2_top = divider_y + 40000
# Left metric: 44/44
m2_half_w = METRIC_W // 2
make_textbox(slide, METRIC_LEFT, m2_top, m2_half_w, 200000,
             "44/44", 20, GREEN, bold=True, align=PP_ALIGN.CENTER)
make_textbox(slide, METRIC_LEFT, m2_top + 180000, m2_half_w, 100000,
             "测试全部通过", 8, DARK_GRAY, align=PP_ALIGN.CENTER)

# Right metric: 20 mappings
make_textbox(slide, METRIC_LEFT + m2_half_w, m2_top, m2_half_w, 200000,
             "20", 20, BLUE_ACCENT, bold=True, align=PP_ALIGN.CENTER)
make_textbox(slide, METRIC_LEFT + m2_half_w, m2_top + 180000, m2_half_w, 100000,
             "需求追踪映射", 8, DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom info line
info_y = m2_top + 300000
make_textbox(slide, METRIC_LEFT, info_y, METRIC_W, 80000,
             "7 test files  ·  vitest v2.1.9  ·  3.16s runtime", 7, MID_GRAY, align=PP_ALIGN.CENTER)

# Spec traceability note
trace_y = info_y + 80000
make_textbox(slide, METRIC_LEFT, trace_y, METRIC_W, 80000,
             "spec-tests.md: 14 ✅ 完整  ·  4 ⚠️ 部分  ·  2 ❌ 待补", 7, DARK_GRAY, align=PP_ALIGN.CENTER)

# === Save ===
prs.save('OpenSpec_技术研究汇报_精简版.pptx')
print("Done! Slide 4 updated with Experiment 4 (Testing Workflow).")
