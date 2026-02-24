"""
Rebuild PPT slides to consolidate two separate binary comparisons into
a unified three-way comparison: Vanilla Claude Code vs 原生 OpenSpec vs OpenSpec-HW2.

Based on ORIGINAL 50-slide PPT from git.

Changes:
  - Slides 9-13  : rebuilt as three-way main body (was slides 9-15)
  - Slides 29-33 : rebuilt as three-way appendix  (was slides 29-40)
  - Slides 14-15, 34-40: deleted (content merged)
  - TOC (slide 2) and Summary (slide 28): text updated
  - All page numbers updated
  - Final total: 50 - 9 = 41 slides
"""
import copy, re, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

PPTX = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报.pptx"
prs = Presentation(PPTX)
SW = prs.slide_width   # 12192000
SH = prs.slide_height  # 6858000

# ── Colours ──────────────────────────────────────────────────────
DARK_BG  = RGBColor(0x16, 0x21, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT  = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MED = RGBColor(0x99, 0x99, 0x99)
BLACK    = RGBColor(0x33, 0x33, 0x33)
BLUE     = RGBColor(0x00, 0x78, 0xD4)
GREEN    = RGBColor(0x10, 0x7C, 0x10)
ORANGE   = RGBColor(0xE8, 0x6C, 0x00)
PURPLE   = RGBColor(0x6B, 0x69, 0xD6)
TEAL     = RGBColor(0x00, 0x8B, 0x8B)
RED      = RGBColor(0xD1, 0x34, 0x38)

COL_VANILLA = RGBColor(0x78, 0x78, 0x78)   # gray header for Vanilla
COL_NATIVE  = BLUE                           # blue header for 原生
COL_HW2     = GREEN                          # green header for HW2
COL_DIM     = RGBColor(0x44, 0x44, 0x44)    # dark gray for dimension col

LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)

print(f"Opened PPT: {len(prs.slides)} slides")

# ── Helper functions ─────────────────────────────────────────────
def clear_slide(slide):
    """Remove all shapes from a slide, preserving background."""
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

def add_textbox(slide, left, top, width, height, text, size=Pt(14),
                color=BLACK, bold=False, align=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_rich_textbox(slide, left, top, width, height, paragraphs_data):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pd in enumerate(paragraphs_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get('align', PP_ALIGN.LEFT)
        p.space_after = pd.get('space_after', Pt(4))
        txt = pd.get('text', '')
        if pd.get('bullet'):
            txt = '• ' + txt
        run = p.add_run()
        run.text = txt
        run.font.size = pd.get('size', Pt(13))
        run.font.color.rgb = pd.get('color', BLACK)
        run.font.bold = pd.get('bold', False)
        run.font.name = 'Microsoft YaHei'
    return txBox

def add_accent_bar(slide, left, top, width, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_colored_box(slide, left, top, width, height, fill_color, border_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box

def add_page_num(slide, page, total, color=GRAY_MED):
    add_textbox(slide, Emu(SW - Inches(1.2)), Emu(SH - Inches(0.5)),
                Inches(1), Inches(0.4), f"{page}/{total}",
                size=Pt(11), color=color, align=PP_ALIGN.RIGHT)

def set_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def style_cell(cell, text, font_size=Pt(11), bold=False, color=BLACK,
               align=PP_ALIGN.CENTER, fill_color=None, font_name='Microsoft YaHei'):
    cell.text = ''
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Reduce cell margins for compactness
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.name = font_name
    run.font.bold = bold
    run.font.color.rgb = color
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

def create_comparison_table(slide, left, top, width, height, data,
                            col_widths=None, font_size=Pt(11)):
    """Create a styled 4-column comparison table.
    data: list of lists; row 0 = header, last row may be totals."""
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    header_colors = [COL_DIM, COL_VANILLA, COL_NATIVE, COL_HW2]
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            txt = str(data[r][c])
            is_header = (r == 0)
            is_total  = (r == rows - 1 and '合计' in str(data[r][0]))
            al = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            if is_header:
                fc = header_colors[c] if c < len(header_colors) else BLUE
                style_cell(cell, txt, font_size=font_size, bold=True,
                           color=WHITE, align=al, fill_color=fc)
            elif is_total:
                style_cell(cell, txt, font_size=font_size, bold=True,
                           color=BLACK, align=al,
                           fill_color=RGBColor(0xE0, 0xE0, 0xE0))
            else:
                bg = RGBColor(0xF8, 0xF8, 0xF8) if r % 2 == 0 else None
                style_cell(cell, txt, font_size=font_size, bold=False,
                           color=BLACK, align=al, fill_color=bg)
    return tbl

# ── TOTAL after changes (50 - 9 deleted = 41) ───────────────────
TOTAL = 41

# ═══════════════════════════════════════════════════════════════════
# SECTION A : Rebuild Main Body Slides 9-13  (0-indexed 8-12)
# ═══════════════════════════════════════════════════════════════════

# ── Slide 9: Three-way experiment setup ──────────────────────────
s5 = prs.slides[8]
clear_slide(s5)

add_textbox(s5, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '实测对比：同一需求，三种方法的量化差异',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s5, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

add_textbox(s5, Inches(0.5), Inches(0.85), Inches(11), Inches(0.35),
            '测试项目：A股有色金属量化选股交易系统 | 相同需求描述，分别用三种方式让 Claude Code 实现',
            size=Pt(12), color=GRAY_MED)

# Three workflow columns
COL_W = Inches(3.6)
COL_GAP = Inches(0.25)
COL_TOP = Inches(1.5)
COL_H = Inches(4.3)

workflows = [
    {
        'title': 'Vanilla Claude Code',
        'color': COL_VANILLA,
        'bg': LIGHT_GRAY,
        'flow': '需求描述 → AI 直接编码 → 代码产出',
        'desc': [
            {'text': 'Vanilla Claude Code', 'size': Pt(16), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(8)},
            {'text': '工作流', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '需求描述 → AI 直接编码 → 代码', 'size': Pt(12), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(12)},
            {'text': '特征', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '无前置规划步骤', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': 'AI 凭自身判断决定架构', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '无 spec / design / task 产出', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '最快交付，但无质量保障', 'size': Pt(12), 'bullet': True, 'space_after': Pt(8)},
            {'text': 'Artifact 数量：0', 'size': Pt(12), 'color': GRAY_MED, 'space_after': Pt(2)},
        ]
    },
    {
        'title': '原生 OpenSpec (4 Artifact)',
        'color': COL_NATIVE,
        'bg': LIGHT_BLUE,
        'flow': 'Proposal → Specs + Design → Tasks → 代码',
        'desc': [
            {'text': '原生 OpenSpec', 'size': Pt(16), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(8)},
            {'text': '工作流', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': 'Proposal → Specs + Design → Tasks → 代码', 'size': Pt(12), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(12)},
            {'text': '特征', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '标准 4-artifact 工作流', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': 'Specs 与 Design 并行生成', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '基本规格约束', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '质量显著优于无规格', 'size': Pt(12), 'bullet': True, 'space_after': Pt(8)},
            {'text': 'Artifact 数量：4', 'size': Pt(12), 'color': GRAY_MED, 'space_after': Pt(2)},
        ]
    },
    {
        'title': 'OpenSpec-HW2 (5 Artifact)',
        'color': COL_HW2,
        'bg': LIGHT_GREEN,
        'flow': 'Proposal → Usecases → Specs → Design → Tasks → 代码',
        'desc': [
            {'text': 'OpenSpec-HW2', 'size': Pt(16), 'color': COL_HW2, 'bold': True, 'space_after': Pt(8)},
            {'text': '工作流', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': 'Proposal → Usecases → Specs → Design → Tasks', 'size': Pt(12), 'color': COL_HW2, 'bold': True, 'space_after': Pt(12)},
            {'text': '特征', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '新增 Usecases 场景建模步骤', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '8 个用例驱动规格细化', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '级联放大效应', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '企业级质量产出', 'size': Pt(12), 'bullet': True, 'space_after': Pt(8)},
            {'text': 'Artifact 数量：5', 'size': Pt(12), 'color': GRAY_MED, 'space_after': Pt(2)},
        ]
    },
]

for i, wf in enumerate(workflows):
    col_left = Inches(0.5) + i * (COL_W + COL_GAP)
    add_colored_box(s5, col_left, COL_TOP, COL_W, COL_H, wf['bg'], wf['color'])
    add_rich_textbox(s5, Emu(col_left + Inches(0.2)), Emu(COL_TOP + Inches(0.15)),
                     Emu(COL_W - Inches(0.4)), Emu(COL_H - Inches(0.3)), wf['desc'])

# Bottom control-variable bar
add_accent_bar(s5, Inches(0.5), Inches(6.05), Inches(11.3), RGBColor(0x33, 0x33, 0x33))
add_textbox(s5, Inches(0.5), Inches(6.15), Inches(11.3), Inches(0.35),
            '控制变量：完全相同的需求描述  |  同一 AI 模型 (Claude Sonnet)  |  独立执行  |  无人工干预',
            size=Pt(11), color=GRAY_MED, align=PP_ALIGN.CENTER)
add_page_num(s5, 9, TOTAL)
print("Slide 9: Three-way experiment setup ✓")

# ── Slide 10: Core metrics comparison table ──────────────────────
s6 = prs.slides[9]
clear_slide(s6)

add_textbox(s6, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '核心指标对比：规格层级决定代码质量',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s6, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

metrics_data = [
    ['对比维度',       'Vanilla Claude Code', '原生 OpenSpec',    'OpenSpec-HW2'],
    ['代码总量',       '~1,200 行',           '~2,824 行',        '~5,015 行'],
    ['文件数',         '5-6 个',              '35 个',             '43 个'],
    ['架构层级',       '扁平结构',            '分层 / 简单类',     '6 层 Protocol 架构'],
    ['设计模式',       '无（直接实现）',       '简单类 / 函数',     'Protocol + 装饰器注册表'],
    ['测试文件',       '0 个',                '8 个',              '12 个'],
    ['测试用例',       '0',                   '67 个',             '88 个'],
    ['风控层级',       '固定止损',            '4 层完整',          '4 层 + 聚合入口'],
    ['Artifact 数',    '0',                   '4',                 '5'],
    ['Design 深度',    '0',                   '93 行 / 6 决策',    '255 行 / 8 决策'],
    ['任务分解',       '0',                   '44 个',             '95 个'],
]

create_comparison_table(
    s6, Inches(0.5), Inches(1.0), Inches(11.3), Inches(4.6),
    metrics_data,
    col_widths=[Inches(2.2), Inches(3.0), Inches(3.0), Inches(3.1)],
    font_size=Pt(11)
)

# Conclusion bar
add_accent_bar(s6, Inches(0.5), Inches(5.85), Inches(11.3), GREEN)
add_rich_textbox(s6, Inches(0.5), Inches(5.95), Inches(11.3), Inches(0.6), [
    {'text': '→ 规格化程度与代码质量强正相关：', 'size': Pt(12), 'bold': True, 'color': GREEN, 'space_after': Pt(2)},
    {'text': 'Artifact 数量 0→4→5 | 代码量 1,200→2,824→5,015 | 测试 0→67→88 | 任务分解 0→44→95',
     'size': Pt(12), 'color': RGBColor(0x55, 0x55, 0x55), 'space_after': Pt(2)},
])
add_page_num(s6, 10, TOTAL)
print("Slide 10: Core metrics table ✓")

# ── Slide 11: Architecture & risk control ────────────────────────
s7 = prs.slides[10]
clear_slide(s7)

add_textbox(s7, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '架构与风控：从扁平实现到企业级分层',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s7, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

# Left half: Architecture
HALF_W = Inches(5.5)
add_textbox(s7, Inches(0.5), Inches(0.95), HALF_W, Inches(0.35),
            '架构对比', size=Pt(15), color=BLUE, bold=True)

arch_items = [
    {
        'title': 'Vanilla', 'color': COL_VANILLA, 'bg': LIGHT_GRAY,
        'lines': [
            {'text': 'Vanilla Claude Code — 扁平原型', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '5-6 个文件，所有逻辑堆在一起', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '无接口抽象，直接调用', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '无数据验证层', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'title': '原生', 'color': COL_NATIVE, 'bg': LIGHT_BLUE,
        'lines': [
            {'text': '原生 OpenSpec — 分层架构', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '35 个文件，数据/因子/策略/风控分层', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '简单类和函数组织', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '基本模块化', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'title': 'HW2', 'color': COL_HW2, 'bg': LIGHT_GREEN,
        'lines': [
            {'text': 'OpenSpec-HW2 — 6 层 Protocol 架构', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '43 个文件，6 层清晰分离', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'Protocol 接口 + 装饰器注册', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '完整依赖注入，可独立测试', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
]
arch_box_h = Inches(1.45)
for j, item in enumerate(arch_items):
    y = Inches(1.35) + j * (arch_box_h + Inches(0.1))
    add_colored_box(s7, Inches(0.5), y, HALF_W, arch_box_h, item['bg'], item['color'])
    add_rich_textbox(s7, Inches(0.7), Emu(y + Inches(0.1)),
                     Emu(HALF_W - Inches(0.4)), Emu(arch_box_h - Inches(0.15)), item['lines'])

# Right half: Risk control
RIGHT_L = Inches(6.3)
RIGHT_W = Inches(5.5)
add_textbox(s7, RIGHT_L, Inches(0.95), RIGHT_W, Inches(0.35),
            '风控对比', size=Pt(15), color=ORANGE, bold=True)

risk_items = [
    {
        'color': COL_VANILLA, 'bg': LIGHT_GRAY,
        'lines': [
            {'text': 'Vanilla — 固定止损', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '硬编码止损阈值（如固定 5%）', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '无动态调整 / 持仓限制 / 异常检测', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '真实交易风险极高', 'size': Pt(11), 'bullet': True, 'color': RED, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_NATIVE, 'bg': LIGHT_BLUE,
        'lines': [
            {'text': '原生 — 4 层完整风控', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': 'L1 ATR动态止损 / L2 持仓限制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L3 行业集中度 / L4 异常波动检测', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '完整风控体系', 'size': Pt(11), 'bullet': True, 'color': BLUE, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_HW2, 'bg': LIGHT_GREEN,
        'lines': [
            {'text': 'HW2 — 4 层 + 聚合入口', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '上述 4 层 + RiskAggregator 统一调度', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '风控规则可配置 + 完整日志', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '生产级风控架构', 'size': Pt(11), 'bullet': True, 'color': GREEN, 'space_after': Pt(2)},
        ]
    },
]
for j, item in enumerate(risk_items):
    y = Inches(1.35) + j * (arch_box_h + Inches(0.1))
    add_colored_box(s7, RIGHT_L, y, RIGHT_W, arch_box_h, item['bg'], item['color'])
    add_rich_textbox(s7, Emu(RIGHT_L + Inches(0.2)), Emu(y + Inches(0.1)),
                     Emu(RIGHT_W - Inches(0.4)), Emu(arch_box_h - Inches(0.15)), item['lines'])

add_accent_bar(s7, Inches(0.5), Inches(6.05), Inches(11.3), BLUE)
add_textbox(s7, Inches(0.5), Inches(6.15), Inches(11.3), Inches(0.35),
            '架构分层和风控完善度与规格化程度高度正相关：规格越详细，AI 的架构决策越成熟',
            size=Pt(11), color=GRAY_MED, align=PP_ALIGN.CENTER)
add_page_num(s7, 11, TOTAL)
print("Slide 11: Architecture & risk control ✓")

# ── Slide 12: Quality dimensions ─────────────────────────────────
s8 = prs.slides[11]
clear_slide(s8)

add_textbox(s8, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '测试 × 回测 × 可维护性：三维质量对比',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s8, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

# Three quality dimension boxes
DIM_W = Inches(3.6)
DIM_GAP = Inches(0.25)
DIM_TOP = Inches(1.1)
DIM_H = Inches(4.7)

quality_dims = [
    {
        'title': '测试覆盖度',
        'color': BLUE,
        'bg': LIGHT_BLUE,
        'content': [
            {'text': '测试覆盖度', 'size': Pt(16), 'color': BLUE, 'bold': True, 'space_after': Pt(10)},
            {'text': 'Vanilla Claude Code', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '0 测试文件，0 测试用例', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': 'AI 完全跳过测试', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': '原生 OpenSpec', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '8 测试文件，67 测试用例', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '单元测试覆盖核心模块', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': 'OpenSpec-HW2', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '12 测试文件，88 测试用例', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '195 行集成测试', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '单元 + 集成 + 端到端全覆盖', 'size': Pt(12), 'bullet': True, 'space_after': Pt(4)},
        ]
    },
    {
        'title': '回测真实性',
        'color': ORANGE,
        'bg': LIGHT_ORANGE,
        'content': [
            {'text': '回测真实性', 'size': Pt(16), 'color': ORANGE, 'bold': True, 'space_after': Pt(10)},
            {'text': 'Vanilla Claude Code', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '无 A 股规则模拟', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '回测结果不可信', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': '原生 OpenSpec', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '完整 A 股规则（T+1, 涨跌停, 印花税）', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '7 个性能评估指标', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': 'OpenSpec-HW2', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '完整 A 股规则 + 防前视偏差', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '13+ 个性能评估指标', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '回测结果接近真实市场', 'size': Pt(12), 'bullet': True, 'space_after': Pt(4)},
        ]
    },
    {
        'title': '可维护性 & 配置管理',
        'color': GREEN,
        'bg': LIGHT_GREEN,
        'content': [
            {'text': '可维护性 & 配置', 'size': Pt(16), 'color': GREEN, 'bold': True, 'space_after': Pt(10)},
            {'text': 'Vanilla Claude Code', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '硬编码常量，0 配置外部化', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '修改参数需改源码', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': '原生 OpenSpec', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '基本模块化，代码分层', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '部分参数仍内嵌', 'size': Pt(12), 'bullet': True, 'space_after': Pt(10)},
            {'text': 'OpenSpec-HW2', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '外部 YAML 配置管理', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '12 类参数外部化', 'size': Pt(12), 'bullet': True, 'space_after': Pt(3)},
            {'text': '改配置不改代码', 'size': Pt(12), 'bullet': True, 'space_after': Pt(4)},
        ]
    },
]

for i, dim in enumerate(quality_dims):
    x = Inches(0.5) + i * (DIM_W + DIM_GAP)
    add_colored_box(s8, x, DIM_TOP, DIM_W, DIM_H, dim['bg'], dim['color'])
    add_rich_textbox(s8, Emu(x + Inches(0.2)), Emu(DIM_TOP + Inches(0.15)),
                     Emu(DIM_W - Inches(0.4)), Emu(DIM_H - Inches(0.3)), dim['content'])

add_accent_bar(s8, Inches(0.5), Inches(6.05), Inches(11.3), GREEN)
add_textbox(s8, Inches(0.5), Inches(6.15), Inches(11.3), Inches(0.35),
            '三个维度均呈现 Vanilla < 原生 < HW2 的阶梯式提升，Usecases 步骤的增量价值在每个维度均可量化',
            size=Pt(11), color=GRAY_MED, align=PP_ALIGN.CENTER)
add_page_num(s8, 12, TOTAL)
print("Slide 12: Quality dimensions ✓")

# ── Slide 13: Comprehensive scoring & scenarios ──────────────────
s9 = prs.slides[12]
clear_slide(s9)

add_textbox(s9, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '综合评分与适用场景',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s9, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

score_data = [
    ['评分维度 (/10)',   'Vanilla', '原生 OpenSpec', 'OpenSpec-HW2'],
    ['架构清晰度',       '6',       '7',            '9'],
    ['测试完备性',       '2',       '7',            '9'],
    ['风控能力',         '5',       '8',            '9'],
    ['回测真实性',       '6',       '7',            '8'],
    ['A 股规则合规',     '4',       '8',            '9'],
    ['可维护性',         '7',       '7',            '8'],
    ['可扩展性',         '5',       '6',            '8'],
    ['开发效率',         '9',       '6',            '7'],
    ['合计 (/80)',       '44',      '56',           '67'],
]

create_comparison_table(
    s9, Inches(0.5), Inches(1.0), Inches(6.8), Inches(4.2),
    score_data,
    col_widths=[Inches(2.0), Inches(1.5), Inches(1.7), Inches(1.6)],
    font_size=Pt(11)
)

# Scenario recommendations (right side)
add_textbox(s9, Inches(7.6), Inches(1.0), Inches(4.2), Inches(0.35),
            '适用场景推荐', size=Pt(15), color=BLACK, bold=True)
add_accent_bar(s9, Inches(7.6), Inches(1.35), Inches(4.2), RGBColor(0x33, 0x33, 0x33))

scenarios = [
    {
        'lines': [
            {'text': 'Vanilla Claude Code', 'size': Pt(13), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '44/80  适用：', 'size': Pt(11), 'bold': True, 'space_after': Pt(3)},
            {'text': '快速原型验证', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '一次性脚本 / 学习用途', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '无长期维护需求', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
        ]
    },
    {
        'lines': [
            {'text': '原生 OpenSpec', 'size': Pt(13), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '56/80  适用：', 'size': Pt(11), 'bold': True, 'space_after': Pt(3)},
            {'text': '中等复杂度项目', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '团队协作基础场景', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '标准业务系统', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
        ]
    },
    {
        'lines': [
            {'text': 'OpenSpec-HW2', 'size': Pt(13), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': '67/80  适用：', 'size': Pt(11), 'bold': True, 'space_after': Pt(3)},
            {'text': '企业级系统 / 多模块', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '合规要求高（金融/医疗）', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '长期维护 / 团队交接', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
        ]
    },
]

sc_y = Inches(1.55)
for sc in scenarios:
    add_rich_textbox(s9, Inches(7.6), sc_y, Inches(4.2), Inches(1.3), sc['lines'])
    sc_y += Inches(1.4)

# Bottom conclusion
add_accent_bar(s9, Inches(0.5), Inches(5.55), Inches(11.3), GREEN)
add_rich_textbox(s9, Inches(0.5), Inches(5.65), Inches(11.3), Inches(0.8), [
    {'text': '核心结论', 'size': Pt(13), 'color': GREEN, 'bold': True, 'space_after': Pt(3)},
    {'text': '规格化程度直接决定最终质量。Vanilla（无规格）→ 原生 OpenSpec（4 artifact）→ HW2（5 artifact + Usecases），'
     '三者总分 44→56→67，呈阶梯式提升。Usecases 步骤是 ROI 最高的前置投入。',
     'size': Pt(12), 'color': RGBColor(0x44, 0x44, 0x44), 'space_after': Pt(2)},
])
add_page_num(s9, 13, TOTAL)
print("Slide 13: Comprehensive scoring ✓")


# ═══════════════════════════════════════════════════════════════════
# SECTION B : Rebuild Appendix Slides 29-33  (0-indexed 28-32)
# ═══════════════════════════════════════════════════════════════════

# ── Slide 29: Appendix cover ─────────────────────────────────────
s25 = prs.slides[28]
clear_slide(s25)
set_dark_bg(s25)

add_textbox(s25, Inches(1.5), Inches(1.5), Inches(9), Inches(1),
            '附录 1：量化交易系统三方深度对比',
            size=Pt(36), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_accent_bar(s25, Inches(4.5), Inches(2.7), Inches(3), BLUE)

cover_lines = [
    {'text': 'Vanilla Claude Code  vs  原生 OpenSpec  vs  OpenSpec-HW2',
     'size': Pt(18), 'color': RGBColor(0xCC, 0xCC, 0xCC), 'align': PP_ALIGN.CENTER,
     'bold': True, 'space_after': Pt(16)},
    {'text': 'A股有色金属多因子量化交易系统  |  同一需求 · 三种实现 · 逐项拆解',
     'size': Pt(14), 'color': RGBColor(0x99, 0x99, 0x99), 'align': PP_ALIGN.CENTER,
     'space_after': Pt(16)},
    {'text': 'OpenSpec-HW2 版：43 源码文件 · 88 测试用例 · 6 层 Protocol 架构',
     'size': Pt(13), 'color': LIGHT_GREEN, 'align': PP_ALIGN.CENTER, 'space_after': Pt(4)},
    {'text': '原生 OpenSpec 版：35 源码文件 · 67 测试用例 · 分层架构',
     'size': Pt(13), 'color': LIGHT_BLUE, 'align': PP_ALIGN.CENTER, 'space_after': Pt(4)},
    {'text': 'Vanilla Claude Code 版：5-6 源码文件 · 0 测试 · 扁平结构',
     'size': Pt(13), 'color': LIGHT_GRAY, 'align': PP_ALIGN.CENTER, 'space_after': Pt(4)},
]
add_rich_textbox(s25, Inches(1.5), Inches(3.1), Inches(9), Inches(3.0), cover_lines)
add_page_num(s25, 29, TOTAL, color=GRAY_MED)
print("Slide 29: Appendix cover ✓")

# ── Slide 30: Three-way architecture detail ──────────────────────
s26 = prs.slides[29]
clear_slide(s26)

add_textbox(s26, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '三方架构深度对比：代码组织与设计模式',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s26, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

# Three columns for detailed architecture
A_COL_W = Inches(3.6)
A_COL_GAP = Inches(0.25)
A_COL_TOP = Inches(1.05)
A_COL_H = Inches(5.1)

arch_details = [
    {
        'color': COL_VANILLA, 'bg': LIGHT_GRAY,
        'content': [
            {'text': 'Vanilla Claude Code', 'size': Pt(15), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(6)},
            {'text': '5-6 个文件 · 扁平结构', 'size': Pt(12), 'bold': True, 'space_after': Pt(8)},
            {'text': '文件清单', 'size': Pt(12), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': 'main.py: 策略逻辑+数据获取 (~400行)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'backtest.py: 回测引擎 (~300行)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'utils.py: 工具函数 (~200行)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'config.py: 硬编码配置 (~100行)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'app.py: Flask Web (~200行)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(8)},
            {'text': '设计模式', 'size': Pt(12), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
            {'text': '无设计模式，过程式代码', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': '函数间直接调用，高耦合', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': '无接口抽象', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_NATIVE, 'bg': LIGHT_BLUE,
        'content': [
            {'text': '原生 OpenSpec', 'size': Pt(15), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(6)},
            {'text': '35 个文件 · 分层架构', 'size': Pt(12), 'bold': True, 'space_after': Pt(8)},
            {'text': '目录结构', 'size': Pt(12), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': 'src/data/: 数据获取与存储', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'src/factors/: 因子计算', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'src/strategy/: 选股策略', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'src/risk/: 风控 (4层)', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'src/backtest/: 回测引擎', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'tests/: 8 个测试文件', 'size': Pt(10), 'bullet': True, 'space_after': Pt(8)},
            {'text': '设计模式', 'size': Pt(12), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
            {'text': '简单类封装 + 函数组织', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': '基本模块化，职责初步分离', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_HW2, 'bg': LIGHT_GREEN,
        'content': [
            {'text': 'OpenSpec-HW2', 'size': Pt(15), 'color': COL_HW2, 'bold': True, 'space_after': Pt(6)},
            {'text': '43 个文件 · 6 层 Protocol 架构', 'size': Pt(12), 'bold': True, 'space_after': Pt(8)},
            {'text': '6 层架构', 'size': Pt(12), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': 'L1 数据层: 多源适配+验证器', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L2 因子层: 注册表+可扩展', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L3 选股层: 多策略组合', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L4 交易层: 订单管理+执行', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L5 风控层: 4层+聚合器', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'L6 回测层: 真实性检验', 'size': Pt(10), 'bullet': True, 'space_after': Pt(8)},
            {'text': '设计模式', 'size': Pt(12), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
            {'text': 'Protocol 接口 + 装饰器注册', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
            {'text': '依赖注入 + 外部 YAML 配置', 'size': Pt(10), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
]

for i, ad in enumerate(arch_details):
    x = Inches(0.5) + i * (A_COL_W + A_COL_GAP)
    add_colored_box(s26, x, A_COL_TOP, A_COL_W, A_COL_H, ad['bg'], ad['color'])
    add_rich_textbox(s26, Emu(x + Inches(0.2)), Emu(A_COL_TOP + Inches(0.12)),
                     Emu(A_COL_W - Inches(0.4)), Emu(A_COL_H - Inches(0.24)), ad['content'])

add_page_num(s26, 30, TOTAL)
print("Slide 30: Three-way architecture detail ✓")

# ── Slide 31: Three-way risk control detail ──────────────────────
s27 = prs.slides[30]
clear_slide(s27)

add_textbox(s27, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '风控能力三方对比：量化交易的生死线',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s27, Inches(0.5), Inches(0.75), Inches(11.3), ORANGE)

add_textbox(s27, Inches(0.5), Inches(0.85), Inches(11), Inches(0.35),
            '在量化交易中，风控能力决定了系统能否在真实市场中存活。这是三个版本差距最大的地方。',
            size=Pt(12), color=GRAY_MED)

# Three risk control columns
R_COL_TOP = Inches(1.35)
R_COL_H = Inches(4.8)

risk_details = [
    {
        'color': COL_VANILLA, 'bg': LIGHT_GRAY,
        'content': [
            {'text': 'Vanilla Claude Code', 'size': Pt(15), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(6)},
            {'text': '风控等级：初级', 'size': Pt(12), 'color': RED, 'bold': True, 'space_after': Pt(8)},
            {'text': '仅固定止损', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '硬编码止损阈值（如固定 5%）', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '不随市场波动调整', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '无持仓限制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '无行业集中度控制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '无异常波动检测', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
            {'text': '风险评估', 'size': Pt(12), 'color': RED, 'bold': True, 'space_after': Pt(3)},
            {'text': '真实交易风险极高', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '极端行情下可能爆仓', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_NATIVE, 'bg': LIGHT_BLUE,
        'content': [
            {'text': '原生 OpenSpec', 'size': Pt(15), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(6)},
            {'text': '风控等级：完整', 'size': Pt(12), 'color': BLUE, 'bold': True, 'space_after': Pt(8)},
            {'text': '4 层风控体系', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': 'L1 ATR 动态止损', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '  根据波动率自动调整止损位', 'size': Pt(10), 'color': GRAY_MED, 'space_after': Pt(3)},
            {'text': 'L2 持仓限制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '  单只上限 + 总仓位控制', 'size': Pt(10), 'color': GRAY_MED, 'space_after': Pt(3)},
            {'text': 'L3 行业集中度控制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '  避免板块过度集中', 'size': Pt(10), 'color': GRAY_MED, 'space_after': Pt(3)},
            {'text': 'L4 异常波动检测', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '  异常行情自动触发保护', 'size': Pt(10), 'color': GRAY_MED, 'space_after': Pt(2)},
        ]
    },
    {
        'color': COL_HW2, 'bg': LIGHT_GREEN,
        'content': [
            {'text': 'OpenSpec-HW2', 'size': Pt(15), 'color': COL_HW2, 'bold': True, 'space_after': Pt(6)},
            {'text': '风控等级：生产级', 'size': Pt(12), 'color': GREEN, 'bold': True, 'space_after': Pt(8)},
            {'text': '4 层 + 聚合入口', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
            {'text': '包含原生全部 4 层风控', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '+ RiskAggregator 统一调度', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '+ 风控规则 YAML 可配置', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '+ 完整风控执行日志', 'size': Pt(11), 'bullet': True, 'space_after': Pt(3)},
            {'text': '+ 风控事件回调机制', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
            {'text': '增量价值', 'size': Pt(12), 'color': GREEN, 'bold': True, 'space_after': Pt(3)},
            {'text': '统一入口让风控层可独立测试', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '配置化支持快速策略调整', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
]

for i, rd in enumerate(risk_details):
    x = Inches(0.5) + i * (A_COL_W + A_COL_GAP)
    add_colored_box(s27, x, R_COL_TOP, A_COL_W, R_COL_H, rd['bg'], rd['color'])
    add_rich_textbox(s27, Emu(x + Inches(0.2)), Emu(R_COL_TOP + Inches(0.1)),
                     Emu(A_COL_W - Inches(0.4)), Emu(R_COL_H - Inches(0.2)), rd['content'])

add_accent_bar(s27, Inches(0.5), Inches(6.35), Inches(11.3), ORANGE)
add_textbox(s27, Inches(0.5), Inches(6.45), Inches(11.3), Inches(0.3),
            '风控能力是规格化程度差异最直观的体现：无规格 = 无风控；规格越详细，风控架构越完善',
            size=Pt(11), color=GRAY_MED, align=PP_ALIGN.CENTER)
add_page_num(s27, 31, TOTAL)
print("Slide 31: Three-way risk control detail ✓")

# ── Slide 32: Three-way backtesting detail ───────────────────────
s28 = prs.slides[31]
clear_slide(s28)

add_textbox(s28, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '回测真实性三方对比：你的回测结果能信多少？',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s28, Inches(0.5), Inches(0.75), Inches(11.3), ORANGE)

add_textbox(s28, Inches(0.5), Inches(0.85), Inches(11), Inches(0.35),
            '回测是量化策略上线前的最后一道关卡。回测越贴近真实市场，上线后的意外越少。',
            size=Pt(12), color=GRAY_MED)

# A-share rules comparison table
bt_data = [
    ['A 股交易规则',      'Vanilla',   '原生 OpenSpec', 'OpenSpec-HW2'],
    ['T+1 交收制度',      '✗',         '✓',             '✓'],
    ['涨跌停限制 (10%)',  '✗',         '✓',             '✓'],
    ['印花税 / 佣金',     '✗',         '✓',             '✓'],
    ['最小买入单位 (100股)', '✗',      '✓',             '✓'],
    ['停牌处理',          '✗',         '✓',             '✓'],
    ['前视偏差检测',      '✗',         '✗',             '✓'],
    ['滑点模拟',          '✗',         '基础',          '完整'],
]

create_comparison_table(
    s28, Inches(0.5), Inches(1.35), Inches(7.0), Inches(3.2),
    bt_data,
    col_widths=[Inches(2.4), Inches(1.4), Inches(1.6), Inches(1.6)],
    font_size=Pt(11)
)

# Performance metrics comparison (right side)
add_textbox(s28, Inches(7.8), Inches(1.35), Inches(4.0), Inches(0.35),
            '性能评估指标对比', size=Pt(14), color=BLACK, bold=True)

perf_lines = [
    {'text': 'Vanilla Claude Code', 'size': Pt(12), 'color': COL_VANILLA, 'bold': True, 'space_after': Pt(3)},
    {'text': '无标准化性能指标输出', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
    {'text': '原生 OpenSpec — 7 个指标', 'size': Pt(12), 'color': COL_NATIVE, 'bold': True, 'space_after': Pt(3)},
    {'text': '年化收益率、最大回撤、夏普比率', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
    {'text': '胜率、盈亏比、波动率、Calmar比率', 'size': Pt(11), 'bullet': True, 'space_after': Pt(8)},
    {'text': 'OpenSpec-HW2 — 13+ 个指标', 'size': Pt(12), 'color': COL_HW2, 'bold': True, 'space_after': Pt(3)},
    {'text': '上述全部 + Sortino 比率', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
    {'text': '信息比率、Alpha、Beta', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
    {'text': '月度收益分布、连续亏损统计', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
    {'text': '基准对比分析', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
]
add_rich_textbox(s28, Inches(7.8), Inches(1.75), Inches(4.0), Inches(3.5), perf_lines)

# Bottom summary box
add_colored_box(s28, Inches(0.5), Inches(4.85), Inches(11.3), Inches(1.3),
                RGBColor(0xFE, 0xF3, 0xE8), ORANGE)
bt_summary = [
    {'text': '回测真实性结论', 'size': Pt(14), 'color': ORANGE, 'bold': True, 'space_after': Pt(4)},
    {'text': 'Vanilla 版完全忽略 A 股交易规则，回测结果无参考价值，直接上线必定亏损。', 'size': Pt(12), 'space_after': Pt(3)},
    {'text': '原生 OpenSpec 版覆盖了核心 A 股规则，回测基本可信；HW2 版进一步增加了前视偏差检测和完整滑点模拟，'
     '回测结果最接近真实市场表现。', 'size': Pt(12), 'space_after': Pt(2)},
]
add_rich_textbox(s28, Inches(0.7), Inches(4.95), Inches(10.9), Inches(1.1), bt_summary)

add_page_num(s28, 32, TOTAL)
print("Slide 32: Three-way backtesting detail ✓")

# ── Slide 33: Three-way comprehensive scoring ────────────────────
s29 = prs.slides[32]
clear_slide(s29)

add_textbox(s29, Inches(0.5), Inches(0.2), Inches(11), Inches(0.55),
            '综合评分：三方八维度逐项打分',
            size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s29, Inches(0.5), Inches(0.75), Inches(11.3), BLUE)

# Full scoring table with comments
full_score_data = [
    ['维度 (/10)',    'Vanilla',  '原生',  'HW2',  '关键差异'],
    ['架构清晰度',    '6',        '7',     '9',    'Protocol 架构 vs 简单类 vs 扁平'],
    ['测试完备性',    '2',        '7',     '9',    '0→67→88 测试用例'],
    ['风控能力',      '5',        '8',     '9',    '固定止损→4层→4层+聚合'],
    ['回测真实性',    '6',        '7',     '8',    '无规则→完整→+防前视偏差'],
    ['A股规则合规',   '4',        '8',     '9',    '无→完整模拟→完整+增强'],
    ['可维护性',      '7',        '7',     '8',    '简单易读→模块化→YAML配置'],
    ['可扩展性',      '5',        '6',     '8',    '改代码→基本扩展→插件式'],
    ['开发效率',      '9',        '6',     '7',    '最快交付（但质量最低）'],
    ['合计 (/80)',    '44',       '56',    '67',   'Vanilla 仅在开发效率领先'],
]

create_comparison_table(
    s29, Inches(0.5), Inches(1.0), Inches(11.3), Inches(4.2),
    full_score_data,
    col_widths=[Inches(1.8), Inches(1.2), Inches(1.2), Inches(1.2), Inches(5.9)],
    font_size=Pt(11)
)

# Conclusion box
add_colored_box(s29, Inches(0.5), Inches(5.4), Inches(11.3), Inches(1.2),
                RGBColor(0xE8, 0xF5, 0xE9), GREEN)
conclusion_lines = [
    {'text': '三方对比核心结论', 'size': Pt(14), 'color': GREEN, 'bold': True, 'space_after': Pt(4)},
    {'text': '① 规格前置 = 质量保障：三者总分 44→56→67，规格化程度越高，产出质量越高', 'size': Pt(12), 'space_after': Pt(3)},
    {'text': '② Usecases 的增量价值：原生→HW2 的 11 分差距（56→67），完全由 Usecases 步骤驱动', 'size': Pt(12), 'space_after': Pt(3)},
    {'text': '③ 不是代码多就好：Vanilla 代码最少但最快，适合原型；企业级系统需要 HW2 级别的规格化', 'size': Pt(12), 'space_after': Pt(2)},
]
add_rich_textbox(s29, Inches(0.7), Inches(5.5), Inches(10.9), Inches(1.0), conclusion_lines)

add_page_num(s29, 33, TOTAL)
print("Slide 33: Three-way comprehensive scoring ✓")


# ═══════════════════════════════════════════════════════════════════
# SECTION C : Update TOC (slide 2) and Summary (slide 28)
# ═══════════════════════════════════════════════════════════════════

# Update TOC slide (index 1) - find text about experiment
s2 = prs.slides[1]
for shape in s2.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if '初始构建 + 需求变更对比' in run.text:
                run.text = '三方对比实验（Vanilla / 原生 / HW2）+ 需求变更'
                print(f"TOC updated: {run.text}")

# Update Summary slide (index 27 = slide 28)
s28_summary = prs.slides[27]
for shape in s28_summary.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if '初始构建' in run.text and '模块化' in run.text:
                run.text = '三方对比：总分 44→56→67（Vanilla→原生→HW2） | 需求变更：边界覆盖 5x、测试 27:0、文档 636:0 行'
                print(f"Summary updated: {run.text}")


# ═══════════════════════════════════════════════════════════════════
# SECTION D : Delete extra slides
# ═══════════════════════════════════════════════════════════════════

# Original 50 slides. Delete:
#   Slides 14-15 (0-indexed 13,14) - old 原生 vs HW2 tail
#   Slides 34-40 (0-indexed 33-39) - old Appendix 1A tail + Appendix 1B
# Total: 2 + 7 = 9 slides deleted → 41 remaining
# Delete in REVERSE order to avoid index shifting

delete_indices = sorted([13, 14, 33, 34, 35, 36, 37, 38, 39], reverse=True)
slide_list = prs.slides._sldIdLst
rNs = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

for idx in delete_indices:
    if idx < len(slide_list):
        elem = slide_list[idx]
        rId = elem.get(rNs)
        slide_list.remove(elem)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        print(f"Deleted slide at old index {idx} (1-indexed: {idx+1})")

print(f"\nAfter deletion: {len(prs.slides)} slides")


# ═══════════════════════════════════════════════════════════════════
# SECTION E : Update ALL page numbers
# ═══════════════════════════════════════════════════════════════════

TOTAL_FINAL = len(prs.slides)
print(f"\nUpdating page numbers to /{TOTAL_FINAL}...")

for idx, slide in enumerate(prs.slides):
    page = idx + 1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Page number shapes are typically at bottom-right
        if shape.top and shape.top > 5800000 and shape.left and shape.left > 9000000:
            txt = shape.text_frame.text.strip()
            if re.match(r'^\d+/\d+$', txt):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = f"{page}/{TOTAL_FINAL}"


# ═══════════════════════════════════════════════════════════════════
# SECTION F : Save
# ═══════════════════════════════════════════════════════════════════

prs.save(PPTX)
print(f"\n{'='*60}")
print(f"DONE! Saved {PPTX}")
print(f"Total slides: {TOTAL_FINAL}")
print(f"{'='*60}")
