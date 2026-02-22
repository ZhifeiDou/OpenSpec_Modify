"""
Rebuild Appendix 2 slides (30-39) with complete, readable content.
Style: between presentation and reading material.
"""
import copy, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml.etree import SubElement
from lxml import etree
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

PPTX = r"C:\Users\Zhifei Dou\Desktop\Openspec_modified\OpenSpec_技术研究汇报.pptx"
prs = Presentation(PPTX)

SW = prs.slide_width   # 12192000
SH = prs.slide_height  # 6858000

# --- Colors ---
DARK_BG   = RGBColor(0x16, 0x21, 0x3E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LT   = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MED  = RGBColor(0x99, 0x99, 0x99)
BLACK     = RGBColor(0x33, 0x33, 0x33)
BLUE      = RGBColor(0x00, 0x78, 0xD4)
GREEN     = RGBColor(0x10, 0x7C, 0x10)
ORANGE    = RGBColor(0xE8, 0x6C, 0x00)
PURPLE    = RGBColor(0x6B, 0x69, 0xD6)
TEAL      = RGBColor(0x00, 0x8B, 0x8B)
RED       = RGBColor(0xD1, 0x34, 0x38)

# --- Delete existing slides 30-38 (0-indexed: 29-37) ---
slide_list = prs.slides._sldIdLst
rNs = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
# Collect elements to remove (reverse order to avoid index shift)
elements_to_remove = []
for i in range(min(38, len(prs.slides)) - 1, 28, -1):
    elements_to_remove.append(slide_list[i])

for elem in elements_to_remove:
    rId = elem.get(rNs)
    slide_list.remove(elem)
    # Remove the relationship
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass

print(f"After deletion: {len(prs.slides)} slides")

# --- Helper functions ---
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_shape(slide, left, top, width, height):
    """Add a rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

def add_textbox(slide, left, top, width, height, text, size=Pt(14), color=BLACK,
                bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a simple text box with one run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_rich_textbox(slide, left, top, width, height, paragraphs_data):
    """
    Add a text box with multiple paragraphs.
    paragraphs_data: list of dicts with keys:
      - text: str
      - size: Pt value (default 14)
      - color: RGBColor (default BLACK)
      - bold: bool (default False)
      - align: PP_ALIGN (default LEFT)
      - space_after: Pt value (default 6)
      - bullet: bool (default False) - add bullet prefix
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = pd.get('align', PP_ALIGN.LEFT)
        p.space_after = pd.get('space_after', Pt(6))

        txt = pd.get('text', '')
        if pd.get('bullet'):
            txt = '• ' + txt

        run = p.add_run()
        run.text = txt
        run.font.size = pd.get('size', Pt(14))
        run.font.color.rgb = pd.get('color', BLACK)
        run.font.bold = pd.get('bold', False)
        run.font.name = 'Microsoft YaHei'

    return txBox

def add_accent_bar(slide, left, top, width, color):
    """Add a thin colored accent bar."""
    bar = add_shape(slide, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_colored_box(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle box with fill."""
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(1)
    else:
        box.line.fill.background()
    return box

def add_page_num(slide, page, total, color=GRAY_MED):
    """Add page number at bottom right."""
    add_textbox(slide, Emu(SW - Inches(1.2)), Emu(SH - Inches(0.5)),
                Inches(1), Inches(0.4), f"{page}/{total}",
                size=Pt(11), color=color, align=PP_ALIGN.RIGHT)

def make_direction_slide(title, accent_color, priority_tag,
                         narrative_paras, evidence_paras, takeaway,
                         page, total):
    """
    Create a direction detail slide with standard layout:
    - Title bar with accent color
    - Left 58%: narrative (complete sentences)
    - Right 38%: evidence box (学术依据 + 商业价值)
    - Bottom: key takeaway
    """
    slide = add_blank_slide()

    # Title
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6),
                title, size=Pt(24), color=BLACK, bold=True)

    # Accent bar under title
    add_accent_bar(slide, Inches(0.6), Inches(0.85), Inches(9.5), accent_color)

    # Priority tag (top right)
    add_textbox(slide, Inches(9.2), Inches(0.3), Inches(2.5), Inches(0.4),
                priority_tag, size=Pt(12), color=accent_color, bold=True,
                align=PP_ALIGN.RIGHT)

    # Left narrative (58% width)
    LEFT_W = Inches(6.2)
    add_rich_textbox(slide, Inches(0.6), Inches(1.1), LEFT_W, Inches(4.8), narrative_paras)

    # Right evidence box (38% width)
    RIGHT_L = Inches(7.1)
    RIGHT_W = Inches(4.6)

    # Evidence background box
    box = add_colored_box(slide, RIGHT_L, Inches(1.1), RIGHT_W, Inches(4.8), GRAY_LT)

    # Evidence content
    add_rich_textbox(slide, Emu(RIGHT_L + Inches(0.25)), Inches(1.2),
                     Emu(RIGHT_W - Inches(0.5)), Inches(4.6), evidence_paras)

    # Bottom takeaway bar
    add_accent_bar(slide, Inches(0.6), Inches(6.2), Inches(11), accent_color)
    add_textbox(slide, Inches(0.6), Inches(6.3), Inches(11), Inches(0.4),
                takeaway, size=Pt(12), color=GRAY_MED)

    add_page_num(slide, page, total)
    return slide


# ========== BUILD NEW SLIDES ==========
TOTAL = 39  # 29 original + 10 new

# --- Slide 30: Cover ---
s30 = add_blank_slide()
set_dark_bg(s30)
add_textbox(s30, Inches(1.5), Inches(1.8), Inches(9), Inches(1.2),
            '附录 2：六大研究改进方向详解', size=Pt(36), color=WHITE, bold=True,
            align=PP_ALIGN.CENTER)
add_accent_bar(s30, Inches(4.5), Inches(3.2), Inches(3), BLUE)

cover_lines = [
    {'text': '以下内容基于 2023-2026 年软件工程与需求工程领域的学术研究，', 'size': Pt(15), 'color': RGBColor(0xCC,0xCC,0xCC), 'align': PP_ALIGN.CENTER, 'space_after': Pt(2)},
    {'text': '结合 17 次定向文献检索，梳理了 OpenSpec 最具潜力的六大改进方向。', 'size': Pt(15), 'color': RGBColor(0xCC,0xCC,0xCC), 'align': PP_ALIGN.CENTER, 'space_after': Pt(20)},
    {'text': '每个方向包含：现状问题 → 解决方案 → 学术依据 → 商业价值', 'size': Pt(14), 'color': RGBColor(0x99,0x99,0x99), 'align': PP_ALIGN.CENTER, 'space_after': Pt(8)},
    {'text': '其中 Domain-Specific Packs 作为公司核心竞争力方向，进行了重点展开。', 'size': Pt(14), 'color': RGBColor(0xFF, 0xD7, 0x00), 'align': PP_ALIGN.CENTER, 'bold': True},
]
add_rich_textbox(s30, Inches(1.5), Inches(3.6), Inches(9), Inches(2.5), cover_lines)
add_page_num(s30, 30, TOTAL, color=GRAY_MED)
print("Slide 30: Cover done")


# --- Slide 31: Spec Quality Gate ---
make_direction_slide(
    title='改进方向 1：Spec Quality Gate（智能需求质量关卡）',
    accent_color=BLUE,
    priority_tag='优先级：最高 | 短期可落地',
    narrative_paras=[
        {'text': '现状问题', 'size': Pt(16), 'color': BLUE, 'bold': True, 'space_after': Pt(4)},
        {'text': '在大多数团队中，需求规格的撰写依赖个人经验，缺乏统一的质量标准。'
                 'PM 写完需求后直接交给开发，但"写了什么"和"开发理解了什么"之间往往存在巨大偏差。'
                 '学术研究表明，30-50% 的软件返工最终可以追溯到需求阶段的模糊、遗漏或矛盾。',
         'size': Pt(13), 'space_after': Pt(14)},
        {'text': '解决方案', 'size': Pt(16), 'color': BLUE, 'bold': True, 'space_after': Pt(4)},
        {'text': '在 OpenSpec 的工作流中增加一个 AI 驱动的"质量关卡"节点。'
                 '每条需求在进入设计和开发阶段之前，自动接受三项检查：可测试性（能否写出对应的测试用例）、'
                 '无歧义性（是否存在多种理解方式）、完整性（是否覆盖了所有必要场景）。',
         'size': Pt(13), 'space_after': Pt(8)},
        {'text': '系统会为每条需求生成一个 0-100 分的 Quality Score，'
                 '未达标的条目会附带具体的改进建议，让作者在源头修正问题。',
         'size': Pt(13), 'space_after': Pt(8)},
    ],
    evidence_paras=[
        {'text': '学术依据', 'size': Pt(14), 'color': BLUE, 'bold': True, 'space_after': Pt(4)},
        {'text': 'Ferrari et al. (2023) 通过 74 篇系统性综述证实，LLM 在需求缺陷检测上已达到实用水平。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'Hymel et al. (2025) 实验显示，GPT-4 检测需求歧义的 F1 值达 0.72，覆盖率超过人工审查。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'ISO 29148 提供了国际通用的需求工程质量框架，可作为评分标准的参照系。', 'size': Pt(12), 'space_after': Pt(16), 'bullet': True},
        {'text': '商业价值', 'size': Pt(14), 'color': BLUE, 'bold': True, 'space_after': Pt(4)},
        {'text': '减少 30-50% 的需求返工成本，这在大型项目中意味着数周的开发时间节省。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '可作为 ISO/CMMI 合规工具，面向有审计需求的企业客户销售。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '作为 OpenSpec 的"第一个可交付功能"，是最快见效的改进点。', 'size': Pt(12), 'bullet': True},
    ],
    takeaway='实施路径：proposal → 校验检查点 → specs 分析 → 评分报告 → design 审查 → Quality Score 仪表盘',
    page=31, total=TOTAL
)
print("Slide 31: Spec Quality Gate done")


# --- Slide 32: Spec-Code Traceability ---
make_direction_slide(
    title='改进方向 2：Spec-Code Traceability（需求-代码自动追溯）',
    accent_color=GREEN,
    priority_tag='优先级：高 | 中期建设',
    narrative_paras=[
        {'text': '现状问题', 'size': Pt(16), 'color': GREEN, 'bold': True, 'space_after': Pt(4)},
        {'text': '随着项目迭代，需求文档和实际代码之间会逐渐"漂移"（Spec Drift）——'
                 '需求改了但代码没跟上，或者代码改了但需求文档还是旧版本。'
                 '在大型团队中，这种漂移往往要到集成测试甚至上线后才被发现。',
         'size': Pt(13), 'space_after': Pt(8)},
        {'text': '更关键的是，EU AI Act 将于 2026 年 8 月正式实施，'
                 '要求所有 AI 系统必须具备完整的需求-代码可追溯性（Article 11）。'
                 '这意味着 traceability 不再是"锦上添花"，而是合规刚需。',
         'size': Pt(13), 'space_after': Pt(14)},
        {'text': '解决方案', 'size': Pt(16), 'color': GREEN, 'bold': True, 'space_after': Pt(4)},
        {'text': '建立从 spec ID 到代码位置的双向映射关系。当需求变更时，系统自动标记受影响的代码文件；'
                 '当代码被修改时，反向检查是否有需求需要同步更新。'
                 '整套机制集成到 CI/CD 流水线中，实现"Living Specs"——需求文档始终反映代码的真实状态。',
         'size': Pt(13), 'space_after': Pt(8)},
    ],
    evidence_paras=[
        {'text': '学术依据', 'size': Pt(14), 'color': GREEN, 'bold': True, 'space_after': Pt(4)},
        {'text': 'Cheng et al. (2026) 提出基于 LLMARE 的需求追踪框架，F1 值达 0.88，显著优于传统方法。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'EU AI Act Article 11 明确要求高风险 AI 系统具备全链路可追溯性。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '多家监管机构（FDA、金融监管）正在跟进类似要求。', 'size': Pt(12), 'space_after': Pt(16), 'bullet': True},
        {'text': '商业价值', 'size': Pt(14), 'color': GREEN, 'bold': True, 'space_after': Pt(4)},
        {'text': '合规驱动型市场：面向金融、医疗、AI 等受监管行业，traceability 是准入门槛。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'SaaS 订阅模式：按项目规模收费的持续性收入。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '与 Quality Gate 形成产品组合：先检查质量，再追踪落地，闭环管理。', 'size': Pt(12), 'bullet': True},
    ],
    takeaway='实施路径：spec ↔ code mapping → git diff drift detection → CI/CD 集成 → Living Specs 仪表盘',
    page=32, total=TOTAL
)
print("Slide 32: Traceability done")


# --- Slide 33: Multi-Agent Orchestration ---
make_direction_slide(
    title='改进方向 3：Multi-Agent Orchestration（多智能体协同编排）',
    accent_color=ORANGE,
    priority_tag='优先级：高 | 中期建设',
    narrative_paras=[
        {'text': '现状问题', 'size': Pt(16), 'color': ORANGE, 'bold': True, 'space_after': Pt(4)},
        {'text': '当前的 AI 编程助手通常是"单 Agent"模式——一个 AI 从头到尾处理整个任务。'
                 '这在简单任务上没问题，但面对包含 20+ 子任务的复杂项目时，'
                 '单个 Agent 容易丢失上下文、遗忘前序步骤，错误率随任务规模急剧上升。',
         'size': Pt(13), 'space_after': Pt(14)},
        {'text': '解决方案', 'size': Pt(16), 'color': ORANGE, 'bold': True, 'space_after': Pt(4)},
        {'text': '利用 OpenSpec 天然的 task DAG（任务有向无环图）结构，将复杂任务拆解为多个子任务，'
                 '分配给不同的专业化 Agent 并行处理。每个 Agent 只需关注自己的子任务，'
                 '而 OpenSpec 的 artifact（需求规格、设计文档）充当 Agent 之间的"协议"，确保最终产出一致。',
         'size': Pt(13), 'space_after': Pt(8)},
        {'text': '这类似于软件团队的分工协作，只不过"团队成员"变成了多个 AI Agent。',
         'size': Pt(13), 'space_after': Pt(8)},
    ],
    evidence_paras=[
        {'text': '学术依据', 'size': Pt(14), 'color': ORANGE, 'bold': True, 'space_after': Pt(4)},
        {'text': 'GitHub 于 2026 年 2 月发布 Agent HQ，支持多 Agent 协同开发，验证了这一方向的产业价值。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'Anthropic 的 Claude Code Agent Teams、SWE-bench 多 Agent 方案均显示并行处理效率提升 3-4 倍。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'DAG 调度是成熟的工程范式（Airflow、Prefect 等已在数据工程领域广泛验证）。', 'size': Pt(12), 'space_after': Pt(16), 'bullet': True},
        {'text': '商业价值', 'size': Pt(14), 'color': ORANGE, 'bold': True, 'space_after': Pt(4)},
        {'text': '效率提升 3-4 倍：复杂任务的端到端时间大幅缩短。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '计算资源按量付费：多 Agent 并行 = 更高的 API 调用量 = 收入增长。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '与 GitHub/Anthropic 的生态对齐，可借势而非逆流。', 'size': Pt(12), 'bullet': True},
    ],
    takeaway='实施路径：tasks.md → DAG 解析 → Agent Worker Pool → 乐观并发 → artifact 冲突检测 → 合并',
    page=33, total=TOTAL
)
print("Slide 33: Multi-Agent done")


# --- Slide 34: Platform / Web Hub ---
make_direction_slide(
    title='改进方向 4：Platform / Web Hub（从命令行到协作平台）',
    accent_color=PURPLE,
    priority_tag='优先级：中 | 中长期',
    narrative_paras=[
        {'text': '现状问题', 'size': Pt(16), 'color': PURPLE, 'bold': True, 'space_after': Pt(4)},
        {'text': 'OpenSpec 目前是一个命令行工具（CLI），这意味着只有会用终端的工程师才能使用它。'
                 '但在实际项目中，需求的利益相关者还包括产品经理、测试工程师、业务方——'
                 '他们需要查看和评审需求文档，但不会使用命令行。',
         'size': Pt(13), 'space_after': Pt(8)},
        {'text': '此外，artifact graph（需求之间的依赖关系图）在 CLI 中难以直观呈现，'
                 '而这恰恰是理解项目全貌的关键信息。',
         'size': Pt(13), 'space_after': Pt(14)},
        {'text': '解决方案', 'size': Pt(16), 'color': PURPLE, 'bold': True, 'space_after': Pt(4)},
        {'text': '分阶段将 OpenSpec 从 CLI 扩展为 Web 平台：'
                 '第一步是提供 Web UI 来可视化 artifact graph 和 spec 文档；'
                 '第二步增加协作功能（评审、评论、审批流）；'
                 '第三步构建 Marketplace，让社区可以分享和交易 spec 模板和 Domain Packs。',
         'size': Pt(13), 'space_after': Pt(8)},
    ],
    evidence_paras=[
        {'text': '学术依据', 'size': Pt(14), 'color': PURPLE, 'bold': True, 'space_after': Pt(4)},
        {'text': 'AI 编程工具市场 2025 年规模达 $202 亿，Cursor ($2.5B+)、Devin ($2B+) 均走产品化路线。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'Kiro (Amazon) 和 GitHub 已推出 spec review 的可视化功能，证明市场需求存在。', 'size': Pt(12), 'space_after': Pt(16), 'bullet': True},
        {'text': '商业价值', 'size': Pt(14), 'color': PURPLE, 'bold': True, 'space_after': Pt(4)},
        {'text': 'PLG（产品驱动增长）：免费 CLI → 付费 Web 协作 → 企业版，经典的 SaaS 增长路径。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'Marketplace 抽成：第三方 spec 模板和 Domain Pack 的交易平台。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '扩大用户群：从"工程师专用"变为"全团队可用"，市场天花板大幅提高。', 'size': Pt(12), 'bullet': True},
    ],
    takeaway='实施路径：CLI 引擎（已有）→ REST API 层 → Web 前端 → 团队协作功能 → Marketplace',
    page=34, total=TOTAL
)
print("Slide 34: Platform done")


# --- Slide 35: Domain-Specific Packs (1/3) - 核心概念 ---
s35 = add_blank_slide()
add_textbox(s35, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6),
            '改进方向 5：Domain-Specific Packs（领域规格智能包）', size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s35, Inches(0.6), Inches(0.85), Inches(11), TEAL)
add_textbox(s35, Inches(9.2), Inches(0.3), Inches(2.5), Inches(0.4),
            '核心竞争力方向 | 1/3', size=Pt(12), color=TEAL, bold=True, align=PP_ALIGN.RIGHT)

# Core insight - highlighted box
box35 = add_colored_box(s35, Inches(0.6), Inches(1.15), Inches(11), Inches(1.5),
                        RGBColor(0xE6, 0xF7, 0xF7), TEAL)
insight_paras = [
    {'text': '核心洞察', 'size': Pt(16), 'color': TEAL, 'bold': True, 'space_after': Pt(6)},
    {'text': '通用大模型（LLM）擅长写代码，但不懂行业。它不知道金融交易必须做风控检查、'
             '医疗系统必须符合 HIPAA、汽车软件必须通过 ISO 26262 认证。'
             '而公司多年积累的行业 Know-How——业务规则、合规标准、最佳实践——恰好填补了这个空白。'
             '将这些 domain knowledge 封装成结构化的"领域规格包"（Domain Pack），'
             '就是把公司的隐性知识变成可复用、可销售的数字资产。',
     'size': Pt(13), 'space_after': Pt(4)},
]
add_rich_textbox(s35, Inches(0.9), Inches(1.25), Inches(10.4), Inches(1.4), insight_paras)

# What is a Domain Pack?
what_paras = [
    {'text': '什么是 Domain Pack？', 'size': Pt(16), 'color': TEAL, 'bold': True, 'space_after': Pt(6)},
    {'text': '一个 Domain Pack 本质上是一组行业专属的"规格模板 + 规则集 + 验证逻辑"，以 OpenSpec 的标准格式封装。'
             '当 AI 在生成需求规格和代码时，Domain Pack 会作为先验知识（prior knowledge）注入工作流，'
             '使 AI 的产出自动遵守该行业的规范和约束条件。',
     'size': Pt(13), 'space_after': Pt(10)},
    {'text': '一个 Pack 包含的具体内容：', 'size': Pt(13), 'bold': True, 'space_after': Pt(4)},
    {'text': 'Rules（行业规则）：如"所有交易必须经过风控引擎"、"患者数据必须加密存储"', 'size': Pt(12), 'space_after': Pt(3), 'bullet': True},
    {'text': 'Spec 模板：预定义的需求规格模板，包含行业必填字段和标准结构', 'size': Pt(12), 'space_after': Pt(3), 'bullet': True},
    {'text': 'Schema 定义：数据结构和接口的行业标准约束', 'size': Pt(12), 'space_after': Pt(3), 'bullet': True},
    {'text': '验证规则（Validators）：自动检查生成的代码是否违反行业规范', 'size': Pt(12), 'space_after': Pt(3), 'bullet': True},
]
add_rich_textbox(s35, Inches(0.6), Inches(2.85), Inches(11), Inches(3.2), what_paras)

add_accent_bar(s35, Inches(0.6), Inches(6.2), Inches(11), TEAL)
add_textbox(s35, Inches(0.6), Inches(6.3), Inches(11), Inches(0.4),
            '类比：如果 OpenSpec 是一个"AI 建筑师"，那么 Domain Pack 就是"行业建筑规范手册"——没有它，AI 盖的房子不合规。',
            size=Pt(12), color=GRAY_MED)
add_page_num(s35, 35, TOTAL)
print("Slide 35: Domain Packs 1/3 done")


# --- Slide 36: Domain-Specific Packs (2/3) - 为什么是护城河 ---
s36 = add_blank_slide()
add_textbox(s36, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6),
            'Domain-Specific Packs：为什么是核心竞争壁垒', size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s36, Inches(0.6), Inches(0.85), Inches(11), TEAL)
add_textbox(s36, Inches(9.2), Inches(0.3), Inches(2.5), Inches(0.4),
            '核心竞争力方向 | 2/3', size=Pt(12), color=TEAL, bold=True, align=PP_ALIGN.RIGHT)

moat_paras = [
    {'text': '护城河分析：为什么别人难以复制', 'size': Pt(16), 'color': TEAL, 'bold': True, 'space_after': Pt(8)},
    {'text': '1. Domain Knowledge 是稀缺资源', 'size': Pt(14), 'bold': True, 'space_after': Pt(3)},
    {'text': '行业规则、合规标准、最佳实践，这些知识散布在行业专家的脑中、监管文件里、和多年项目经验中。'
             '它们无法通过简单的网页爬取或模型训练获得。公司已有的行业积累，就是天然的数据壁垒。',
     'size': Pt(13), 'space_after': Pt(10)},
    {'text': '2. 先发优势带来网络效应', 'size': Pt(14), 'bold': True, 'space_after': Pt(3)},
    {'text': '第一个做出高质量金融 Pack 的团队，会吸引金融领域的用户；用户的使用反馈又会帮助改进 Pack 质量，'
             '形成"越用越好→越好越多人用"的飞轮。后来者即使有技术能力，也缺乏这个冷启动的数据积累。',
     'size': Pt(13), 'space_after': Pt(10)},
    {'text': '3. 合规认证构成额外壁垒', 'size': Pt(14), 'bold': True, 'space_after': Pt(3)},
    {'text': '如果 Domain Pack 能通过 ISO 26262（汽车）、HIPAA（医疗）、SOX（金融）等行业认证，'
             '这本身就是一个极高的准入门槛。竞争对手不仅要做出 Pack，还要花时间和成本通过认证。',
     'size': Pt(13), 'space_after': Pt(10)},
    {'text': '4. 与 OpenSpec 生态深度绑定', 'size': Pt(14), 'bold': True, 'space_after': Pt(3)},
    {'text': 'Domain Pack 是基于 OpenSpec 的 spec 格式和 schema 体系构建的，'
             '这意味着使用 Pack 的客户同时也被锁定在 OpenSpec 生态中。Pack 越多，生态越强，切换成本越高。',
     'size': Pt(13), 'space_after': Pt(4)},
]
add_rich_textbox(s36, Inches(0.6), Inches(1.1), Inches(11), Inches(5.0), moat_paras)

add_accent_bar(s36, Inches(0.6), Inches(6.2), Inches(11), TEAL)
add_textbox(s36, Inches(0.6), Inches(6.3), Inches(11), Inches(0.4),
            '一句话总结：AI 能力人人都有，但 domain knowledge + 结构化封装 + 合规认证 = 难以逾越的竞争壁垒。',
            size=Pt(12), color=GRAY_MED)
add_page_num(s36, 36, TOTAL)
print("Slide 36: Domain Packs 2/3 done")


# --- Slide 37: Domain-Specific Packs (3/3) - 落地路线图 ---
s37 = add_blank_slide()
add_textbox(s37, Inches(0.6), Inches(0.3), Inches(9), Inches(0.6),
            'Domain-Specific Packs：行业示例与落地路线图', size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s37, Inches(0.6), Inches(0.85), Inches(11), TEAL)
add_textbox(s37, Inches(9.2), Inches(0.3), Inches(2.5), Inches(0.4),
            '核心竞争力方向 | 3/3', size=Pt(12), color=TEAL, bold=True, align=PP_ALIGN.RIGHT)

# Three columns for three industries
COL_W = Inches(3.4)
COL_GAP = Inches(0.3)
COL_Y = Inches(1.2)
COL_H = Inches(3.4)

industries = [
    {
        'name': '金融 / 量化交易',
        'color': RGBColor(0x00, 0x78, 0xD4),
        'bg': RGBColor(0xE8, 0xF0, 0xFE),
        'timeline': 'Q1 启动（已有量化交易基础）',
        'content': [
            {'text': '金融 / 量化交易 Pack', 'size': Pt(15), 'color': RGBColor(0x00, 0x78, 0xD4), 'bold': True, 'space_after': Pt(6)},
            {'text': '为什么先做金融？', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': '公司在量化交易项目中已积累了完整的风控规则、回测验证逻辑和交易引擎架构经验。'
                     '这些就是现成的 domain knowledge，只需结构化封装即可变成 Pack。',
             'size': Pt(11), 'space_after': Pt(8)},
            {'text': 'Pack 包含的内容：', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': '风控规则集（止损、持仓限制、异常检测）', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '回测真实性验证逻辑（防前视偏差）', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': '交易引擎 spec 模板', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'SOX / MiFID II 合规检查', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'name': '医疗 / 生命科学',
        'color': RGBColor(0x10, 0x7C, 0x10),
        'bg': RGBColor(0xE8, 0xF5, 0xE9),
        'timeline': 'Q2 启动（合规驱动）',
        'content': [
            {'text': '医疗 / 生命科学 Pack', 'size': Pt(15), 'color': RGBColor(0x10, 0x7C, 0x10), 'bold': True, 'space_after': Pt(6)},
            {'text': '市场驱动力：', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': 'FDA 对医疗软件的监管日趋严格，HIPAA 合规是所有医疗 IT 系统的准入门槛。'
                     '医疗机构愿意为确定性的合规方案支付高溢价。',
             'size': Pt(11), 'space_after': Pt(8)},
            {'text': 'Pack 包含的内容：', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': 'HIPAA 数据保护规则集', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'FDA 510(k) 软件验证流程模板', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'HL7 FHIR 接口 schema', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'EU MDR 合规检查', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
    {
        'name': '汽车 / 工业控制',
        'color': RGBColor(0xE8, 0x6C, 0x00),
        'bg': RGBColor(0xFE, 0xF3, 0xE8),
        'timeline': 'Q3 启动（安全关键）',
        'content': [
            {'text': '汽车 / 工业控制 Pack', 'size': Pt(15), 'color': RGBColor(0xE8, 0x6C, 0x00), 'bold': True, 'space_after': Pt(6)},
            {'text': '市场驱动力：', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': '汽车软件正从嵌入式向智能化转型，ISO 26262 功能安全认证是强制要求。'
                     '传统车企缺乏 AI 辅助的 spec 工具，这是一个蓝海市场。',
             'size': Pt(11), 'space_after': Pt(8)},
            {'text': 'Pack 包含的内容：', 'size': Pt(12), 'bold': True, 'space_after': Pt(3)},
            {'text': 'ASIL 等级 spec 模板', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'AUTOSAR 架构约束', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'OTA 更新安全验证', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
            {'text': 'IEC 61508 工业安全标准', 'size': Pt(11), 'bullet': True, 'space_after': Pt(2)},
        ]
    },
]

for i, ind in enumerate(industries):
    col_left = Inches(0.6) + i * (COL_W + COL_GAP)
    # Column background box
    add_colored_box(s37, col_left, COL_Y, COL_W, COL_H, ind['bg'], ind['color'])
    # Column content
    add_rich_textbox(s37, Emu(col_left + Inches(0.2)), Emu(COL_Y + Inches(0.15)),
                     Emu(COL_W - Inches(0.4)), Emu(COL_H - Inches(0.3)), ind['content'])

# Business model summary
biz_paras = [
    {'text': '商业模式：Vertical SaaS', 'size': Pt(14), 'color': TEAL, 'bold': True, 'space_after': Pt(4)},
    {'text': '每个 Domain Pack 按行业 + 规模定价（年订阅制），目标客单价 $10K-50K/年。'
             '先从公司最熟悉的金融领域切入（冷启动成本最低），再向医疗和汽车扩展。'
             '每进入一个新行业，都需要与该领域的行业专家合作，这进一步巩固了护城河。',
     'size': Pt(12), 'space_after': Pt(4)},
]
add_rich_textbox(s37, Inches(0.6), Inches(4.8), Inches(11), Inches(1.3), biz_paras)

add_accent_bar(s37, Inches(0.6), Inches(6.2), Inches(11), TEAL)
add_textbox(s37, Inches(0.6), Inches(6.3), Inches(11), Inches(0.4),
            '落地路线：金融 Pack（Q1，已有基础）→ 医疗 Pack（Q2，合规驱动）→ 汽车 Pack（Q3，安全关键）→ 更多行业',
            size=Pt(12), color=GRAY_MED)
add_page_num(s37, 37, TOTAL)
print("Slide 37: Domain Packs 3/3 done")


# --- Slide 38: Formal Verification ---
make_direction_slide(
    title='改进方向 6：Formal Verification（规格的形式化验证）',
    accent_color=RED,
    priority_tag='优先级：低（长期） | 技术储备',
    narrative_paras=[
        {'text': '现状问题', 'size': Pt(16), 'color': RED, 'bold': True, 'space_after': Pt(4)},
        {'text': '目前验证"代码是否满足需求"主要依赖测试——但测试只能覆盖有限的场景。'
                 '对于关键代码路径，"测试没报错"不等于"一定没有 bug"。'
                 '在航空、金融、医疗等领域，这种不确定性是不可接受的。',
         'size': Pt(13), 'space_after': Pt(14)},
        {'text': '解决方案', 'size': Pt(16), 'color': RED, 'bold': True, 'space_after': Pt(4)},
        {'text': '将 OpenSpec 的 Given/When/Then 格式的需求规格，'
                 '自动转换为形式化语言（如 LTL、CTL、SMT 公式），'
                 '然后用数学方法证明代码在所有可能的输入下都满足规格要求——而不是仅在测试用例覆盖的输入下。',
         'size': Pt(13), 'space_after': Pt(8)},
        {'text': '这是一个技术门槛很高的长期方向，但一旦实现，'
                 '将是 OpenSpec 最具技术壁垒的差异化功能。',
         'size': Pt(13), 'space_after': Pt(8)},
    ],
    evidence_paras=[
        {'text': '学术依据', 'size': Pt(14), 'color': RED, 'bold': True, 'space_after': Pt(4)},
        {'text': 'Klopmann (2025) 提出 AI 辅助形式化验证框架，将验证通过率从 68% 提升至 96%。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'Vericode benchmark 显示 LLM 生成的形式化 spec 质量已接近人工水平。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'ReaLis、Zhipu 等团队正在该方向快速推进。', 'size': Pt(12), 'space_after': Pt(16), 'bullet': True},
        {'text': '商业价值', 'size': Pt(14), 'color': RED, 'bold': True, 'space_after': Pt(4)},
        {'text': '安全关键行业（航空、医疗、金融）的形式化验证咨询费用通常在 $10K-$100K+ 量级。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': 'FDA/ISO 认证越来越认可形式化方法，长期具有合规价值。', 'size': Pt(12), 'space_after': Pt(4), 'bullet': True},
        {'text': '技术壁垒极高，一旦建立先发优势，竞争对手追赶成本巨大。', 'size': Pt(12), 'bullet': True},
    ],
    takeaway='实施路径：Given/When/Then → 属性提取 → LTL/CTL 公式生成 → SMT 求解器验证 → 验证报告',
    page=38, total=TOTAL
)
print("Slide 38: Formal Verification done")


# --- Slide 39: Priority Overview ---
s39 = add_blank_slide()
add_textbox(s39, Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.6),
            '优先级排序与推荐实施路径', size=Pt(24), color=BLACK, bold=True)
add_accent_bar(s39, Inches(0.6), Inches(0.85), Inches(11), RGBColor(0x33,0x33,0x33))

# Simplified priority list
priority_data = [
    ('1', 'Spec Quality Gate', '短期（1-3月）', BLUE, '最快见效，直接提升需求质量，降低返工成本'),
    ('2', 'Domain-Specific Packs', '短期启动，持续建设', TEAL, '核心竞争壁垒，公司 domain knowledge 的商业化载体'),
    ('3', 'Multi-Agent Orchestration', '中期（3-6月）', ORANGE, '效率倍增器，与行业趋势对齐'),
    ('4', 'Spec-Code Traceability', '中期（3-6月）', GREEN, '合规刚需，EU AI Act 驱动'),
    ('5', 'Platform / Web Hub', '中长期（6-12月）', PURPLE, 'PLG 增长引擎，扩大用户群'),
    ('6', 'Formal Verification', '长期（12月+）', RED, '技术储备，极高壁垒但实现周期长'),
]

ROW_Y_START = Inches(1.15)
ROW_H = Inches(0.75)
ROW_GAP = Inches(0.08)

for i, (num, name, timeline, color, desc) in enumerate(priority_data):
    row_y = ROW_Y_START + i * (ROW_H + ROW_GAP)

    # Number circle (colored box)
    num_box = add_colored_box(s39, Inches(0.6), row_y, Inches(0.5), ROW_H, color)
    # number text
    add_textbox(s39, Inches(0.6), Emu(row_y + Inches(0.15)), Inches(0.5), Inches(0.4),
                num, size=Pt(20), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Name
    add_textbox(s39, Inches(1.2), Emu(row_y + Inches(0.05)), Inches(3), Inches(0.35),
                name, size=Pt(15), color=BLACK, bold=True)

    # Timeline
    add_textbox(s39, Inches(1.2), Emu(row_y + Inches(0.38)), Inches(2.5), Inches(0.3),
                timeline, size=Pt(11), color=GRAY_MED)

    # Description
    add_textbox(s39, Inches(4.2), Emu(row_y + Inches(0.12)), Inches(7.0), Inches(0.5),
                desc, size=Pt(13), color=RGBColor(0x55,0x55,0x55))

    # Special tag for Domain Packs
    if name == 'Domain-Specific Packs':
        add_textbox(s39, Inches(3.6), Emu(row_y + Inches(0.05)), Inches(0.5), Inches(0.3),
                    '★', size=Pt(14), color=RGBColor(0xFF, 0xD7, 0x00), bold=True)

# Bottom recommendation
rec_paras = [
    {'text': '推荐策略', 'size': Pt(14), 'color': RGBColor(0x33,0x33,0x33), 'bold': True, 'space_after': Pt(4)},
    {'text': '短期以 Spec Quality Gate 作为 MVP 快速验证产品价值，同步启动 Domain-Specific Packs（金融领域）的建设。'
             '中期补齐 Multi-Agent 和 Traceability 能力，长期向平台化和形式化验证演进。'
             'Domain Packs 贯穿始终，是从短期到长期持续积累竞争壁垒的核心抓手。',
     'size': Pt(12), 'space_after': Pt(4)},
]
add_rich_textbox(s39, Inches(0.6), Inches(6.2) - Inches(0.7), Inches(11), Inches(1.0), rec_paras)

add_page_num(s39, 39, TOTAL)
print("Slide 39: Priority Overview done")


# --- Update ALL page numbers to /39 ---
print("\nUpdating page numbers...")
for idx, slide in enumerate(prs.slides):
    page = idx + 1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top > 6000000 and shape.left > 9000000:
            txt = shape.text_frame.text.strip()
            if re.match(r'^\d+/\d+$', txt):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = f"{page}/{TOTAL}"
                        print(f"  Slide {page}: updated page number to {page}/{TOTAL}")

prs.save(PPTX)
print(f"\nSaved! Total slides: {len(prs.slides)}")
