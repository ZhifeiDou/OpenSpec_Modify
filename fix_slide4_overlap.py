"""
Fix: Move conclusion box below the new Exp4 content to prevent overlap.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

prs = Presentation('OpenSpec_技术研究汇报_精简版.pptx')
slide = prs.slides[3]  # Slide 4

# The Exp4 table + total row ends at y=6456810.
# Move conclusion to below that.
NEW_CONCLUSION_TOP = 6490000
NEW_CONCLUSION_H = 340000

for shape in slide.shapes:
    if shape.name == 'Rectangle 191':  # Conclusion background
        shape.top = NEW_CONCLUSION_TOP
        shape.height = NEW_CONCLUSION_H
    elif shape.name == 'TextBox 192':  # Conclusion text
        shape.top = NEW_CONCLUSION_TOP + 15000
        shape.height = NEW_CONCLUSION_H - 30000
        # Also update font size to fit in smaller box
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > Pt(9):
                    run.font.size = Pt(9)

prs.save('OpenSpec_技术研究汇报_精简版.pptx')
print("Fixed: conclusion moved below Exp4 content.")
