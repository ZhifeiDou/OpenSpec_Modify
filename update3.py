import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ICONS_DIR = 'icons'
prs = Presentation('OpenSpec_\u6280\u672f\u7814\u7a76\u6c47\u62a5_\u7cbe\u7b80\u7248_update.pptx')
slide = prs.slides[2]

# Clear
spTree = slide.shapes._spTree
for tag in ['p:sp','p:grpSp','p:cxnSp','p:pic']:
    for el in list(spTree.findall(qn(tag))): spTree.remove(el)

def set_ea(runs):
    for r in runs:
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None: ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', 'Microsoft YaHei')

def add_tb(sl, l, t, w, h, text, sz=Pt(12), bold=False, color=RGBColor(0x44,0x44,0x44), align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(l, t, w, h); tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.name = 'Microsoft YaHei'; p.font.size = sz; p.font.bold = bold; p.font.color.rgb = color
    set_ea(p.runs); return box

def add_rect(sl, l, t, w, h, fc):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc; s.line.fill.background(); return s

def add_rrect(sl, l, t, w, h, fc, bc=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_icon_tag(sl, x, y, icon_file, text, fill_color, tag_w, tag_h=Emu(210000),
                 icon_sz=Emu(170000), text_color=RGBColor(0xFF,0xFF,0xFF), fsz=Pt(8)):
    tag = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, tag_w, tag_h)
    tag.fill.solid(); tag.fill.fore_color.rgb = fill_color; tag.line.fill.background()
    tf = tag.text_frame; tf.word_wrap = False
    tf.margin_left = Emu(25000); tf.margin_right = Emu(25000)
    tf.margin_top = Emu(10000); tf.margin_bottom = Emu(10000)
    icon_path = os.path.join(ICONS_DIR, icon_file) if icon_file else None
    has_icon = icon_path and os.path.exists(icon_path)
    if has_icon:
        pad = Emu(15000); iy = y + (tag_h - icon_sz) // 2
        sl.shapes.add_picture(icon_path, x + pad, iy, icon_sz, icon_sz)
        tx2 = x + pad + icon_sz + Emu(10000); tw2 = tag_w - pad - icon_sz - Emu(35000)
        tb = sl.shapes.add_textbox(tx2, y, tw2, tag_h)
        p = tb.text_frame.paragraphs[0]; p.text = text
        p.font.name = 'Microsoft YaHei'; p.font.size = fsz; p.font.bold = True; p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.LEFT; set_ea(p.runs)
        tb.text_frame._txBody.find(qn('a:bodyPr')).set('anchor', 'ctr')
    else:
        p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Microsoft YaHei'; p.font.size = fsz; p.font.bold = True; p.font.color.rgb = text_color
        set_ea(p.runs)

# DIMENSIONS
LX=Emu(300000); LW=Emu(4200000); MH=Emu(1400000); LH=Emu(1200000)
LGAP=Emu(120000); SY=Emu(1060000)
LABEL_W=Emu(1300000); TS=LX+LABEL_W+Emu(50000)
TH=Emu(210000); TGX=Emu(50000); TGY=Emu(35000); ISZ=Emu(170000)
TF=Pt(8); CW=Emu(72000); TP=Emu(100000); IE=Emu(200000)
nx=LX+Emu(60000)

# Positions
y1=SY
y2=y1+MH+LGAP
y3=y2+LH+LGAP

# GLOBAL
add_rect(slide, 0, 0, Emu(12191695), Emu(54864), RGBColor(0x00,0x78,0xD4))
add_tb(slide, Emu(300000), Emu(240000), Emu(4500000), Emu(380000),
       'Agentic Coding \u6280\u672f\u6808\u5168\u666f', sz=Pt(22), bold=True, color=RGBColor(0x1A,0x1A,0x2E))
add_rect(slide, Emu(300000), Emu(660000), Emu(650000), Emu(30000), RGBColor(0x00,0x78,0xD4))
add_tb(slide, Emu(300000), Emu(730000), Emu(3800000), Emu(230000),
       '\u4e09\u5c42\u67b6\u6784\u652f\u6491 AI \u8f85\u52a9\u5f00\u53d1\u5168\u6d41\u7a0b', sz=Pt(10), color=RGBColor(0x44,0x44,0x44))

# ===== LAYER 1: merged framework (LEFT-RIGHT split) =====
fw_accent = RGBColor(0x10,0x7C,0x10)
add_rrect(slide, LX, y1, LW, MH, RGBColor(0xE8,0xF5,0xE9), fw_accent)
add_rect(slide, LX, y1+Emu(6000), Emu(38000), MH-Emu(12000), fw_accent)

# Main label
add_tb(slide, nx, y1+Emu(50000), Emu(2000000), Emu(220000),
       '\u2460 \u6846\u67b6\u5c42', sz=Pt(12), bold=True, color=fw_accent)
add_tb(slide, nx, y1+Emu(260000), Emu(2000000), Emu(160000),
       'Framework', sz=Pt(7), color=RGBColor(0x99,0x99,0x99))

# Left-right split dimensions
center_x = LX + LW // 2
left_x = LX + Emu(60000)
left_mx = center_x - Emu(40000)
right_x = center_x + Emu(40000)
right_mx = LX + LW - Emu(35000)

sub_y = y1 + Emu(350000)   # sub-label y
tag_y = y1 + Emu(550000)   # tags start y

# Left sub-section: process methods
add_tb(slide, left_x, sub_y, Emu(1800000), Emu(180000),
       '\u8fc7\u7a0b\u65b9\u6cd5\u5c42', sz=Pt(9), bold=True, color=RGBColor(0x2E,0x7D,0x32))
add_tb(slide, left_x, sub_y+Emu(160000), Emu(1800000), Emu(130000),
       'Process & Methodology', sz=Pt(6), color=RGBColor(0x99,0x99,0x99))

tools1 = [
    ('OpenSpec','openspec.png',RGBColor(0x10,0x7C,0x10)),
    ('Speckit','speckit.png',RGBColor(0x2E,0x7D,0x32)),
    ('Superpowers',None,RGBColor(0x38,0x8E,0x3C)),
    ('BMAD','bmad.png',RGBColor(0x43,0xA0,0x47)),
    ('Agile Artifacts',None,RGBColor(0x4C,0xAF,0x50)),
]
tx = left_x; ty = tag_y
for name, icon, color in tools1:
    hi = icon and os.path.exists(os.path.join(ICONS_DIR, icon))
    bw = max(Emu(480000), len(name)*CW+TP)
    tw = bw + (IE if hi else Emu(0))
    if tx + tw > left_mx: tx = left_x; ty += TH + TGY
    add_icon_tag(slide, tx, ty, icon, name, color, tw, TH, ISZ, fsz=TF)
    tx += tw + TGX

# Vertical dashed separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x, y1+Emu(320000), Emu(10), MH-Emu(360000))
sep.fill.background()
sep.line.color.rgb = RGBColor(0xAA,0xAA,0xAA)
sep.line.width = Pt(1)
ln_el = sep._element.find('.//'+qn('a:ln'))
if ln_el is not None: etree.SubElement(ln_el, qn('a:prstDash')).set('val','dash')

# Right sub-section: tool extensions
add_tb(slide, right_x, sub_y, Emu(1800000), Emu(180000),
       '\u5de5\u5177\u6269\u5c55\u5c42', sz=Pt(9), bold=True, color=RGBColor(0x00,0x78,0xD4))
add_tb(slide, right_x, sub_y+Emu(160000), Emu(1800000), Emu(130000),
       'Tool Extensions & MCP', sz=Pt(6), color=RGBColor(0x99,0x99,0x99))

tools2 = [
    ('Everything',None,RGBColor(0x00,0x78,0xD4)),
    ('Context7',None,RGBColor(0x15,0x65,0xC0)),
    ('MCP Servers',None,RGBColor(0x1E,0x88,0xE5)),
    ('Playwright','playwright.png',RGBColor(0x21,0x96,0xF3)),
    ('Seq. Thinking',None,RGBColor(0x42,0xA5,0xF5)),
]
tx = right_x; ty = tag_y
for name, icon, color in tools2:
    hi = icon and os.path.exists(os.path.join(ICONS_DIR, icon))
    bw = max(Emu(480000), len(name)*CW+TP)
    tw = bw + (IE if hi else Emu(0))
    if tx + tw > right_mx: tx = right_x; ty += TH + TGY
    add_icon_tag(slide, tx, ty, icon, name, color, tw, TH, ISZ, fsz=TF)
    tx += tw + TGX

# ===== LAYER 2: Code Agent =====
ca_accent = RGBColor(0xE8,0x6C,0x00)
add_rrect(slide, LX, y2, LW, LH, RGBColor(0xFF,0xF3,0xE0), ca_accent)
add_rect(slide, LX, y2+Emu(6000), Emu(38000), LH-Emu(12000), ca_accent)
add_tb(slide, nx, y2+Emu(65000), LABEL_W, Emu(210000),
       '\u2461 Code Agent\u5c42', sz=Pt(11), bold=True, color=ca_accent)
add_tb(slide, nx, y2+Emu(275000), LABEL_W, Emu(160000),
       'Coding Agents', sz=Pt(7), color=RGBColor(0x99,0x99,0x99))

tools3 = [
    ('Claude Code','claude.png',RGBColor(0xE6,0x5C,0x00)),
    ('Cursor','cursor.png',RGBColor(0xEF,0x6C,0x00)),
    ('Windsurf','windsurf.png',RGBColor(0xF5,0x7C,0x00)),
    ('Cline','cline.png',RGBColor(0xFB,0x8C,0x00)),
    ('OpenCode',None,RGBColor(0xFF,0x9E,0x00)),
    ('Aider','aider.png',RGBColor(0xFF,0xA7,0x26)),
]
mx3 = LX + LW - Emu(35000)
tx = TS; ty = y2 + Emu(60000)
for name, icon, color in tools3:
    hi = icon and os.path.exists(os.path.join(ICONS_DIR, icon))
    bw = max(Emu(480000), len(name)*CW+TP)
    tw = bw + (IE if hi else Emu(0))
    if tx + tw > mx3: tx = TS; ty += TH + TGY
    add_icon_tag(slide, tx, ty, icon, name, color, tw, TH, ISZ, fsz=TF)
    tx += tw + TGX

# ===== LAYER 3: Foundation Models =====
fm_accent = RGBColor(0xD1,0x34,0x38)
add_rrect(slide, LX, y3, LW, LH, RGBColor(0xFD,0xE8,0xE8), fm_accent)
add_rect(slide, LX, y3+Emu(6000), Emu(38000), LH-Emu(12000), fm_accent)
add_tb(slide, nx, y3+Emu(65000), LABEL_W, Emu(210000),
       '\u2462 \u5e95\u5ea7\u5927\u6a21\u578b\u5c42', sz=Pt(11), bold=True, color=fm_accent)
add_tb(slide, nx, y3+Emu(275000), LABEL_W, Emu(160000),
       'Foundation Models', sz=Pt(7), color=RGBColor(0x99,0x99,0x99))

tools4 = [
    ('Claude 4.6','claude.png',RGBColor(0xC6,0x28,0x28)),
    ('GPT-5.1','openai.png',RGBColor(0xD3,0x2F,0x2F)),
    ('Gemini 3','gemini.png',RGBColor(0xE5,0x39,0x35)),
    ('DeepSeek V3','deepseek.png',RGBColor(0xEF,0x53,0x50)),
    ('Kimi K2.5','kimi.png',RGBColor(0xE5,0x73,0x73)),
    ('GLM-5',None,RGBColor(0xEF,0x9A,0x9A)),
]
tx = TS; ty = y3 + Emu(60000)
for name, icon, color in tools4:
    hi = icon and os.path.exists(os.path.join(ICONS_DIR, icon))
    bw = max(Emu(480000), len(name)*CW+TP)
    tw = bw + (IE if hi else Emu(0))
    if tx + tw > mx3: tx = TS; ty += TH + TGY
    add_icon_tag(slide, tx, ty, icon, name, color, tw, TH, ISZ, fsz=TF)
    tx += tw + TGX

# ===== NO BRACES =====

# ===== RIGHT PLACEHOLDER =====
add_tb(slide, Emu(5600000), Emu(730000), Emu(6100000), Emu(280000),
       '\u7cfb\u7edf\u5de5\u7a0b\u6846\u67b6 vs Code Agent', sz=Pt(15), bold=True, color=RGBColor(0x1A,0x1A,0x2E))
ph=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(5600000), Emu(1060000), Emu(6140730), Emu(5400000))
ph.fill.solid(); ph.fill.fore_color.rgb=RGBColor(0xF8,0xF9,0xFA)
ph.line.color.rgb=RGBColor(0xBB,0xBB,0xBB); ph.line.width=Pt(1.5)
ln=ph._element.find('.//'+qn('a:ln'))
if ln is not None: etree.SubElement(ln, qn('a:prstDash')).set('val','dash')
tf=ph.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.text='\u5728\u6b64\u5904\u6dfb\u52a0\u5185\u5bb9...'
p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(14); p.font.color.rgb=RGBColor(0xCC,0xCC,0xCC)
p.font.name='Microsoft YaHei'; ph.text_frame._txBody.find(qn('a:bodyPr')).set('anchor','ctr')

add_tb(slide, Emu(11094415), Emu(6400800), Emu(914400), Emu(365760), '', sz=Pt(11),
       color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)

prs.save('OpenSpec_\u6280\u672f\u7814\u7a76\u6c47\u62a5_\u7cbe\u7b80\u7248_update.pptx')
print('Saved!')

import comtypes.client
comtypes.CoInitialize()
try:
    pp=comtypes.client.CreateObject('PowerPoint.Application'); pp.Visible=1
    pres=pp.Presentations.Open(os.path.abspath('OpenSpec_\u6280\u672f\u7814\u7a76\u6c47\u62a5_\u7cbe\u7b80\u7248_update.pptx'),WithWindow=False)
    pres.Slides(3).Export(os.path.abspath('slide3_v5.png'),'PNG',2560,1440)
    pres.Close(); pp.Quit(); print('Exported')
except Exception as e: print(f'Error: {e}')
finally: comtypes.CoUninitialize()
